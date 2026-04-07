import os
from pptx import Presentation
import cv2
import numpy as np
from tqdm import tqdm

def create_output_folder(pptx_file):
    """Create an output folder with the same name as the PPTX file if it doesn't exist."""
    folder_name = os.path.splitext(pptx_file)[0]
    os.makedirs(folder_name, exist_ok=True)
    return folder_name

def extract_and_upscale_images(pptx_file, output_folder, target_resolution=(3840, 2160)):
    """Extract images from a PPTX file, upscale them to 4K, and save them in the output folder."""
    prs = Presentation(pptx_file)
    image_count = 1

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'image'):
                image = shape.image
                image_bytes = image.blob
                
                # Convert image bytes to numpy array
                nparr = np.frombuffer(image_bytes, np.uint8)
                img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
                
                if img is None:
                    continue

                height, width = img.shape[:2]
                
                # Calculate scaling factors to reach the target resolution
                scale_x = target_resolution[0] / width
                scale_y = target_resolution[1] / height
                scale = max(scale_x, scale_y)
                
                # Only upscale if the image is smaller than the target resolution
                if scale > 1:
                    new_width = int(width * scale)
                    new_height = int(height * scale)
                    
                    # Resize the image using Lanczos interpolation for better quality
                    img = cv2.resize(img, (new_width, new_height), interpolation=cv2.INTER_LANCZOS4)
                
                output_path = os.path.join(output_folder, f"{image_count:02d}.png")
                cv2.imwrite(output_path, img)
                
                image_count += 1

    return image_count - 1

def process_presentations(root_dir):
    """Process all PPTX files starting with 'Lecture' in each section folder within course folders."""
    for course_folder in os.scandir(root_dir):
        if not course_folder.is_dir():
            continue
        
        for section_folder in os.scandir(course_folder.path):
            if not section_folder.is_dir():
                continue
            
            # Ensure all_pptx folder exists
            all_pptx_folder = os.path.join(section_folder.path, 'all_pptx')
            os.makedirs(all_pptx_folder, exist_ok=True)
            
            # Get list of already processed files
            processed_files = set(os.listdir(all_pptx_folder))
            
            # Process PPTX files in this section folder
            pptx_files = [
                f for f in os.listdir(section_folder.path)
                if f.startswith("Lecture") and f.endswith(".pptx") and f not in processed_files
            ]
            
            for filename in tqdm(pptx_files, desc=f"Processing {section_folder.name}"):
                pptx_file = os.path.join(section_folder.path, filename)
                
                output_folder = create_output_folder(pptx_file)
                extract_and_upscale_images(pptx_file, output_folder)
                
                # Move the file to all_pptx folder
                new_pptx_path = os.path.join(all_pptx_folder, filename)
                os.rename(pptx_file, new_pptx_path)

    print("All presentations processed and moved.")

if __name__ == "__main__":
    # Set the root directory to 'input' in the current working directory
    folder = os.path.join(os.getcwd(), "input")
    # Start processing presentations in the input folder
    process_presentations(folder)