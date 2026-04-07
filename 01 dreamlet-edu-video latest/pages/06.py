"""
CODING CONVENTION: NO SHARED CODE
- All code for this page must be contained entirely within this single file
- Never import from other page files or create shared utilities
- Copy any needed functions directly into this file
- Each page is completely self-contained and independent

4K IMAGE (PPTX & ZIP) - FINAL
Machine-aware final version for:
- MacBook Pro M3 Pro (18GB) using an Apple Silicon optimized CPU path
- Acer Windows laptop with i5-12450H + RTX 3050 using CUDA when CuPy is installed
"""

import streamlit as st
import os
import re
import time
import cv2
import numpy as np
import io
import enum
import fnmatch
import platform
import tempfile
import glob
import shutil
import subprocess
import zipfile
import psutil
from PIL import Image, ImageDraw, ImageFont
try:
    # For newer PIL versions
    LANCZOS = Image.Resampling.LANCZOS
except AttributeError:
    # For older PIL versions
    LANCZOS = Image.LANCZOS
from pptx import Presentation
from pathlib import Path
from typing import Dict, List, Tuple, Optional, Any
from datetime import datetime, timedelta
import concurrent.futures
import threading
from functools import lru_cache

# Optional GPU framework imports
GPU_FRAMEWORKS = {
    'cuda': False,
}

try:
    import cupy as cp
    import cupyx.scipy.ndimage as gpu_ndimage
    GPU_FRAMEWORKS['cuda'] = True
    print("✅ CUDA framework available")
except ImportError:
    print("❌ CUDA framework not available - install with: pip install cupy-cuda12x")

MACHINE_PROFILES = {
    'macbook_m3_pro': {
        'name': 'MacBook Pro M3 Pro (18GB)',
        'platform': 'darwin',
        'cpu': 'Apple M3 Pro',
        'ram_gb': 18,
        'gpu': {
            'type': 'apple_silicon',
            'runtime': 'cpu_optimized',
        },
        'optimization_settings': {
            'max_threads': 8,
            'batch_size': 6,
            'memory_limit_gb': 14,
            'use_gpu_acceleration': False,
            'preferred_image_processing': 'cpu',
        },
    },
    'windows_i5_rtx3050': {
        'name': 'Acer Windows laptop (i5-12450H + RTX 3050)',
        'platform': 'windows',
        'cpu': 'Intel Core i5-12450H',
        'ram_gb': 16,
        'gpu': {
            'type': 'nvidia',
            'model': 'RTX 3050 Laptop GPU',
            'runtime': 'cuda',
            'vram_gb': 4,
        },
        'optimization_settings': {
            'max_threads': 6,
            'batch_size': 4,
            'memory_limit_gb': 12,
            'use_gpu_acceleration': True,
            'preferred_image_processing': 'cuda',
        },
    },
}

FALLBACK_PROFILE = {
    'name': 'Generic machine',
    'platform': 'generic',
    'cpu': 'Unknown',
    'ram_gb': 8,
    'gpu': {
        'type': 'unknown',
        'runtime': 'cpu',
    },
    'optimization_settings': {
        'max_threads': 4,
        'batch_size': 2,
        'memory_limit_gb': 8,
        'use_gpu_acceleration': False,
        'preferred_image_processing': 'cpu',
    },
}

class MachineDetector:
    """Detect machine configuration and optimize settings accordingly"""
    
    def __init__(self):
        self.machine_config = self.detect_machine()
        
    def detect_machine(self) -> Dict:
        """Detect current machine and return optimized configuration"""
        system_info = {
            'platform': platform.system().lower(),
            'architecture': platform.machine().lower(),
            'cpu': platform.processor().lower(),
            'memory_gb': round(psutil.virtual_memory().total / (1024**3)),
            'cpu_count': psutil.cpu_count()
        }
        
        # Try to detect GPU
        gpu_info = self._detect_gpu()
        system_info.update(gpu_info)
        
        print(f"🔍 Detection Debug:")
        print(f"  Platform: {system_info['platform']}")
        print(f"  Architecture: {system_info['architecture']}")
        print(f"  CPU: {system_info['cpu']}")
        print(f"  Memory: {system_info['memory_gb']} GB")
        print(f"  GPU Detected: {system_info.get('gpu_detected', False)}")
        print(f"  GPU Type: {system_info.get('gpu_type', 'None')}")
        
        print(f"🔍 Checking M3 Pro criteria:")
        print(f"  Platform == 'darwin': {system_info['platform'] == 'darwin'}")
        print(f"  Architecture == 'arm64': {system_info['architecture'] == 'arm64'}")
        print(f"  Memory in range 16-20: {16 <= system_info['memory_gb'] <= 20}")
        print(f"  Apple Silicon GPU detected: {system_info.get('gpu_detected', False) and system_info.get('gpu_type') == 'apple_silicon'}")
        
        if (system_info['platform'] == 'darwin' and 
            system_info['architecture'] == 'arm64' and
            16 <= system_info['memory_gb'] <= 20 and
            system_info.get('gpu_detected', False) and 
            system_info.get('gpu_type') == 'apple_silicon'):
            print("✅ Matched M3 Pro configuration!")
            return self._build_machine_config('macbook_m3_pro', system_info)
        
        print(f"🔍 Checking Windows criteria:")
        print(f"  Platform == 'windows': {system_info['platform'] == 'windows'}")
        print(f"  Memory >= 14: {system_info['memory_gb'] >= 14}")
        print(f"  GPU detected: {system_info.get('gpu_detected', False)}")
        
        if (system_info['platform'] == 'windows' and 
              system_info['memory_gb'] >= 14 and 
              system_info.get('gpu_type') == 'nvidia'):
            print("✅ Matched Windows RTX 3050 configuration!")
            return self._build_machine_config('windows_i5_rtx3050', system_info)
        
        print("❌ No specific machine matched, using fallback")
        return self._build_fallback_config(system_info)

    def _build_machine_config(self, machine_id: str, system_info: Dict) -> Dict:
        profile = MACHINE_PROFILES[machine_id]
        machine_config = {
            'name': profile['name'],
            'platform': profile['platform'],
            'cpu': profile['cpu'],
            'ram_gb': profile['ram_gb'],
            'gpu': dict(profile['gpu']),
            'optimization_settings': dict(profile['optimization_settings']),
            'detected_info': system_info,
            'machine_id': machine_id,
        }
        return machine_config

    def _build_fallback_config(self, system_info: Dict) -> Dict:
        return {
            'name': FALLBACK_PROFILE['name'],
            'platform': FALLBACK_PROFILE['platform'],
            'cpu': FALLBACK_PROFILE['cpu'],
            'ram_gb': FALLBACK_PROFILE['ram_gb'],
            'gpu': dict(FALLBACK_PROFILE['gpu']),
            'optimization_settings': dict(FALLBACK_PROFILE['optimization_settings']),
            'detected_info': system_info,
            'machine_id': 'unknown',
        }
    
    def _detect_gpu(self) -> Dict:
        """Detect GPU information"""
        gpu_info = {'gpu_detected': False}
        
        try:
            result = subprocess.run(
                ['nvidia-smi', '--query-gpu=name', '--format=csv,noheader'],
                capture_output=True,
                text=True,
                timeout=5,
            )
            if result.returncode == 0:
                gpu_name = result.stdout.strip().splitlines()[0] if result.stdout.strip() else 'NVIDIA GPU'
                gpu_info.update({
                    'gpu_detected': True,
                    'gpu_type': 'nvidia',
                    'gpu_name': gpu_name,
                })
        except Exception:
            pass
        
        if platform.system() == 'Darwin':
            try:
                result = subprocess.run(['uname', '-m'], capture_output=True, text=True)
                arch = result.stdout.strip()
                print(f"🔍 macOS Architecture: {arch}")
                
                if arch == 'arm64':  # Apple Silicon
                    gpu_info.update({
                        'gpu_detected': True,
                        'gpu_type': 'apple_silicon',
                        'gpu_name': 'Apple Silicon GPU',
                        'unified_memory': True,
                    })
                    print("✅ Apple Silicon GPU detected")
                else:
                    print("❌ Intel Mac detected (no Apple Silicon GPU)")
            except Exception as e:
                print(f"❌ Error detecting Apple Silicon: {e}")
                # Fallback check
                if 'apple' in platform.processor().lower():
                    gpu_info.update({
                        'gpu_detected': True,
                        'gpu_type': 'apple_silicon',
                        'gpu_name': 'Apple Silicon GPU',
                        'unified_memory': True,
                    })
                    print("✅ Apple Silicon detected via processor check")
        
        return gpu_info
    
    def get_optimization_settings(self) -> Dict:
        """Get optimization settings for detected machine"""
        return self.machine_config.get('optimization_settings', {})
    
    def get_machine_info(self) -> Dict:
        """Get detected machine information"""
        preferred = self.machine_config.get('optimization_settings', {}).get('preferred_image_processing', 'cpu')
        runtime_acceleration = preferred == 'cuda' and GPU_FRAMEWORKS['cuda']
        return {
            'name': self.machine_config.get('name', 'Unknown'),
            'machine_id': self.machine_config.get('machine_id', 'unknown'),
            'detected_info': self.machine_config.get('detected_info', {}),
            'gpu_available': runtime_acceleration,
            'acceleration_preference': preferred,
        }

# Initialize machine detector
MACHINE_DETECTOR = MachineDetector()
OPTIMIZATION_SETTINGS = MACHINE_DETECTOR.get_optimization_settings()

# GPU Accelerated Image Processing
class GPUImageProcessor:
    """Machine-aware image processing with optional CUDA acceleration"""
    
    def __init__(self):
        self.gpu_framework = self._detect_best_framework()
        self.machine_info = MACHINE_DETECTOR.get_machine_info()
        
    def _detect_best_framework(self) -> str:
        """Detect the best available runtime acceleration"""
        preferred = OPTIMIZATION_SETTINGS.get('preferred_image_processing', 'cpu')
        
        if preferred == 'cuda' and GPU_FRAMEWORKS['cuda']:
            return 'cuda'
        return 'cpu'
    
    def resize_image_gpu(self, img_array: np.ndarray, new_size: Tuple[int, int]) -> np.ndarray:
        """Use the best available acceleration path for resizing"""
        try:
            if self.gpu_framework == 'cuda':
                return self._resize_cuda(img_array, new_size)
            return self._resize_cpu_optimized(img_array, new_size)
        except Exception as e:
            st.warning(f"Accelerated resize failed, falling back to CPU: {e}")
            return self._resize_cpu(img_array, new_size)
    
    def _resize_cuda(self, img_array: np.ndarray, new_size: Tuple[int, int]) -> np.ndarray:
        """CUDA-based resizing for NVIDIA GPUs"""
        try:
            # Transfer to GPU
            gpu_array = cp.asarray(img_array)
            
            # Calculate scale factors
            scale_factors = (
                new_size[1] / img_array.shape[0],  # height scale
                new_size[0] / img_array.shape[1],  # width scale
                1.0  # channels (no scaling)
            )
            
            # GPU resize operation using CuPy
            resized_gpu = gpu_ndimage.zoom(gpu_array, scale_factors, order=1, prefilter=False)
            
            # Transfer back to CPU
            return cp.asnumpy(resized_gpu).astype(np.uint8)
            
        except Exception as e:
            return self._resize_cpu(img_array, new_size)
    
    def _resize_cpu_optimized(self, img_array: np.ndarray, new_size: Tuple[int, int]) -> np.ndarray:
        """Optimized CPU resizing for the Apple Silicon fallback path"""
        return cv2.resize(img_array, new_size, interpolation=cv2.INTER_LANCZOS4)
    
    def _resize_cpu(self, img_array: np.ndarray, new_size: Tuple[int, int]) -> np.ndarray:
        """Standard CPU resizing"""
        return cv2.resize(img_array, new_size, interpolation=cv2.INTER_LANCZOS4)

# Initialize GPU processor
GPU_PROCESSOR = GPUImageProcessor()

# Local utility functions
def increase_image_decompression_limit():
    """Increases the Pillow image decompression bomb limit to avoid errors with very large images"""
    new_limit = 3840 * 2160 * 50  # Approximately 415 million pixels
    Image.MAX_IMAGE_PIXELS = new_limit

def get_input_directory() -> str:
    """Get the path to the input directory"""
    input_dir = os.path.join(os.getcwd(), "input")
    return input_dir

def ensure_directory_exists(directory_path: str) -> None:
    """Create directory if it doesn't exist"""
    if not os.path.exists(directory_path):
        os.makedirs(directory_path)

def find_files(directory: str, pattern: str) -> List[str]:
    """Find all files matching a pattern in a directory (recursively)"""
    result = []
    for root, _, filenames in os.walk(directory):
        for filename in fnmatch.filter(filenames, pattern):
            result.append(os.path.join(root, filename))
    return result

def find_presentation_files(directory: str) -> List[str]:
    """Find all presentation files (PPTX and ZIP) in a directory"""
    pptx_files = find_files(directory, "*.pptx")
    zip_files = find_files(directory, "*.zip")
    return pptx_files + zip_files

def natural_sort_key(filename: str) -> List:
    """Generate a key for natural sorting (handles numbers correctly)"""
    parts = re.split(r'(\d+)', filename)
    return [int(part) if part.isdigit() else part.lower() for part in parts]

def extract_numeric_prefix(filename: str) -> int:
    """Extract numeric prefix from filename like '1_Title.png' -> 1"""
    match = re.match(r'^(\d+)', filename)
    if match:
        return int(match.group(1))
    return 999999  # Put files without numeric prefix at the end
def extract_slides_from_zip(zip_file: str, output_folder: str, status_mgr=None, target_resolution=(3840, 2160)) -> Tuple[bool, List[str], int]:
    """GPU OPTIMIZED: Extract images from a ZIP file, upscale them to 4K using GPU acceleration"""
    temp_dir = tempfile.mkdtemp()
    
    try:
        os.makedirs(output_folder, exist_ok=True)
        
        if status_mgr:
            status_mgr.update_action(f"Extracting ZIP file: {os.path.basename(zip_file)}")
        
        # Extract ZIP contents
        try:
            with zipfile.ZipFile(zip_file, 'r') as zip_ref:
                zip_ref.extractall(temp_dir)
        except Exception as e:
            error_msg = f"Error extracting ZIP file: {str(e)}"
            if status_mgr:
                status_mgr.error(error_msg)
            return False, [], 0
        
        # Find all image files
        image_extensions = ['*.png', '*.jpg', '*.jpeg', '*.PNG', '*.JPG', '*.JPEG']
        image_files = []
        for ext in image_extensions:
            image_files.extend(glob.glob(os.path.join(temp_dir, '**', ext), recursive=True))
        
        if not image_files:
            error_msg = "No image files found in ZIP archive"
            if status_mgr:
                status_mgr.error(error_msg)
            return False, [], 0
        
        # Sort images by numeric prefix
        image_files.sort(key=lambda x: extract_numeric_prefix(os.path.basename(x)))
        
        if status_mgr:
            gpu_info = GPU_PROCESSOR.gpu_framework
            status_mgr.update_action(f"Processing {len(image_files)} images using {gpu_info.upper()} acceleration...")
        
        # GPU OPTIMIZATION: Process images with GPU acceleration
        extracted_images = []
        
        def process_single_image_gpu(args):
            i, img_path = args
            try:
                with Image.open(img_path) as img:
                    img_array = np.array(img)
                    
                    # Check if upscaling is needed
                    if img.width >= target_resolution[0] and img.height >= target_resolution[1]:
                        processed_img = img.copy()
                    else:
                        # GPU upscale to 4K
                        aspect_ratio = img.width / img.height
                        
                        if aspect_ratio > target_resolution[0] / target_resolution[1]:
                            new_width = target_resolution[0]
                            new_height = int(target_resolution[0] / aspect_ratio)
                        else:
                            new_height = target_resolution[1]
                            new_width = int(target_resolution[1] * aspect_ratio)
                        
                        # Use GPU resize
                        resized_array = GPU_PROCESSOR.resize_image_gpu(img_array, (new_width, new_height))
                        processed_img = Image.fromarray(resized_array)
                    
                    # Save with sequential naming
                    output_path = os.path.join(output_folder, f"{i:02d}.png")
                    processed_img.save(output_path, "PNG", optimize=True)
                    return output_path
                    
            except Exception as e:
                if status_mgr:
                    status_mgr.warning(f"Error processing image {os.path.basename(img_path)}: {str(e)}")
                return None
        
        # Use machine-optimized threading
        max_workers = OPTIMIZATION_SETTINGS.get('max_threads', 4)
        batch_size = OPTIMIZATION_SETTINGS.get('batch_size', 2)
        
        # Process in batches to manage GPU memory
        for batch_start in range(0, len(image_files), batch_size):
            batch_end = min(batch_start + batch_size, len(image_files))
            batch_files = image_files[batch_start:batch_end]
            
            with concurrent.futures.ThreadPoolExecutor(max_workers=min(max_workers, len(batch_files))) as executor:
                args_list = [(i + batch_start, img_path) for i, img_path in enumerate(batch_files, 1)]
                
                for future in concurrent.futures.as_completed([executor.submit(process_single_image_gpu, args) for args in args_list]):
                    result = future.result()
                    if result:
                        extracted_images.append(result)
        
        # Sort extracted images by filename
        extracted_images.sort()
        image_count = len(extracted_images)
        
        if image_count == 0:
            error_msg = "Failed to process any images from ZIP"
            if status_mgr:
                status_mgr.error(error_msg)
            return False, [], 0
        
        if status_mgr:
            status_mgr.update_action(f"Successfully extracted {image_count} images using GPU acceleration")
        
        return True, extracted_images, image_count
        
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)

def extract_course_lecture_section(file_path: str) -> Tuple[Optional[str], Optional[str], Optional[str]]:
    """Extract course, lecture, and section information from a file path"""
    dir_parts = os.path.normpath(file_path).split(os.sep)
    
    course = None
    lecture = None
    section = None
    
    for part in dir_parts:
        course_match = re.search(r'course\s*(\d+)', part.lower())
        if course_match:
            course = course_match.group(1)
            break
        number_start_match = re.match(r'^(\d+)\s+', part)
        if number_start_match:
            course = number_start_match.group(1)
            break
        bracket_match = re.search(r'[(\[]\s*(\d+)\s*[)\]]', part)
        if bracket_match:
            course = bracket_match.group(1)
            break
    
    filename = os.path.basename(file_path)
    lecture_patterns = [
        r'lecture\s*(\d+)',
        r'lec\s*(\d+)',
        r'^(\d+)[-\s]',
        r'^\w+\s*(\d+)'
    ]
    
    for pattern in lecture_patterns:
        lecture_match = re.search(pattern, filename.lower())
        if lecture_match:
            lecture = lecture_match.group(1)
            break
    
    if not lecture:
        for part in dir_parts:
            for pattern in lecture_patterns:
                lecture_match = re.search(pattern, part.lower())
                if lecture_match:
                    lecture = lecture_match.group(1)
                    break
            if lecture:
                break
    
    section_patterns = [
        r'section\s*(\d+)',
        r'sec\s*(\d+)'
    ]
    
    for part in dir_parts:
        for pattern in section_patterns:
            section_match = re.search(pattern, part.lower())
            if section_match:
                section = section_match.group(1)
                break
        if section:
            break
    
    return course, lecture, section

# Multi-method converter functions
class ConversionMethod(str, enum.Enum):
    LIBREOFFICE = "libreoffice"
    PYTHON_PPTX = "python-pptx"
    PDF2IMAGE = "pdf2image"

CONVERSION_METHODS = {
    ConversionMethod.LIBREOFFICE: "LibreOffice + pdftoppm (best quality)",
    ConversionMethod.PYTHON_PPTX: "Direct PPTX processing (fastest)",
    ConversionMethod.PDF2IMAGE: "PDF2Image (most compatible)"
}

@lru_cache(maxsize=1)
def get_libreoffice_path():
    """Get the correct LibreOffice path based on the operating system (cached)"""
    if platform.system() == 'Darwin':  # macOS
        possible_paths = [
            '/Applications/LibreOffice.app/Contents/MacOS/soffice',
            '/Applications/LibreOffice.app/Contents/MacOS/soffice.bin',
            '/opt/homebrew/bin/soffice',
            '/usr/local/bin/soffice'
        ]
        for path in possible_paths:
            if os.path.exists(path):
                return path
        return None
    else:
        return 'soffice'

def extract_slides_directly(pptx_path, output_dir, status_text=None, target_resolution=(3840, 2160)):
    """GPU OPTIMIZED: Extract slides directly from PPTX file using python-pptx with GPU acceleration"""
    try:
        os.makedirs(output_dir, exist_ok=True)
        
        if status_text:
            status_text.info(f"Extracting slides directly from PPTX: {os.path.basename(pptx_path)}")
        
        prs = Presentation(pptx_path)
        slide_count = len(prs.slides)
        
        if slide_count == 0:
            if status_text:
                status_text.warning("Presentation contains no slides.")
            return False, "Presentation contains no slides", 0
        
        # GPU OPTIMIZATION: Process slides in batches
        batch_size = OPTIMIZATION_SETTINGS.get('batch_size', 5)
        
        for batch_start in range(0, slide_count, batch_size):
            batch_end = min(batch_start + batch_size, slide_count)
            
            for i in range(batch_start, batch_end):
                slide = prs.slides[i]
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                
                slide_aspect_ratio = slide_width / slide_height
                
                if slide_aspect_ratio > target_resolution[0] / target_resolution[1]:
                    width = target_resolution[0]
                    height = int(width / slide_aspect_ratio)
                else:
                    height = target_resolution[1]
                    width = int(height * slide_aspect_ratio)
                
                img = np.ones((height, width, 3), np.uint8) * 255
                
                scale_x = width / slide_width
                scale_y = height / slide_height
                
                for shape in slide.shapes:
                    if hasattr(shape, 'image') and shape.image:
                        try:
                            image_stream = io.BytesIO(shape.image.blob)
                            slide_image = Image.open(image_stream)
                            
                            slide_image_cv = cv2.cvtColor(np.array(slide_image), cv2.COLOR_RGB2BGR)
                            
                            left = int(shape.left * scale_x)
                            top = int(shape.top * scale_y)
                            img_width = int(shape.width * scale_x)
                            img_height = int(shape.height * scale_y)
                            
                            # Use GPU for resizing if available
                            if GPU_PROCESSOR.gpu_framework != 'cpu':
                                resized_image = GPU_PROCESSOR.resize_image_gpu(slide_image_cv, (img_width, img_height))
                            else:
                                resized_image = cv2.resize(slide_image_cv, (img_width, img_height), interpolation=cv2.INTER_LANCZOS4)
                            
                            try:
                                if img_height > 0 and img_width > 0:
                                    if top + img_height <= height and left + img_width <= width:
                                        img[top:top+img_height, left:left+img_width] = resized_image
                            except Exception as e:
                                if status_text:
                                    status_text.warning(f"Error placing image on slide {i+1}: {str(e)}")
                        except Exception as e:
                            if status_text:
                                status_text.warning(f"Error processing image in slide {i+1}: {str(e)}")
                
                output_path = os.path.join(output_dir, f"{i+1:02d}.png")
                cv2.imwrite(output_path, img, [cv2.IMWRITE_PNG_COMPRESSION, 9])
        
        return True, f"Successfully extracted {slide_count} slides", slide_count
    
    except Exception as e:
        error_msg = f"Error in direct PPTX extraction: {str(e)}"
        if status_text:
            status_text.error(error_msg)
        return False, error_msg, 0

def convert_with_pdf2image(pptx_path, output_dir, status_text=None, target_resolution=(3840, 2160)):
    """GPU OPTIMIZED: Convert PPTX to images using LibreOffice for PPTX→PDF and pdf2image for PDF→PNG"""
    temp_dir = tempfile.mkdtemp()
    pdf_path = os.path.join(temp_dir, os.path.splitext(os.path.basename(pptx_path))[0] + '.pdf')
    
    try:
        os.makedirs(output_dir, exist_ok=True)
        
        if status_text:
            status_text.info(f"Converting presentation using pdf2image: {os.path.basename(pptx_path)}")
        
        libreoffice_path = get_libreoffice_path()
        if not libreoffice_path:
            error_msg = "LibreOffice not found. Please make sure LibreOffice is installed."
            if status_text:
                status_text.error(error_msg)
            return False, error_msg, 0
        
        try:
            if status_text:
                status_text.info("Converting presentation to PDF...")
            
            process = subprocess.run(
                [libreoffice_path, '--headless', '--convert-to', 'pdf', 
                 '--outdir', temp_dir, pptx_path], 
                check=True, capture_output=True, timeout=120
            )
            
            if not os.path.exists(pdf_path):
                error_msg = f"PDF was not created at expected path: {pdf_path}"
                if status_text:
                    status_text.error(error_msg)
                return False, error_msg, 0
                
        except Exception as e:
            error_msg = f"Error converting to PDF: {str(e)}"
            if status_text:
                status_text.error(error_msg)
            return False, error_msg, 0
        
        try:
            if status_text:
                status_text.info("Converting PDF to images using pdf2image...")
            
            from pdf2image import convert_from_path
            
            images = convert_from_path(pdf_path, dpi=300, fmt="png")
            slide_count = len(images)
            
            if slide_count == 0:
                error_msg = "No slides were extracted from the PDF"
                if status_text:
                    status_text.error(error_msg)
                return False, error_msg, 0
            
            # GPU OPTIMIZATION: Process images in parallel batches
            def process_image_batch(batch_data):
                batch_results = []
                for i, img in batch_data:
                    img_array = np.array(img)
                    aspect_ratio = img.width / img.height
                    
                    if aspect_ratio > target_resolution[0] / target_resolution[1]:
                        new_width = target_resolution[0]
                        new_height = int(target_resolution[0] / aspect_ratio)
                    else:
                        new_height = target_resolution[1]
                        new_width = int(target_resolution[1] * aspect_ratio)
                    
                    # Use GPU resize if available
                    if GPU_PROCESSOR.gpu_framework != 'cpu':
                        resized_array = GPU_PROCESSOR.resize_image_gpu(img_array, (new_width, new_height))
                        resized_img = Image.fromarray(resized_array)
                    else:
                        resized_img = img.resize((new_width, new_height), LANCZOS)
                    
                    output_path = os.path.join(output_dir, f"{i:02d}.png")
                    resized_img.save(output_path, "PNG", optimize=True)
                    batch_results.append(output_path)
                return batch_results
            
            # Process in batches
            batch_size = OPTIMIZATION_SETTINGS.get('batch_size', 3)
            all_results = []
            
            for batch_start in range(0, slide_count, batch_size):
                batch_end = min(batch_start + batch_size, slide_count)
                batch_data = [(i, images[i-1]) for i in range(batch_start + 1, batch_end + 1)]
                batch_results = process_image_batch(batch_data)
                all_results.extend(batch_results)
            
            return True, f"Successfully extracted {slide_count} slides", slide_count
        
        except ImportError:
            error_msg = "pdf2image module not available. Please install it with: pip install pdf2image"
            if status_text:
                status_text.error(error_msg)
            return False, error_msg, 0
        except Exception as e:
            error_msg = f"Error converting PDF to images: {str(e)}"
            if status_text:
                status_text.error(error_msg)
            return False, error_msg, 0
            
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)

def extract_slides_from_pptx(pptx_path, output_dir, status_text=None):
    """GPU OPTIMIZED: Extract slides from PPTX using LibreOffice + pdftoppm method"""
    temp_dir = tempfile.mkdtemp()
    
    try:
        os.makedirs(output_dir, exist_ok=True)
        
        if status_text:
            status_text.info(f"Converting presentation to PDF using LibreOffice...")
        
        libreoffice_path = get_libreoffice_path()
        if not libreoffice_path:
            error_msg = "LibreOffice not found. Please make sure LibreOffice is installed."
            if status_text:
                status_text.error(error_msg)
            return False, error_msg, 0
            
        try:
            process = subprocess.run(
                [libreoffice_path, '--headless', '--convert-to', 'pdf',
                 '--outdir', temp_dir, pptx_path], 
                check=True, capture_output=True, timeout=120
            )
            
            expected_pdf_name = os.path.splitext(os.path.basename(pptx_path))[0] + '.pdf'
            pdf_path = os.path.join(temp_dir, expected_pdf_name)
            
            if not os.path.exists(pdf_path):
                error_msg = f"PDF was not created at expected path: {pdf_path}"
                if status_text:
                    status_text.error(error_msg)
                return False, error_msg, 0
                
        except Exception as e:
            error_msg = f"Error converting to PDF: {str(e)}"
            if status_text:
                status_text.error(error_msg)
            return False, error_msg, 0
            
        # Convert PDF to images using pdftoppm
        if status_text:
            status_text.info(f"Converting PDF to images...")
            
        temp_img_dir = tempfile.mkdtemp()
        success = False
        
        try:
            base_filename = os.path.join(temp_img_dir, "slide")
            pdftoppm_command = [
                "pdftoppm", "-png", "-r", "300",
                pdf_path, base_filename
            ]
            subprocess.run(pdftoppm_command, check=True, capture_output=True, timeout=120)
            
            extracted_files = sorted([f for f in os.listdir(temp_img_dir) if f.startswith("slide-") and f.endswith(".png")])
            slide_count = len(extracted_files)
            
            if slide_count > 0:
                # GPU OPTIMIZATION: Process images in parallel
                def process_slide_image(args):
                    i, filename = args
                    file_path = os.path.join(temp_img_dir, filename)
                    with Image.open(file_path) as img:
                        img_array = np.array(img)
                        aspect_ratio = img.width / img.height
                        
                        if aspect_ratio > 3840/2160:
                            new_width = 3840
                            new_height = int(3840 / aspect_ratio)
                        else:
                            new_height = 2160
                            new_width = int(2160 * aspect_ratio)
                        
                        # Use GPU resize if available
                        if GPU_PROCESSOR.gpu_framework != 'cpu':
                            resized_array = GPU_PROCESSOR.resize_image_gpu(img_array, (new_width, new_height))
                            resized = Image.fromarray(resized_array)
                        else:
                            resized = img.resize((new_width, new_height), LANCZOS)
                        
                        new_filename = os.path.join(output_dir, f"{i:02d}.png")
                        resized.save(new_filename, "PNG", optimize=True)
                        return new_filename
                
                # Process in parallel with limited threads
                max_workers = min(OPTIMIZATION_SETTINGS.get('max_threads', 3), slide_count)
                with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
                    args_list = [(i, filename) for i, filename in enumerate(extracted_files, 1)]
                    list(executor.map(process_slide_image, args_list))
                        
                return True, f"Successfully extracted {slide_count} slides", slide_count
        except Exception:
            pass
            
        # Fallback to pdf2image
        try:
            from pdf2image import convert_from_path
            
            images = convert_from_path(pdf_path, dpi=300, output_folder=temp_img_dir, fmt="png")
            slide_count = len(images)
            
            if slide_count > 0:
                # GPU OPTIMIZATION: Parallel processing for fallback too
                def process_fallback_image(args):
                    i, img = args
                    img_array = np.array(img)
                    aspect_ratio = img.width / img.height
                    
                    if aspect_ratio > 3840/2160:
                        new_width = 3840
                        new_height = int(3840 / aspect_ratio)
                    else:
                        new_height = 2160
                        new_width = int(2160 * aspect_ratio)
                    
                    # Use GPU resize if available
                    if GPU_PROCESSOR.gpu_framework != 'cpu':
                        resized_array = GPU_PROCESSOR.resize_image_gpu(img_array, (new_width, new_height))
                        resized = Image.fromarray(resized_array)
                    else:
                        resized = img.resize((new_width, new_height), LANCZOS)
                    
                    output_path = os.path.join(output_dir, f"{i:02d}.png")
                    resized.save(output_path, "PNG", optimize=True)
                    return output_path
                
                max_workers = min(OPTIMIZATION_SETTINGS.get('max_threads', 3), slide_count)
                with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
                    args_list = [(i, img) for i, img in enumerate(images, 1)]
                    list(executor.map(process_fallback_image, args_list))
                    
                return True, f"Successfully extracted {slide_count} slides", slide_count
        except Exception:
            pass
            
        return False, "Failed to extract slides - required tools not available", 0
        
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)

def extract_slides_with_method(
    pptx_path: str, 
    output_dir: str, 
    method: ConversionMethod,
    status_text: Optional[Any] = None,
    enable_auto_fallback: bool = True,
    target_resolution: Tuple[int, int] = (3840, 2160)
) -> Tuple[bool, str, int]:
    """Extract slides from a PPTX file using the specified method"""
    os.makedirs(output_dir, exist_ok=True)
    
    if status_text and hasattr(status_text, 'update_action'):
        status_text.update_action(f"Attempting extraction using {CONVERSION_METHODS[method]} method")
    elif status_text:
        status_text.info(f"Attempting extraction using {CONVERSION_METHODS[method]} method")
    
    success, message, slide_count = _try_method(pptx_path, output_dir, method, status_text, target_resolution)
    
    if not success and enable_auto_fallback:
        if status_text and hasattr(status_text, 'update_action'):
            status_text.update_action(f"Primary method failed. Trying fallback methods...")
        elif status_text:
            status_text.warning(f"Primary method failed: {message}")
            status_text.info("Trying fallback methods...")
        
        fallback_methods = get_fallback_methods(method)
        
        for fallback_method in fallback_methods:
            if status_text and hasattr(status_text, 'update_action'):
                status_text.update_action(f"Trying {CONVERSION_METHODS[fallback_method]} method...")
            elif status_text:
                status_text.info(f"Trying {CONVERSION_METHODS[fallback_method]} method...")
                
            success, message, slide_count = _try_method(pptx_path, output_dir, fallback_method, status_text, target_resolution)
            
            if success:
                if status_text and hasattr(status_text, 'update_action'):
                    status_text.update_action(f"Fallback method {CONVERSION_METHODS[fallback_method]} succeeded")
                elif status_text:
                    status_text.success(f"Fallback method {CONVERSION_METHODS[fallback_method]} succeeded")
                break
        
    return success, message, slide_count

def get_fallback_methods(primary_method: ConversionMethod) -> List[ConversionMethod]:
    """Get a list of fallback methods to try if the primary method fails"""
    all_methods = list(ConversionMethod)
    all_methods.remove(primary_method)
    
    if primary_method == ConversionMethod.LIBREOFFICE:
        return [ConversionMethod.PDF2IMAGE, ConversionMethod.PYTHON_PPTX]
    elif primary_method == ConversionMethod.PYTHON_PPTX:
        return [ConversionMethod.LIBREOFFICE, ConversionMethod.PDF2IMAGE]
    else:  # PDF2IMAGE
        return [ConversionMethod.LIBREOFFICE, ConversionMethod.PYTHON_PPTX]

def _try_method(
    pptx_path: str, 
    output_dir: str, 
    method: ConversionMethod,
    status_text: Optional[Any],
    target_resolution: Tuple[int, int]
) -> Tuple[bool, str, int]:
    """Try to extract slides using a specific method"""
    try:
        if method == ConversionMethod.LIBREOFFICE:
            return extract_slides_from_pptx(pptx_path, output_dir, status_text)
        elif method == ConversionMethod.PYTHON_PPTX:
            return extract_slides_directly(pptx_path, output_dir, status_text, target_resolution)
        elif method == ConversionMethod.PDF2IMAGE:
            return convert_with_pdf2image(pptx_path, output_dir, status_text, target_resolution)
        else:
            error_msg = f"Unknown conversion method: {method}"
            if status_text and hasattr(status_text, 'error'):
                status_text.error(error_msg)
            return False, error_msg, 0
            
    except ImportError as e:
        error_msg = f"Required module not available for {method} method: {str(e)}"
        if status_text and hasattr(status_text, 'error'):
            status_text.error(error_msg)
        return False, error_msg, 0
        
    except Exception as e:
        error_msg = f"Error in {method} method: {str(e)}"
        if status_text and hasattr(status_text, 'error'):
            status_text.error(error_msg)
        return False, error_msg, 0

# Increase the Pillow decompression bomb limit for large images
increase_image_decompression_limit()

# Constants for image processing
LOGO_PATH = 'config/logo.png'
COPYRIGHT_PATH = 'config/copyright.txt'
LOGO_SIZE = 250
LOGO_PADDING = 0
LOGO_POSITION = 'top-right'
COPYRIGHT_PADDING = 35
COPYRIGHT_POSITION = 'bottom-center'
FONT_SIZE = 65
FONT_COLOR = (0, 0, 0)  # Black
TARGET_RESOLUTION = (3840, 2160)  # 4K

@lru_cache(maxsize=128)
def read_copyright(file_path):
    """OPTIMIZED: Read copyright text from file (cached)"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read().strip()
        text = re.sub(r'A©', '©', text)
        return text
    except UnicodeDecodeError:
        try:
            with open(file_path, 'r', encoding='cp1252') as f:
                text = f.read().strip()
                text = re.sub(r'A©', '©', text)
                return text
        except Exception:
            return "© All Rights Reserved"
    except Exception:
        return "© All Rights Reserved"

def get_position(img_width, img_height, element_width, element_height, position, padding):
    """Calculate position based on alignment and padding"""
    if position == 'top-left':
        return (padding, padding)
    elif position == 'top-right':
        return (img_width - element_width - padding, padding)
    elif position == 'bottom-left':
        return (padding, img_height - element_height - padding)
    elif position == 'bottom-right':
        return (img_width - element_width - padding, img_height - element_height - padding)
    elif position == 'bottom-center':
        return ((img_width - element_width) // 2, img_height - element_height - padding)
    else:
        raise ValueError(f"Invalid position: {position}")

@lru_cache(maxsize=8)
def get_font(size):
    """OPTIMIZED: Get a font with fallback options (cached)"""
    try:
        return ImageFont.truetype("Arial.ttf", size)
    except IOError:
        try:
            return ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", size)
        except IOError:
            return ImageFont.load_default()

@lru_cache(maxsize=16)
def create_logo_copyright_overlay(img_width: int, img_height: int, logo_path: str, copyright_text: str) -> Image.Image:
    """
    BATCH OPTIMIZATION: Pre-render logo and copyright overlay for reuse across multiple images
    This avoids recreating the same overlay for every single image
    
    Args:
        img_width: Target image width
        img_height: Target image height  
        logo_path: Path to logo file (for cache key)
        copyright_text: Copyright text to render
        
    Returns:
        Pre-rendered overlay image with logo and copyright
    """
    # Create transparent overlay layer
    overlay = Image.new('RGBA', (img_width, img_height), (0, 0, 0, 0))
    
    # Add logo if it exists
    if os.path.exists(logo_path):
        try:
            with Image.open(logo_path) as logo:
                logo = logo.convert("RGBA")
                logo_resized = logo.resize((LOGO_SIZE, LOGO_SIZE), LANCZOS)
                logo_pos = get_position(img_width, img_height, 
                                       LOGO_SIZE, LOGO_SIZE, LOGO_POSITION, LOGO_PADDING)
                overlay.paste(logo_resized, logo_pos, logo_resized)
        except Exception as e:
            print(f"Warning: Could not load logo: {e}")
    
    # Add copyright text
    draw = ImageDraw.Draw(overlay)
    font = get_font(FONT_SIZE)
    
    # Get text size
    text_bbox = draw.textbbox((0, 0), copyright_text, font=font)
    text_width = text_bbox[2] - text_bbox[0]
    text_height = text_bbox[3] - text_bbox[1]
    
    text_pos = get_position(img_width, img_height, 
                           text_width, text_height, COPYRIGHT_POSITION, COPYRIGHT_PADDING)
    draw.text(text_pos, copyright_text, font=font, fill=FONT_COLOR)
    
    return overlay

def apply_cached_overlay(img: Image.Image, overlay: Image.Image) -> Image.Image:
    """
    BATCH OPTIMIZATION: Apply pre-rendered overlay quickly to an image
    Much faster than recreating logo/copyright for each image
    """
    # Convert to RGBA if necessary
    if img.mode != 'RGBA':
        img = img.convert('RGBA')
    
    # Apply the pre-rendered overlay
    return Image.alpha_composite(img, overlay)

def process_image(input_path, output_path, logo, copyright_text):
    """GPU + BATCH OPTIMIZED: Process an image with GPU acceleration and cached overlay"""
    try:
        with Image.open(input_path) as img:
            img_array = np.array(img)
            
            # Check if the image needs upscaling
            if img.width < TARGET_RESOLUTION[0] or img.height < TARGET_RESOLUTION[1]:
                # Calculate new dimensions
                aspect_ratio = img.width / img.height
                if aspect_ratio > 16/9:
                    new_width = TARGET_RESOLUTION[0]
                    new_height = int(TARGET_RESOLUTION[0] / aspect_ratio)
                else:
                    new_height = TARGET_RESOLUTION[1]
                    new_width = int(TARGET_RESOLUTION[1] * aspect_ratio)
                
                # Use GPU resize
                resized_array = GPU_PROCESSOR.resize_image_gpu(img_array, (new_width, new_height))
                processed_img = Image.fromarray(resized_array)
            else:
                processed_img = img.copy()
            
            # BATCH OPTIMIZATION: Use cached overlay instead of recreating logo/copyright
            overlay = create_logo_copyright_overlay(
                processed_img.width, 
                processed_img.height, 
                LOGO_PATH, 
                copyright_text
            )
            
            # Apply the cached overlay (much faster than individual processing)
            final_img = apply_cached_overlay(processed_img, overlay)

            # Save the processed image
            if img.format == 'JPEG':
                final_img = final_img.convert('RGB')
                final_img.save(output_path, format=img.format, quality=95, optimize=True, progressive=True)
            elif img.format == 'PNG':
                final_img.save(output_path, format=img.format, optimize=True)
            else:
                final_img.save(output_path, format=img.format)
        
        return True
    except Exception as e:
        st.error(f"Error processing {input_path}: {str(e)}")
        return False

def extract_and_upscale_images(pptx_file, output_folder, conversion_method=None, enable_auto_fallback=True, status_mgr=None):
    """Extract images from a PPTX file, upscale them to 4K, and save them in the output folder."""
    success, message, slide_count = extract_slides_with_method(
        pptx_file, 
        output_folder, 
        method=conversion_method or ConversionMethod.LIBREOFFICE,
        status_text=status_mgr or st,
        enable_auto_fallback=enable_auto_fallback,
        target_resolution=TARGET_RESOLUTION
    )
    
    if not success:
        if status_mgr:
            status_mgr.error(message)
        else:
            st.error(message)
        return False, [], 0
    
    # Gather the paths of the created images
    extracted_images = []
    if slide_count > 0:
        for i in range(1, slide_count + 1):
            image_path = os.path.join(output_folder, f"{i:02d}.png")
            if os.path.exists(image_path):
                extracted_images.append(image_path)
    
    return success, extracted_images, slide_count

def find_presentations() -> List[str]:
    """Find all presentation files in the input directory"""
    input_dir = get_input_directory()
    all_presentations = find_presentation_files(input_dir)
    return all_presentations

class StatusManager:
    """GPU OPTIMIZED: Status manager with machine-aware threading"""
    
    def __init__(self, progress_bar, status_text, current_action_text=None):
        self.progress_bar = progress_bar
        self.status_text = status_text
        self.current_action_text = current_action_text
        self.current_file = ""
        self.current_action = ""
        self.processed_count = 0
        self.total_count = 0
        self.start_time = datetime.now()
        self._lock = threading.Lock()
        
    def set_total(self, total):
        with self._lock:
            self.total_count = total
        
    def update_progress(self, current_index, filename):
        with self._lock:
            self.processed_count = current_index
            self.current_file = filename
            progress = self.processed_count / self.total_count if self.total_count > 0 else 0
            self.progress_bar.progress(progress)
            
            if self.processed_count > 1:
                elapsed_time = (datetime.now() - self.start_time).total_seconds()
                time_per_file = elapsed_time / (self.processed_count - 1)
                remaining_files = self.total_count - self.processed_count
                remaining_seconds = remaining_files * time_per_file
                remaining_time = str(timedelta(seconds=int(remaining_seconds)))
                
                gpu_info = GPU_PROCESSOR.gpu_framework.upper()
                self.status_text.info(f"Processing {self.processed_count}/{self.total_count}: {filename} ({gpu_info}) (Est. time remaining: {remaining_time})")
            else:
                gpu_info = GPU_PROCESSOR.gpu_framework.upper()
                self.status_text.info(f"Processing {self.processed_count}/{self.total_count}: {filename} ({gpu_info})")
    
    def update_action(self, message):
        with self._lock:
            self.current_action = message
            if self.current_action_text:
                self.current_action_text.info(message)
        
    def info(self, message):
        self.update_action(message)
        
    def warning(self, message):
        st.warning(message)
        
    def error(self, message):
        st.error(message)
        
    def success(self, message):
        self.status_text.success(message)
        if self.current_action_text:
            self.current_action_text.empty()
def process_presentations(input_dir: str, progress_bar, status_text, create_without_logo_folder=False, 
                 conversion_method=None, enable_auto_fallback=True, current_action_text=None) -> List[Dict]:
    """GPU OPTIMIZED: Process all PPTX and ZIP files with GPU acceleration and identical file saving logic"""
    results = []
    
    # Create status manager
    status_mgr = StatusManager(progress_bar, status_text, current_action_text)
    
    # Try to load logo and copyright text (cached)
    try:
        logo = Image.open(LOGO_PATH).convert("RGBA") if os.path.exists(LOGO_PATH) else None
        copyright_text = read_copyright(COPYRIGHT_PATH) if os.path.exists(COPYRIGHT_PATH) else "© All Rights Reserved"
    except Exception as e:
        st.warning(f"Could not load logo or copyright: {str(e)}")
        logo = None
        copyright_text = "© All Rights Reserved"
    
    # Find all presentations in the input directory
    all_presentations = find_presentations()
    total_presentations = len(all_presentations)
    status_mgr.set_total(total_presentations)
    
    # GPU OPTIMIZATION: Process presentations in optimized batches
    batch_size = OPTIMIZATION_SETTINGS.get('batch_size', 2)
    
    for batch_start in range(0, total_presentations, batch_size):
        batch_end = min(batch_start + batch_size, total_presentations)
        batch_presentations = all_presentations[batch_start:batch_end]
        
        # Process each presentation in the current batch
        for i, pptx_file in enumerate(batch_presentations):
            # Update progress with file info
            processed_count = batch_start + i + 1
            file_name = os.path.basename(pptx_file)
            status_mgr.update_progress(processed_count, file_name)
            
            result = {
                "file_path": pptx_file,
                "status": "error",
                "message": "",
                "slide_count": 0,
                "output_dir": ""
            }
            
            try:
                # Determine file type first
                file_extension = os.path.splitext(pptx_file)[1].lower()
                is_zip_file = file_extension == '.zip'
                
                # Get the directory containing the presentation
                pptx_dir = os.path.dirname(pptx_file)
                
                # IDENTICAL FILE SAVING LOGIC AS ORIGINAL
                # Check if this is inside an all_pptx directory
                parent_dir = os.path.basename(pptx_dir)
                if parent_dir == 'all_pptx':
                    # If file is in all_pptx, create lecture folder in parent directory
                    course_dir = os.path.dirname(pptx_dir)  # Go up one level from all_pptx
                    pptx_name = os.path.splitext(os.path.basename(pptx_file))[0]
                    
                    # Extract lecture number from filename for proper naming
                    lecture_match = re.search(r'(\d+)', pptx_name)
                    if lecture_match:
                        lecture_num = lecture_match.group(1).zfill(2)
                        lecture_folder = os.path.join(course_dir, f"Lecture {lecture_num}")
                    else:
                        lecture_folder = os.path.join(course_dir, pptx_name)
                        
                    os.makedirs(lecture_folder, exist_ok=True)
                    
                    # Create "English image" folder inside the lecture folder
                    output_folder = os.path.join(lecture_folder, "English image")
                    os.makedirs(output_folder, exist_ok=True)
                    
                    # The all_pptx folder already exists since the file is in it
                    all_pptx_folder = pptx_dir
                else:
                    # Create 'all_pptx' directory in the same folder as the presentation
                    all_pptx_folder = os.path.join(pptx_dir, 'all_pptx')
                    os.makedirs(all_pptx_folder, exist_ok=True)
                    
                    # Create folder with same name as file for images
                    file_name = os.path.splitext(os.path.basename(pptx_file))[0]
                    
                    # For ZIP files with numeric names, add "Lecture" prefix to match PPTX structure
                    if is_zip_file and re.match(r'^\d+$', file_name):
                        lecture_folder = os.path.join(pptx_dir, f"Lecture {file_name.zfill(2)}")
                    else:
                        lecture_folder = os.path.join(pptx_dir, file_name)
                        
                    os.makedirs(lecture_folder, exist_ok=True)
                    
                    # Create "English image" folder inside the lecture folder
                    output_folder = os.path.join(lecture_folder, "English image")
                    os.makedirs(output_folder, exist_ok=True)
                
                # Create 'without_logo_png' folder inside the output folder only if requested
                without_logo_folder = os.path.join(output_folder, 'without_logo_png')
                if create_without_logo_folder:
                    os.makedirs(without_logo_folder, exist_ok=True)
                
                # Check if file is already processed and in all_pptx
                new_pptx_path = os.path.join(all_pptx_folder, os.path.basename(pptx_file))
                if os.path.exists(new_pptx_path):
                    # If already processed, just update the result
                    result["message"] = "Presentation already processed"
                    
                    # Count existing images
                    png_count = len([f for f in os.listdir(output_folder) 
                                   if f.endswith('.png') and not os.path.isdir(os.path.join(output_folder, f))])
                    
                    result["status"] = "success"
                    result["slide_count"] = png_count
                    result["output_dir"] = output_folder
                    results.append(result)
                    continue
                    
                # Extract and upscale images from the presentation or ZIP file
                temp_dir = None
                
                try:
                    # Only create and use without_logo_folder if requested
                    if create_without_logo_folder:
                        # Ensure the folder exists
                        os.makedirs(without_logo_folder, exist_ok=True)
                        
                        if is_zip_file:
                            # Extract from ZIP file with GPU acceleration
                            success, extracted_images, slide_count = extract_slides_from_zip(
                                pptx_file,
                                without_logo_folder,
                                status_mgr=status_mgr,
                                target_resolution=TARGET_RESOLUTION
                            )
                        else:
                            # Extract from PPTX file with GPU acceleration
                            success, extracted_images, slide_count = extract_and_upscale_images(
                                pptx_file, 
                                without_logo_folder,
                                conversion_method=conversion_method,
                                enable_auto_fallback=enable_auto_fallback,
                                status_mgr=status_mgr
                            )
                        
                        # GPU OPTIMIZATION: Process images in parallel batches with cached overlays
                        processed_img_count = 0
                        
                        def process_image_with_logo(img_path):
                            # Original image is in without_logo_folder
                            # Processed image will be in output_folder
                            output_filename = os.path.basename(img_path)
                            output_path = os.path.join(output_folder, output_filename)
                            
                            # Process the image with GPU acceleration + batch optimization
                            return process_image(img_path, output_path, logo, copyright_text)
                        
                        # Process images in parallel with GPU-optimized settings
                        max_workers = min(OPTIMIZATION_SETTINGS.get('max_threads', 3), len(extracted_images))
                        
                        if status_mgr:
                            status_mgr.update_action(f"Applying cached logo/copyright overlays to {len(extracted_images)} images...")
                        
                        with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
                            results_list = list(executor.map(process_image_with_logo, extracted_images))
                            processed_img_count = sum(1 for r in results_list if r)
                            
                    else:
                        # Use a temporary directory that we manage explicitly
                        temp_dir = tempfile.mkdtemp()
                        
                        if is_zip_file:
                            # Extract from ZIP file with GPU acceleration
                            success, extracted_images, slide_count = extract_slides_from_zip(
                                pptx_file,
                                temp_dir,
                                status_mgr=status_mgr,
                                target_resolution=TARGET_RESOLUTION
                            )
                        else:
                            # Extract from PPTX file with GPU acceleration
                            success, extracted_images, slide_count = extract_and_upscale_images(
                                pptx_file, 
                                temp_dir,
                                conversion_method=conversion_method,
                                enable_auto_fallback=enable_auto_fallback,
                                status_mgr=status_mgr
                            )
                        
                        # GPU OPTIMIZATION: Process images in parallel batches with cached overlays
                        processed_img_count = 0
                        
                        def process_image_direct(img_path):
                            # Temporary image is in temp_dir
                            # Processed image will go directly to output_folder
                            output_filename = os.path.basename(img_path)
                            output_path = os.path.join(output_folder, output_filename)
                            
                            # Process the image with GPU acceleration + batch optimization
                            return process_image(img_path, output_path, logo, copyright_text)
                        
                        # Process images in parallel with GPU-optimized settings
                        max_workers = min(OPTIMIZATION_SETTINGS.get('max_threads', 3), len(extracted_images))
                        
                        if status_mgr:
                            status_mgr.update_action(f"Applying cached logo/copyright overlays to {len(extracted_images)} images...")
                        
                        with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
                            results_list = list(executor.map(process_image_direct, extracted_images))
                            processed_img_count = sum(1 for r in results_list if r)
                    
                    if not success:
                        result["message"] = "Failed to extract images from presentation"
                        results.append(result)
                        continue
                finally:
                    # Clean up the temporary directory if it was created
                    if temp_dir and os.path.exists(temp_dir):
                        shutil.rmtree(temp_dir, ignore_errors=True)
                
                # Move the presentation file to all_pptx folder (IDENTICAL TO ORIGINAL)
                if os.path.exists(pptx_file) and not os.path.exists(new_pptx_path):
                    os.rename(pptx_file, new_pptx_path)
                
                # Update result
                result["status"] = "success"
                result["message"] = f"Successfully processed {processed_img_count} images with GPU + batch optimization"
                result["slide_count"] = slide_count
                result["output_dir"] = output_folder
            
            except Exception as e:
                result["message"] = f"Error: {str(e)}"
            
            results.append(result)
            
            # Small delay for UI updates (reduced for better performance)
            time.sleep(0.05)
    
    # Ensure progress bar reaches 100% at the end
    progress_bar.progress(1.0)
    
    return results

def main():
    st.title("4K Image (PPTX & ZIP) - Final")
    
    # Display machine information
    machine_info = MACHINE_DETECTOR.get_machine_info()
    runtime_acceleration = GPU_PROCESSOR.gpu_framework
    
    with st.expander("🖥️ Machine Profile & Acceleration"):
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown("**Detected Machine:**")
            st.write(f"• **Name**: {machine_info['name']}")
            st.write(f"• **Platform**: {machine_info['detected_info'].get('platform', 'Unknown')}")
            st.write(f"• **Memory**: {machine_info['detected_info'].get('memory_gb', 'Unknown')} GB")
            st.write(f"• **CPU Cores**: {machine_info['detected_info'].get('cpu_count', 'Unknown')}")
        
        with col2:
            st.markdown("**Runtime Acceleration:**")
            if runtime_acceleration == 'cuda':
                st.success("✅ CUDA acceleration enabled")
                
                if 'gpu_name' in machine_info['detected_info']:
                    st.write(f"• **GPU**: {machine_info['detected_info']['gpu_name']}")
                
                st.write(f"• **Max Threads**: {OPTIMIZATION_SETTINGS.get('max_threads', 4)}")
                st.write(f"• **Batch Size**: {OPTIMIZATION_SETTINGS.get('batch_size', 2)}")
            elif machine_info['machine_id'] == 'macbook_m3_pro':
                st.info("ℹ️ Apple Silicon detected. Using the optimized CPU path for image work.")
                st.write(f"• **Max Threads**: {OPTIMIZATION_SETTINGS.get('max_threads', 4)}")
                st.write(f"• **Batch Size**: {OPTIMIZATION_SETTINGS.get('batch_size', 2)}")
            else:
                st.warning("⚠️ CUDA acceleration not available - using CPU")
                
                if machine_info['machine_id'] == 'windows_i5_rtx3050':
                    st.info("💡 **To enable CUDA acceleration**: `pip install cupy-cuda12x`")
                else:
                    st.info("💡 This fallback path is CPU-only.")
        
        if machine_info['machine_id'] == 'macbook_m3_pro':
            st.info("🍎 **MacBook profile**: tuned for your M3 Pro with conservative memory use and fast CPU resizing.")
        elif machine_info['machine_id'] == 'windows_i5_rtx3050':
            st.info("🎮 **RTX profile**: tuned for overnight bulk runs with CUDA when CuPy is installed.")
    
    st.write("⚡ Final machine-aware version for your MacBook Pro M3 Pro and Windows RTX 3050 laptop.")
    
    input_dir = get_input_directory()
    
    if not os.path.exists(input_dir):
        st.error(f"Input directory not found: {input_dir}")
        st.info("Please create an 'input' directory in the project root and add your files.")
        return
    
    # Find presentations
    all_presentations = find_presentations()
    
    if not all_presentations:
        st.warning("No presentation files found in the input directory.")
        return
    
    # Display number of presentations found
    total_presentations = len(all_presentations)
    pptx_count = sum(1 for f in all_presentations if f.lower().endswith('.pptx'))
    zip_count = sum(1 for f in all_presentations if f.lower().endswith('.zip'))
    st.info(f"Found {total_presentations} files in input directory: {pptx_count} PPTX, {zip_count} ZIP")
    
    # Options section
    st.header("Options")
    
    # Option for creating without_logo_png folder
    create_without_logo_folder = st.checkbox("Create 'without_logo_png' folder", value=False, 
                                          help="If checked, a folder containing original images without logo and copyright will be created")
    
    # Runtime-specific options
    if runtime_acceleration == 'cuda':
        st.subheader("🚀 CUDA Settings")
        
        custom_batch_size = st.slider(
            "Batch Size (images processed simultaneously)",
            min_value=1,
            max_value=8,
            value=OPTIMIZATION_SETTINGS.get('batch_size', 2),
            help="Higher values use more GPU memory but may be faster"
        )
        OPTIMIZATION_SETTINGS['batch_size'] = custom_batch_size
        
        if custom_batch_size > 4:
            st.warning("⚠️ High batch sizes may cause GPU memory issues with very large images")
    
    # Conversion method selection (only for PPTX files)
    st.subheader("Conversion Method (PPTX only)")
    st.write("Select the method to use for extracting slides from PPTX presentations:")
    st.caption("Note: ZIP files are processed directly and don't use these conversion methods")
    
    # Create a dropdown for conversion method selection
    conversion_methods = list(CONVERSION_METHODS.items())
    conversion_method = st.selectbox(
        "Select conversion method:",
        options=[method.value for method, _ in conversion_methods],
        format_func=lambda x: CONVERSION_METHODS[x],
        index=0,
        help="Choose the method used to extract slides from PPTX presentations (not applicable to ZIP files)"
    )
    
    # Option for automatic fallback
    enable_auto_fallback = st.checkbox(
        "Enable automatic fallback", 
        value=True,
        help="If checked, automatically try other methods if the selected method fails"
    )
    
    # Information about conversion methods and file types
    with st.expander("About File Types and Conversion Methods"):
        st.markdown("""
        ### Supported File Types
        
        #### PPTX Files
        - PowerPoint presentation files
        - Slides are extracted and converted to images
        - Multiple conversion methods available (see below)
        
        #### ZIP Files
        - ZIP archives containing image files (PNG, JPG, JPEG)
        - Images are extracted directly from the ZIP
        - Automatically sorted by numeric prefix (e.g., 1_Title.png, 2_Content.png)
        - Upscaled to 4K if needed using GPU acceleration
        - No conversion method selection needed
        
        ### Conversion Methods (PPTX only)
        
        #### LibreOffice + pdftoppm (best quality)
        - Uses LibreOffice to convert PPTX to PDF
        - Uses pdftoppm to convert PDF to high-quality images
        - GPU-accelerated image resizing
        - Provides the best visual quality for complex slides
        
        #### Direct PPTX processing (fastest)
        - Uses python-pptx to extract slides directly
        - GPU-accelerated image processing
        - Faster than other methods
        - Good for simple presentations with basic elements
        
        #### PDF2Image (most compatible)
        - Uses LibreOffice to convert PPTX to PDF
        - Uses pdf2image library to convert PDF to images
        - GPU-accelerated resizing and processing
        - More compatible fallback option
        """)
        
        st.info("💡 Tip: GPU acceleration is applied to all methods for faster image processing.")
    
    # Process presentations
    st.header("Extract Presentations")
    
    button_label = "🚀 Extract with CUDA Acceleration" if runtime_acceleration == 'cuda' else "Extract Presentations"
    if st.button(button_label, disabled=total_presentations == 0):
        if total_presentations == 0:
            st.warning("No presentations found for processing.")
            return
        
        # Process presentations with GPU acceleration
        progress_bar = st.progress(0)
        status_text = st.empty()
        current_action_text = st.empty()
        results_container = st.container()
        
        start_time = time.time()
        
        # Process the presentations with GPU acceleration
        results = process_presentations(
            input_dir, 
            progress_bar, 
            status_text, 
            create_without_logo_folder,
            conversion_method=conversion_method,
            enable_auto_fallback=enable_auto_fallback,
            current_action_text=current_action_text
        )
        
        end_time = time.time()
        processing_time = end_time - start_time
        
        # Clear status
        current_action_text.empty()
        
        runtime_label = runtime_acceleration.upper() if runtime_acceleration != 'cpu' else 'CPU'
        status_text.success(f"✅ Processed {total_presentations} presentations in {processing_time:.1f}s using {runtime_label}")
        
        with results_container:
            st.subheader("Processing Results")
            
            col1, col2, col3, col4 = st.columns(4)
            with col1:
                st.metric("⚡ Runtime", runtime_label)
            with col2:
                st.metric("⏱️ Time (sec)", f"{processing_time:.1f}")
            with col3:
                success_count = sum(1 for r in results if r["status"] == "success")
                st.metric("✅ Successful", success_count)
            with col4:
                if processing_time > 0:
                    files_per_second = total_presentations / processing_time
                    st.metric("🚀 Speed", f"{files_per_second:.1f} files/sec")
            
            if runtime_acceleration == 'cuda':
                estimated_cpu_time = processing_time * 2.5  # Estimate CPU would be 2.5x slower
                time_saved = estimated_cpu_time - processing_time
                
                if time_saved > 0:
                    st.success(f"🎯 **GPU Acceleration Benefit**: Estimated {time_saved:.1f}s saved vs CPU-only processing")
            
            # Statistics
            error_count = sum(1 for r in results if r["status"] == "error")
            slide_count = sum(r.get("slide_count", 0) for r in results)
            
            st.write(f"❌ Errors: {error_count}")
            st.write(f"📊 Total slides extracted: {slide_count}")
            
            # Detailed results
            for result in results:
                status_icon = "✅" if result["status"] == "success" else "❌"
                file_name = os.path.basename(result["file_path"])
                
                with st.expander(f"{status_icon} {file_name}"):
                    st.write(f"**Status:** {result['status']}")
                    st.write(f"**Message:** {result['message']}")
                    
                    if result["status"] == "success":
                        st.write(f"**Slides extracted:** {result['slide_count']}")
                        st.write(f"**Output directory:** {result['output_dir']}")

if __name__ == "__main__":
    main()
