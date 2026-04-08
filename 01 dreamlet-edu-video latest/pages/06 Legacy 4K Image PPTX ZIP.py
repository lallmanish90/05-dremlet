"""
CODING CONVENTION: NO SHARED CODE
- All code for this page must be contained entirely within this single file
- Never import from other page files or create shared utilities
- Copy any needed functions directly into this file
- Each page is completely self-contained and independent

STATUS: LEGACY
PURPOSE: Older combined PPTX and ZIP 4K image generation page.
MAIN INPUTS:
- PPTX presentations and ZIP image archives under `input/`
MAIN OUTPUTS:
- 4K image folders written into lecture working directories
REQUIRED CONFIG / ASSETS:
- `config/logo.png`
- `config/copyright.txt`
EXTERNAL SERVICES:
- local `libreoffice`
- `pdf2image` / Poppler for PDF conversion paths
HARDWARE ASSUMPTIONS:
- none
REPLACED BY:
- `pages/06 Generate 4K Images.py`
"""

import streamlit as st
import os
import re
import time
import cv2
import numpy as np
import io
import enum
import fnmatch
import platform
import tempfile
import glob
import shutil
import subprocess
import zipfile
from PIL import Image, ImageDraw, ImageFont
try:
    # For newer PIL versions
    LANCZOS = Image.Resampling.LANCZOS
except AttributeError:
    # For older PIL versions
    LANCZOS = Image.LANCZOS
from pptx import Presentation
from pathlib import Path
from typing import Dict, List, Tuple, Optional, Any
from datetime import datetime, timedelta

# Local utility functions (moved from multiple utils modules)
def increase_image_decompression_limit():
    """Increases the Pillow image decompression bomb limit to avoid errors with very large images"""
    new_limit = 3840 * 2160 * 50  # Approximately 415 million pixels
    Image.MAX_IMAGE_PIXELS = new_limit

def get_input_directory() -> str:
    """Get the path to the input directory"""
    input_dir = os.path.join(os.getcwd(), "input")
    return input_dir

def ensure_directory_exists(directory_path: str) -> None:
    """Create directory if it doesn't exist"""
    if not os.path.exists(directory_path):
        os.makedirs(directory_path)

def find_files(directory: str, pattern: str) -> List[str]:
    """Find all files matching a pattern in a directory (recursively)"""
    result = []
    for root, _, filenames in os.walk(directory):
        for filename in fnmatch.filter(filenames, pattern):
            result.append(os.path.join(root, filename))
    return result

def find_presentation_files(directory: str) -> List[str]:
    """Find all presentation files (PPTX and ZIP) in a directory"""
    pptx_files = find_files(directory, "*.pptx")
    zip_files = find_files(directory, "*.zip")
    return pptx_files + zip_files

def natural_sort_key(filename: str) -> List:
    """Generate a key for natural sorting (handles numbers correctly)"""
    # Extract numbers and text separately for proper sorting
    parts = re.split(r'(\d+)', filename)
    return [int(part) if part.isdigit() else part.lower() for part in parts]

def extract_numeric_prefix(filename: str) -> int:
    """Extract numeric prefix from filename like '1_Title.png' -> 1"""
    match = re.match(r'^(\d+)', filename)
    if match:
        return int(match.group(1))
    return 999999  # Put files without numeric prefix at the end

def extract_slides_from_zip(zip_file: str, output_folder: str, status_mgr=None, target_resolution=(3840, 2160)) -> Tuple[bool, List[str], int]:
    """
    Extract images from a ZIP file, upscale them to 4K, and save them WITHOUT logo/copyright.
    The main processing loop will add logo and copyright later.
    
    Args:
        zip_file: Path to the ZIP file
        output_folder: Directory to save the extracted images (without logo/copyright)
        status_mgr: Optional StatusManager to use for status updates
        target_resolution: Target resolution for upscaling (default: 4K)
        
    Returns:
        Tuple of (success, extracted_images, image_count)
    """
    temp_dir = tempfile.mkdtemp()
    
    try:
        os.makedirs(output_folder, exist_ok=True)
        
        if status_mgr:
            status_mgr.update_action(f"Extracting ZIP file: {os.path.basename(zip_file)}")
        
        # Extract ZIP contents
        try:
            with zipfile.ZipFile(zip_file, 'r') as zip_ref:
                zip_ref.extractall(temp_dir)
        except Exception as e:
            error_msg = f"Error extracting ZIP file: {str(e)}"
            if status_mgr:
                status_mgr.error(error_msg)
            return False, [], 0
        
        # Find all image files in the extracted contents
        image_extensions = ['*.png', '*.jpg', '*.jpeg', '*.PNG', '*.JPG', '*.JPEG']
        image_files = []
        for ext in image_extensions:
            image_files.extend(glob.glob(os.path.join(temp_dir, '**', ext), recursive=True))
        
        if not image_files:
            error_msg = "No image files found in ZIP archive"
            if status_mgr:
                status_mgr.error(error_msg)
            return False, [], 0
        
        # Sort images by numeric prefix to maintain slide order
        image_files.sort(key=lambda x: extract_numeric_prefix(os.path.basename(x)))
        
        if status_mgr:
            status_mgr.update_action(f"Found {len(image_files)} images in ZIP, processing...")
        
        # Process each image
        extracted_images = []
        for i, img_path in enumerate(image_files, 1):
            try:
                with Image.open(img_path) as img:
                    # Check if upscaling is needed
                    aspect_ratio = img.width / img.height
                    
                    if img.width >= target_resolution[0] and img.height >= target_resolution[1]:
                        # Image is already 4K or larger, just copy it
                        processed_img = img.copy()
                    else:
                        # Upscale to 4K while maintaining aspect ratio
                        if aspect_ratio > target_resolution[0] / target_resolution[1]:
                            new_width = target_resolution[0]
                            new_height = int(target_resolution[0] / aspect_ratio)
                        else:
                            new_height = target_resolution[1]
                            new_width = int(target_resolution[1] * aspect_ratio)
                        
                        processed_img = img.resize((new_width, new_height), LANCZOS)
                    
                    # Save with sequential naming
                    output_path = os.path.join(output_folder, f"{i:02d}.png")
                    processed_img.save(output_path, "PNG", optimize=True)
                    extracted_images.append(output_path)
                    
            except Exception as e:
                if status_mgr:
                    status_mgr.warning(f"Error processing image {os.path.basename(img_path)}: {str(e)}")
                continue
        
        image_count = len(extracted_images)
        
        if image_count == 0:
            error_msg = "Failed to process any images from ZIP"
            if status_mgr:
                status_mgr.error(error_msg)
            return False, [], 0
        
        if status_mgr:
            status_mgr.update_action(f"Successfully extracted {image_count} images from ZIP")
        
        return True, extracted_images, image_count
        
    finally:
        # Clean up temporary directory
        shutil.rmtree(temp_dir, ignore_errors=True)

def extract_course_lecture_section(file_path: str) -> Tuple[Optional[str], Optional[str], Optional[str]]:
    """Extract course, lecture, and section information from a file path"""
    dir_parts = os.path.normpath(file_path).split(os.sep)
    
    course = None
    lecture = None
    section = None
    
    for part in dir_parts:
        course_match = re.search(r'course\s*(\d+)', part.lower())
        if course_match:
            course = course_match.group(1)
            break
        number_start_match = re.match(r'^(\d+)\s+', part)
        if number_start_match:
            course = number_start_match.group(1)
            break
        bracket_match = re.search(r'[(\[]\s*(\d+)\s*[)\]]', part)
        if bracket_match:
            course = bracket_match.group(1)
            break
    
    filename = os.path.basename(file_path)
    lecture_patterns = [
        r'lecture\s*(\d+)',
        r'lec\s*(\d+)',
        r'^(\d+)[-\s]',
        r'^\w+\s*(\d+)'
    ]
    
    for pattern in lecture_patterns:
        lecture_match = re.search(pattern, filename.lower())
        if lecture_match:
            lecture = lecture_match.group(1)
            break
    
    if not lecture:
        for part in dir_parts:
            for pattern in lecture_patterns:
                lecture_match = re.search(pattern, part.lower())
                if lecture_match:
                    lecture = lecture_match.group(1)
                    break
            if lecture:
                break
    
    section_patterns = [
        r'section\s*(\d+)',
        r'sec\s*(\d+)'
    ]
    
    for part in dir_parts:
        for pattern in section_patterns:
            section_match = re.search(pattern, part.lower())
            if section_match:
                section = section_match.group(1)
                break
        if section:
            break
    
    return course, lecture, section

# Multi-method converter functions
class ConversionMethod(str, enum.Enum):
    LIBREOFFICE = "libreoffice"
    PYTHON_PPTX = "python-pptx"
    PDF2IMAGE = "pdf2image"

CONVERSION_METHODS = {
    ConversionMethod.LIBREOFFICE: "LibreOffice + pdftoppm (best quality)",
    ConversionMethod.PYTHON_PPTX: "Direct PPTX processing (fastest)",
    ConversionMethod.PDF2IMAGE: "PDF2Image (most compatible)"
}

def get_libreoffice_path():
    """Get the correct LibreOffice path based on the operating system"""
    if platform.system() == 'Darwin':  # macOS
        possible_paths = [
            '/Applications/LibreOffice.app/Contents/MacOS/soffice',
            '/Applications/LibreOffice.app/Contents/MacOS/soffice.bin',
            '/opt/homebrew/bin/soffice',
            '/usr/local/bin/soffice'
        ]
        for path in possible_paths:
            if os.path.exists(path):
                return path
        return None
    else:
        return 'soffice'

def extract_slides_directly(pptx_path, output_dir, status_text=None, target_resolution=(3840, 2160)):
    """Extract slides directly from PPTX file using python-pptx"""
    try:
        os.makedirs(output_dir, exist_ok=True)
        
        if status_text:
            status_text.info(f"Extracting slides directly from PPTX: {os.path.basename(pptx_path)}")
        
        prs = Presentation(pptx_path)
        slide_count = len(prs.slides)
        
        if slide_count == 0:
            if status_text:
                status_text.warning("Presentation contains no slides.")
            return False, "Presentation contains no slides", 0
        
        for i, slide in enumerate(prs.slides, 1):
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            slide_aspect_ratio = slide_width / slide_height
            
            if slide_aspect_ratio > target_resolution[0] / target_resolution[1]:
                width = target_resolution[0]
                height = int(width / slide_aspect_ratio)
            else:
                height = target_resolution[1]
                width = int(height * slide_aspect_ratio)
            
            img = np.ones((height, width, 3), np.uint8) * 255
            
            scale_x = width / slide_width
            scale_y = height / slide_height
            
            for shape in slide.shapes:
                if hasattr(shape, 'image') and shape.image:
                    try:
                        image_stream = io.BytesIO(shape.image.blob)
                        slide_image = Image.open(image_stream)
                        
                        slide_image_cv = cv2.cvtColor(np.array(slide_image), cv2.COLOR_RGB2BGR)
                        
                        left = int(shape.left * scale_x)
                        top = int(shape.top * scale_y)
                        img_width = int(shape.width * scale_x)
                        img_height = int(shape.height * scale_y)
                        
                        resized_image = cv2.resize(slide_image_cv, (img_width, img_height), interpolation=cv2.INTER_LANCZOS4)
                        
                        try:
                            if img_height > 0 and img_width > 0:
                                if top + img_height <= height and left + img_width <= width:
                                    img[top:top+img_height, left:left+img_width] = resized_image
                        except Exception as e:
                            if status_text:
                                status_text.warning(f"Error placing image on slide {i}: {str(e)}")
                    except Exception as e:
                        if status_text:
                            status_text.warning(f"Error processing image in slide {i}: {str(e)}")
            
            output_path = os.path.join(output_dir, f"{i:02d}.png")
            cv2.imwrite(output_path, img, [cv2.IMWRITE_PNG_COMPRESSION, 9])
        
        return True, f"Successfully extracted {slide_count} slides", slide_count
    
    except Exception as e:
        error_msg = f"Error in direct PPTX extraction: {str(e)}"
        if status_text:
            status_text.error(error_msg)
        return False, error_msg, 0

def convert_with_pdf2image(pptx_path, output_dir, status_text=None, target_resolution=(3840, 2160)):
    """Convert PPTX to images using LibreOffice for PPTX→PDF and pdf2image for PDF→PNG"""
    temp_dir = tempfile.mkdtemp()
    pdf_path = os.path.join(temp_dir, os.path.splitext(os.path.basename(pptx_path))[0] + '.pdf')
    
    try:
        os.makedirs(output_dir, exist_ok=True)
        
        if status_text:
            status_text.info(f"Converting presentation using pdf2image: {os.path.basename(pptx_path)}")
        
        libreoffice_path = get_libreoffice_path()
        if not libreoffice_path:
            error_msg = "LibreOffice not found. Please make sure LibreOffice is installed."
            if status_text:
                status_text.error(error_msg)
            return False, error_msg, 0
        
        try:
            if status_text:
                status_text.info("Converting presentation to PDF...")
            
            process = subprocess.run(
                [libreoffice_path, '--headless', '--convert-to', 'pdf', 
                 '--outdir', temp_dir, pptx_path], 
                check=True, capture_output=True, timeout=120
            )
            
            if not os.path.exists(pdf_path):
                error_msg = f"PDF was not created at expected path: {pdf_path}"
                if status_text:
                    status_text.error(error_msg)
                return False, error_msg, 0
                
        except Exception as e:
            error_msg = f"Error converting to PDF: {str(e)}"
            if status_text:
                status_text.error(error_msg)
            return False, error_msg, 0
        
        try:
            if status_text:
                status_text.info("Converting PDF to images using pdf2image...")
            
            from pdf2image import convert_from_path
            
            images = convert_from_path(pdf_path, dpi=300, fmt="png")
            slide_count = len(images)
            
            if slide_count == 0:
                error_msg = "No slides were extracted from the PDF"
                if status_text:
                    status_text.error(error_msg)
                return False, error_msg, 0
            
            for i, img in enumerate(images, 1):
                aspect_ratio = img.width / img.height
                
                if aspect_ratio > target_resolution[0] / target_resolution[1]:
                    new_width = target_resolution[0]
                    new_height = int(target_resolution[0] / aspect_ratio)
                else:
                    new_height = target_resolution[1]
                    new_width = int(target_resolution[1] * aspect_ratio)
                
                resized_img = img.resize((new_width, new_height), LANCZOS)
                
                output_path = os.path.join(output_dir, f"{i:02d}.png")
                resized_img.save(output_path, "PNG", optimize=True)
            
            return True, f"Successfully extracted {slide_count} slides", slide_count
        
        except ImportError:
            error_msg = "pdf2image module not available. Please install it with: pip install pdf2image"
            if status_text:
                status_text.error(error_msg)
            return False, error_msg, 0
        except Exception as e:
            error_msg = f"Error converting PDF to images: {str(e)}"
            if status_text:
                status_text.error(error_msg)
            return False, error_msg, 0
            
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)

def extract_slides_from_pptx(pptx_path, output_dir, status_text=None):
    """Extract slides from PPTX using LibreOffice + pdftoppm method"""
    temp_dir = tempfile.mkdtemp()
    
    try:
        os.makedirs(output_dir, exist_ok=True)
        
        if status_text:
            status_text.info(f"Converting presentation to PDF using LibreOffice...")
        
        libreoffice_path = get_libreoffice_path()
        if not libreoffice_path:
            error_msg = "LibreOffice not found. Please make sure LibreOffice is installed."
            if status_text:
                status_text.error(error_msg)
            return False, error_msg, 0
            
        try:
            process = subprocess.run(
                [libreoffice_path, '--headless', '--convert-to', 'pdf',
                 '--outdir', temp_dir, pptx_path], 
                check=True, capture_output=True, timeout=120
            )
            
            expected_pdf_name = os.path.splitext(os.path.basename(pptx_path))[0] + '.pdf'
            pdf_path = os.path.join(temp_dir, expected_pdf_name)
            
            if not os.path.exists(pdf_path):
                error_msg = f"PDF was not created at expected path: {pdf_path}"
                if status_text:
                    status_text.error(error_msg)
                return False, error_msg, 0
                
        except Exception as e:
            error_msg = f"Error converting to PDF: {str(e)}"
            if status_text:
                status_text.error(error_msg)
            return False, error_msg, 0
            
        # Convert PDF to images using pdftoppm
        if status_text:
            status_text.info(f"Converting PDF to images...")
            
        temp_img_dir = tempfile.mkdtemp()
        success = False
        
        try:
            base_filename = os.path.join(temp_img_dir, "slide")
            pdftoppm_command = [
                "pdftoppm", "-png", "-r", "300",
                pdf_path, base_filename
            ]
            subprocess.run(pdftoppm_command, check=True, capture_output=True, timeout=120)
            
            extracted_files = sorted([f for f in os.listdir(temp_img_dir) if f.startswith("slide-") and f.endswith(".png")])
            slide_count = len(extracted_files)
            
            if slide_count > 0:
                for i, filename in enumerate(extracted_files, 1):
                    file_path = os.path.join(temp_img_dir, filename)
                    with Image.open(file_path) as img:
                        aspect_ratio = img.width / img.height
                        
                        if aspect_ratio > 3840/2160:
                            new_width = 3840
                            new_height = int(3840 / aspect_ratio)
                        else:
                            new_height = 2160
                            new_width = int(2160 * aspect_ratio)
                        
                        resized = img.resize((new_width, new_height), LANCZOS)
                        new_filename = os.path.join(output_dir, f"{i:02d}.png")
                        resized.save(new_filename, "PNG", optimize=True)
                        
                return True, f"Successfully extracted {slide_count} slides", slide_count
        except Exception:
            pass
            
        # Fallback to pdf2image
        try:
            from pdf2image import convert_from_path
            
            images = convert_from_path(pdf_path, dpi=300, output_folder=temp_img_dir, fmt="png")
            slide_count = len(images)
            
            if slide_count > 0:
                for i, img in enumerate(images, 1):
                    aspect_ratio = img.width / img.height
                    
                    if aspect_ratio > 3840/2160:
                        new_width = 3840
                        new_height = int(3840 / aspect_ratio)
                    else:
                        new_height = 2160
                        new_width = int(2160 * aspect_ratio)
                    
                    resized = img.resize((new_width, new_height), LANCZOS)
                    output_path = os.path.join(output_dir, f"{i:02d}.png")
                    resized.save(output_path, "PNG", optimize=True)
                    
                return True, f"Successfully extracted {slide_count} slides", slide_count
        except Exception:
            pass
            
        return False, "Failed to extract slides - required tools not available", 0
        
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)

def extract_slides_with_method(
    pptx_path: str, 
    output_dir: str, 
    method: ConversionMethod,
    status_text: Optional[Any] = None,
    enable_auto_fallback: bool = True,
    target_resolution: Tuple[int, int] = (3840, 2160)
) -> Tuple[bool, str, int]:
    """Extract slides from a PPTX file using the specified method"""
    os.makedirs(output_dir, exist_ok=True)
    
    if status_text and hasattr(status_text, 'update_action'):
        status_text.update_action(f"Attempting extraction using {CONVERSION_METHODS[method]} method")
    elif status_text:
        status_text.info(f"Attempting extraction using {CONVERSION_METHODS[method]} method")
    
    success, message, slide_count = _try_method(pptx_path, output_dir, method, status_text, target_resolution)
    
    if not success and enable_auto_fallback:
        if status_text and hasattr(status_text, 'update_action'):
            status_text.update_action(f"Primary method failed. Trying fallback methods...")
        elif status_text:
            status_text.warning(f"Primary method failed: {message}")
            status_text.info("Trying fallback methods...")
        
        fallback_methods = get_fallback_methods(method)
        
        for fallback_method in fallback_methods:
            if status_text and hasattr(status_text, 'update_action'):
                status_text.update_action(f"Trying {CONVERSION_METHODS[fallback_method]} method...")
            elif status_text:
                status_text.info(f"Trying {CONVERSION_METHODS[fallback_method]} method...")
                
            success, message, slide_count = _try_method(pptx_path, output_dir, fallback_method, status_text, target_resolution)
            
            if success:
                if status_text and hasattr(status_text, 'update_action'):
                    status_text.update_action(f"Fallback method {CONVERSION_METHODS[fallback_method]} succeeded")
                elif status_text:
                    status_text.success(f"Fallback method {CONVERSION_METHODS[fallback_method]} succeeded")
                break
        
    return success, message, slide_count

def get_fallback_methods(primary_method: ConversionMethod) -> List[ConversionMethod]:
    """Get a list of fallback methods to try if the primary method fails"""
    all_methods = list(ConversionMethod)
    all_methods.remove(primary_method)
    
    if primary_method == ConversionMethod.LIBREOFFICE:
        return [ConversionMethod.PDF2IMAGE, ConversionMethod.PYTHON_PPTX]
    elif primary_method == ConversionMethod.PYTHON_PPTX:
        return [ConversionMethod.LIBREOFFICE, ConversionMethod.PDF2IMAGE]
    else:  # PDF2IMAGE
        return [ConversionMethod.LIBREOFFICE, ConversionMethod.PYTHON_PPTX]

def _try_method(
    pptx_path: str, 
    output_dir: str, 
    method: ConversionMethod,
    status_text: Optional[Any],
    target_resolution: Tuple[int, int]
) -> Tuple[bool, str, int]:
    """Try to extract slides using a specific method"""
    try:
        if method == ConversionMethod.LIBREOFFICE:
            return extract_slides_from_pptx(pptx_path, output_dir, status_text)
        elif method == ConversionMethod.PYTHON_PPTX:
            return extract_slides_directly(pptx_path, output_dir, status_text, target_resolution)
        elif method == ConversionMethod.PDF2IMAGE:
            return convert_with_pdf2image(pptx_path, output_dir, status_text, target_resolution)
        else:
            error_msg = f"Unknown conversion method: {method}"
            if status_text and hasattr(status_text, 'error'):
                status_text.error(error_msg)
            return False, error_msg, 0
            
    except ImportError as e:
        error_msg = f"Required module not available for {method} method: {str(e)}"
        if status_text and hasattr(status_text, 'error'):
            status_text.error(error_msg)
        return False, error_msg, 0
        
    except Exception as e:
        error_msg = f"Error in {method} method: {str(e)}"
        if status_text and hasattr(status_text, 'error'):
            status_text.error(error_msg)
        return False, error_msg, 0

# Increase the Pillow decompression bomb limit for large images
increase_image_decompression_limit()

st.set_page_config(page_title="06 4K Image (PPTX & ZIP) - Dreamlet", page_icon="🖼️")

# Constants for image processing
LOGO_PATH = 'config/logo.png'
COPYRIGHT_PATH = 'config/copyright.txt'
LOGO_SIZE = 250
LOGO_PADDING = 0
LOGO_POSITION = 'top-right'
COPYRIGHT_PADDING = 35
COPYRIGHT_POSITION = 'bottom-center'
FONT_SIZE = 65
FONT_COLOR = (0, 0, 0)  # Black
TARGET_RESOLUTION = (3840, 2160)  # 4K

def read_copyright(file_path):
    """Read copyright text from file"""
    try:
        # Try UTF-8 encoding first
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read().strip()
            
        # Clean up any encoding issues with the copyright symbol
        # Remove any 'A' characters that appear immediately before '©'
        text = re.sub(r'A©', '©', text)
        
        return text
    except UnicodeDecodeError:
        # If UTF-8 fails, try with different encoding
        try:
            with open(file_path, 'r', encoding='cp1252') as f:
                text = f.read().strip()
                # Clean up any encoding issues
                text = re.sub(r'A©', '©', text)
                return text
        except Exception:
            return "© All Rights Reserved"
    except Exception:
        return "© All Rights Reserved"

def get_position(img_width, img_height, element_width, element_height, position, padding):
    """Calculate position based on alignment and padding"""
    if position == 'top-left':
        return (padding, padding)
    elif position == 'top-right':
        return (img_width - element_width - padding, padding)
    elif position == 'bottom-left':
        return (padding, img_height - element_height - padding)
    elif position == 'bottom-right':
        return (img_width - element_width - padding, img_height - element_height - padding)
    elif position == 'bottom-center':
        return ((img_width - element_width) // 2, img_height - element_height - padding)
    else:
        raise ValueError(f"Invalid position: {position}")

def get_font(size):
    """Get a font with fallback options"""
    try:
        # Try to use a common sans-serif font
        return ImageFont.truetype("Arial.ttf", size)
    except IOError:
        try:
            # Fallback to DejaVuSans if Arial is not available
            return ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", size)
        except IOError:
            # If both fail, use the default font
            return ImageFont.load_default()

def process_image(input_path, output_path, logo, copyright_text):
    """Process an image by adding logo and copyright text"""
    try:
        with Image.open(input_path) as img:
            # Check if the image is already 4K or larger
            if img.width >= TARGET_RESOLUTION[0] and img.height >= TARGET_RESOLUTION[1]:
                processed_img = img.copy()
            else:
                # Upscale the image to 4K while maintaining aspect ratio
                aspect_ratio = img.width / img.height
                if aspect_ratio > 16/9:
                    new_width = TARGET_RESOLUTION[0]
                    new_height = int(TARGET_RESOLUTION[0] / aspect_ratio)
                else:
                    new_height = TARGET_RESOLUTION[1]
                    new_width = int(TARGET_RESOLUTION[1] * aspect_ratio)
                processed_img = img.resize((new_width, new_height), LANCZOS)
            
            # Convert to RGBA if necessary
            if processed_img.mode != 'RGBA':
                processed_img = processed_img.convert('RGBA')

            # Add logo if it exists
            if logo:
                logo_resized = logo.resize((LOGO_SIZE, LOGO_SIZE), LANCZOS)
                logo_pos = get_position(processed_img.width, processed_img.height, 
                                       LOGO_SIZE, LOGO_SIZE, LOGO_POSITION, LOGO_PADDING)
                
                # Create a new image for the logo with an alpha channel
                logo_img = Image.new('RGBA', processed_img.size, (0, 0, 0, 0))
                logo_img.paste(logo_resized, logo_pos, logo_resized)
                
                # Composite the logo onto the processed image
                processed_img = Image.alpha_composite(processed_img, logo_img)

            # Add copyright
            draw = ImageDraw.Draw(processed_img)
            font = get_font(FONT_SIZE)
            
            # Get text size
            text_bbox = draw.textbbox((0, 0), copyright_text, font=font)
            text_width = text_bbox[2] - text_bbox[0]
            text_height = text_bbox[3] - text_bbox[1]
            
            text_pos = get_position(processed_img.width, processed_img.height, 
                                   text_width, text_height, COPYRIGHT_POSITION, COPYRIGHT_PADDING)
            draw.text(text_pos, copyright_text, font=font, fill=FONT_COLOR)

            # Save the processed image with the same format and quality as the original
            if img.format == 'JPEG':
                processed_img = processed_img.convert('RGB')
                processed_img.save(output_path, format=img.format, quality=95, optimize=True, progressive=True)
            elif img.format == 'PNG':
                processed_img.save(output_path, format=img.format, optimize=True)
            else:
                processed_img.save(output_path, format=img.format)
        
        return True
    except Exception as e:
        st.error(f"Error processing {input_path}: {str(e)}")
        return False

def extract_and_upscale_images(pptx_file, output_folder, conversion_method=None, enable_auto_fallback=True, status_mgr=None):
    """
    Extract images from a PPTX file, upscale them to 4K, and save them in the output folder.
    
    Args:
        pptx_file: Path to the PPTX file
        output_folder: Directory to save the extracted images
        conversion_method: Method to use for conversion (libreoffice, python-pptx, pdf2image)
        enable_auto_fallback: Whether to try other methods if the selected one fails
        status_mgr: Optional StatusManager to use for status updates
        
    Returns:
        Tuple of (success, extracted_images, slide_count)
    """
    # Use the multi-method converter
    # If we have a status manager, use it, otherwise fall back to st
    success, message, slide_count = extract_slides_with_method(
        pptx_file, 
        output_folder, 
        method=conversion_method or ConversionMethod.LIBREOFFICE,
        status_text=status_mgr or st,
        enable_auto_fallback=enable_auto_fallback,
        target_resolution=TARGET_RESOLUTION
    )
    
    if not success:
        if status_mgr:
            status_mgr.error(message)
        else:
            st.error(message)
        return False, [], 0
    
    # Gather the paths of the created images
    extracted_images = []
    if slide_count > 0:
        for i in range(1, slide_count + 1):
            image_path = os.path.join(output_folder, f"{i:02d}.png")
            if os.path.exists(image_path):
                extracted_images.append(image_path)
    
    return success, extracted_images, slide_count

def find_presentations() -> List[str]:
    """
    Find all presentation files in the input directory
    
    Returns:
        List of presentation file paths
    """
    input_dir = get_input_directory()
    
    # Find all presentation files
    all_presentations = find_presentation_files(input_dir)
    
    return all_presentations

class StatusManager:
    """A class to manage and simplify status messages during presentation processing"""
    
    def __init__(self, progress_bar, status_text, current_action_text=None):
        self.progress_bar = progress_bar
        self.status_text = status_text
        self.current_action_text = current_action_text
        self.current_file = ""
        self.current_action = ""
        self.processed_count = 0
        self.total_count = 0
        self.start_time = datetime.now()
        
    def set_total(self, total):
        """Set the total number of presentations to process"""
        self.total_count = total
        
    def update_progress(self, current_index, filename):
        """Update the progress bar and main status message"""
        self.processed_count = current_index
        self.current_file = filename
        progress = self.processed_count / self.total_count if self.total_count > 0 else 0
        self.progress_bar.progress(progress)
        
        # Calculate and display estimated time remaining
        if self.processed_count > 1:
            elapsed_time = (datetime.now() - self.start_time).total_seconds()
            time_per_file = elapsed_time / (self.processed_count - 1)  # Exclude current file that just started
            remaining_files = self.total_count - self.processed_count
            remaining_seconds = remaining_files * time_per_file
            remaining_time = str(timedelta(seconds=int(remaining_seconds)))
            
            self.status_text.info(f"Processing {self.processed_count}/{self.total_count}: {filename} (Est. time remaining: {remaining_time})")
        else:
            self.status_text.info(f"Processing {self.processed_count}/{self.total_count}: {filename}")
    
    def update_action(self, message):
        """Update only the current action, not the main status display"""
        self.current_action = message
        if self.current_action_text:
            self.current_action_text.info(message)
        
    def info(self, message):
        """For compatibility with status_text.info but only updates the current action"""
        self.update_action(message)
        
    def warning(self, message):
        """For compatibility with status_text.warning"""
        st.warning(message)
        
    def error(self, message):
        """For compatibility with status_text.error"""
        st.error(message)
        
    def success(self, message):
        """For compatibility with status_text.success"""
        self.status_text.success(message)
        if self.current_action_text:
            self.current_action_text.empty()  # Clear the current action text when we're done

def process_presentations(input_dir: str, progress_bar, status_text, create_without_logo_folder=False, 
                 conversion_method=None, enable_auto_fallback=True, current_action_text=None) -> List[Dict]:
    """
    Process all PPTX and ZIP files in the input directory structure
    
    Args:
        input_dir: Input directory path
        progress_bar: Streamlit progress bar object
        status_text: Streamlit text area for status updates
        create_without_logo_folder: Whether to create without_logo_png folder for original images
        conversion_method: Method to use for PPTX conversion (libreoffice, python-pptx, pdf2image) - not used for ZIP files
        enable_auto_fallback: Whether to try other methods if the selected one fails (PPTX only)
        current_action_text: Streamlit text element for displaying current action
        
    Returns:
        List of dictionaries with processing results
    """
    results = []
    
    # Create status manager
    status_mgr = StatusManager(progress_bar, status_text, current_action_text)
    
    # Try to load logo and copyright text
    try:
        logo = Image.open(LOGO_PATH).convert("RGBA") if os.path.exists(LOGO_PATH) else None
        copyright_text = read_copyright(COPYRIGHT_PATH) if os.path.exists(COPYRIGHT_PATH) else "© All Rights Reserved"
    except Exception as e:
        st.warning(f"Could not load logo or copyright: {str(e)}")
        logo = None
        copyright_text = "© All Rights Reserved"
    
    # Find all presentations in the input directory
    all_presentations = find_presentations()
    total_presentations = len(all_presentations)
    status_mgr.set_total(total_presentations)
    
    # Process each presentation
    for i, pptx_file in enumerate(all_presentations):
        # Update progress with file info
        processed_count = i + 1
        file_name = os.path.basename(pptx_file)
        status_mgr.update_progress(processed_count, file_name)
        
        result = {
            "file_path": pptx_file,
            "status": "error",
            "message": "",
            "slide_count": 0,
            "output_dir": ""
        }
        
        try:
            # Determine file type first
            file_extension = os.path.splitext(pptx_file)[1].lower()
            is_zip_file = file_extension == '.zip'
            
            # Get the directory containing the presentation
            pptx_dir = os.path.dirname(pptx_file)
            
            # Check if this is inside an all_pptx directory
            parent_dir = os.path.basename(pptx_dir)
            if parent_dir == 'all_pptx':
                # If file is in all_pptx, create lecture folder in parent directory
                course_dir = os.path.dirname(pptx_dir)  # Go up one level from all_pptx
                pptx_name = os.path.splitext(os.path.basename(pptx_file))[0]
                
                # Extract lecture number from filename for proper naming
                lecture_match = re.search(r'(\d+)', pptx_name)
                if lecture_match:
                    lecture_num = lecture_match.group(1).zfill(2)
                    lecture_folder = os.path.join(course_dir, f"Lecture {lecture_num}")
                else:
                    lecture_folder = os.path.join(course_dir, pptx_name)
                    
                os.makedirs(lecture_folder, exist_ok=True)
                
                # Create "English image" folder inside the lecture folder
                output_folder = os.path.join(lecture_folder, "English image")
                os.makedirs(output_folder, exist_ok=True)
                
                # The all_pptx folder already exists since the file is in it
                all_pptx_folder = pptx_dir
            else:
                # Create 'all_pptx' directory in the same folder as the presentation
                all_pptx_folder = os.path.join(pptx_dir, 'all_pptx')
                os.makedirs(all_pptx_folder, exist_ok=True)
                
                # Create folder with same name as file for images
                file_name = os.path.splitext(os.path.basename(pptx_file))[0]
                
                # For ZIP files with numeric names, add "Lecture" prefix to match PPTX structure
                if is_zip_file and re.match(r'^\d+$', file_name):
                    lecture_folder = os.path.join(pptx_dir, f"Lecture {file_name.zfill(2)}")
                else:
                    lecture_folder = os.path.join(pptx_dir, file_name)
                    
                os.makedirs(lecture_folder, exist_ok=True)
                
                # Create "English image" folder inside the lecture folder
                output_folder = os.path.join(lecture_folder, "English image")
                os.makedirs(output_folder, exist_ok=True)
            
            # Create 'without_logo_png' folder inside the output folder only if requested
            without_logo_folder = os.path.join(output_folder, 'without_logo_png')
            if create_without_logo_folder:
                os.makedirs(without_logo_folder, exist_ok=True)
            
            # Check if file is already processed and in all_pptx
            new_pptx_path = os.path.join(all_pptx_folder, os.path.basename(pptx_file))
            if os.path.exists(new_pptx_path):
                # If already processed, just update the result
                result["message"] = "Presentation already processed"
                
                # Count existing images
                png_count = len([f for f in os.listdir(output_folder) 
                               if f.endswith('.png') and not os.path.isdir(os.path.join(output_folder, f))])
                
                result["status"] = "success"
                result["slide_count"] = png_count
                result["output_dir"] = output_folder
                results.append(result)
                continue
                
            # Extract and upscale images from the presentation or ZIP file
            import tempfile
            temp_dir = None
            
            try:
                # Only create and use without_logo_folder if requested
                if create_without_logo_folder:
                    # Ensure the folder exists
                    os.makedirs(without_logo_folder, exist_ok=True)
                    
                    if is_zip_file:
                        # Extract from ZIP file
                        success, extracted_images, slide_count = extract_slides_from_zip(
                            pptx_file,
                            without_logo_folder,
                            status_mgr=status_mgr,
                            target_resolution=TARGET_RESOLUTION
                        )
                    else:
                        # Extract from PPTX file
                        success, extracted_images, slide_count = extract_and_upscale_images(
                            pptx_file, 
                            without_logo_folder,
                            conversion_method=conversion_method,
                            enable_auto_fallback=enable_auto_fallback,
                            status_mgr=status_mgr
                        )
                    
                    # Process each extracted image (add logo and copyright)
                    processed_img_count = 0
                    for img_path in extracted_images:
                        # Original image is in without_logo_folder
                        # Processed image will be in output_folder
                        output_filename = os.path.basename(img_path)
                        output_path = os.path.join(output_folder, output_filename)
                        
                        # Process the image
                        if process_image(img_path, output_path, logo, copyright_text):
                            processed_img_count += 1
                else:
                    # Use a temporary directory that we manage explicitly
                    temp_dir = tempfile.mkdtemp()
                    
                    if is_zip_file:
                        # Extract from ZIP file
                        success, extracted_images, slide_count = extract_slides_from_zip(
                            pptx_file,
                            temp_dir,
                            status_mgr=status_mgr,
                            target_resolution=TARGET_RESOLUTION
                        )
                    else:
                        # Extract from PPTX file
                        success, extracted_images, slide_count = extract_and_upscale_images(
                            pptx_file, 
                            temp_dir,
                            conversion_method=conversion_method,
                            enable_auto_fallback=enable_auto_fallback,
                            status_mgr=status_mgr
                        )
                    
                    # Process each extracted image (add logo and copyright)
                    processed_img_count = 0
                    for img_path in extracted_images:
                        # Temporary image is in temp_dir
                        # Processed image will go directly to output_folder
                        output_filename = os.path.basename(img_path)
                        output_path = os.path.join(output_folder, output_filename)
                        
                        # Process the image
                        if process_image(img_path, output_path, logo, copyright_text):
                            processed_img_count += 1
                
                if not success:
                    result["message"] = "Failed to extract images from presentation"
                    results.append(result)
                    continue
            finally:
                # Clean up the temporary directory if it was created
                if temp_dir and os.path.exists(temp_dir):
                    import shutil
                    shutil.rmtree(temp_dir, ignore_errors=True)
            
            # Move the presentation file to all_pptx folder
            if os.path.exists(pptx_file) and not os.path.exists(new_pptx_path):
                os.rename(pptx_file, new_pptx_path)
            
            # Update result
            result["status"] = "success"
            result["message"] = f"Successfully processed {processed_img_count} images"
            result["slide_count"] = slide_count
            result["output_dir"] = output_folder
        
        except Exception as e:
            result["message"] = f"Error: {str(e)}"
        
        results.append(result)
        
        # Small delay for UI updates
        time.sleep(0.1)
    
    # Ensure progress bar reaches 100% at the end
    progress_bar.progress(1.0)
    
    return results

def main():
    st.title("4K Image (PPTX & ZIP)")
    st.write("Generate high-resolution 4K images from presentation slides (PPTX) or ZIP archives containing images.")
    
    input_dir = get_input_directory()
    
    if not os.path.exists(input_dir):
        st.error(f"Input directory not found: {input_dir}")
        st.info("Please create an 'input' directory in the project root and add your files.")
        return
    
    # Find presentations
    all_presentations = find_presentations()
    
    if not all_presentations:
        st.warning("No presentation files found in the input directory.")
        return
    
    # Display number of presentations found
    total_presentations = len(all_presentations)
    pptx_count = sum(1 for f in all_presentations if f.lower().endswith('.pptx'))
    zip_count = sum(1 for f in all_presentations if f.lower().endswith('.zip'))
    st.info(f"Found {total_presentations} files in input directory: {pptx_count} PPTX, {zip_count} ZIP")
    
    # Options section
    st.header("Options")
    
    # Option for creating without_logo_png folder
    create_without_logo_folder = st.checkbox("Create 'without_logo_png' folder", value=False, 
                                          help="If checked, a folder containing original images without logo and copyright will be created")
    
    # Conversion method selection (only for PPTX files)
    st.subheader("Conversion Method (PPTX only)")
    st.write("Select the method to use for extracting slides from PPTX presentations:")
    st.caption("Note: ZIP files are processed directly and don't use these conversion methods")
    
    # Create a dropdown for conversion method selection
    conversion_methods = list(CONVERSION_METHODS.items())
    conversion_method = st.selectbox(
        "Select conversion method:",
        options=[method.value for method, _ in conversion_methods],
        format_func=lambda x: CONVERSION_METHODS[x],
        index=0,
        help="Choose the method used to extract slides from PPTX presentations (not applicable to ZIP files)"
    )
    
    # Option for automatic fallback
    enable_auto_fallback = st.checkbox(
        "Enable automatic fallback", 
        value=True,
        help="If checked, automatically try other methods if the selected method fails"
    )
    
    # Information about conversion methods and file types
    with st.expander("About File Types and Conversion Methods"):
        st.markdown("""
        ### Supported File Types
        
        #### PPTX Files
        - PowerPoint presentation files
        - Slides are extracted and converted to images
        - Multiple conversion methods available (see below)
        
        #### ZIP Files
        - ZIP archives containing image files (PNG, JPG, JPEG)
        - Images are extracted directly from the ZIP
        - Automatically sorted by numeric prefix (e.g., 1_Title.png, 2_Content.png)
        - Upscaled to 4K if needed
        - No conversion method selection needed
        
        ### Conversion Methods (PPTX only)
        
        #### LibreOffice + pdftoppm (best quality)
        - Uses LibreOffice to convert PPTX to PDF
        - Uses pdftoppm to convert PDF to high-quality images
        - Provides the best visual quality for complex slides
        - Requires LibreOffice and poppler-utils to be installed
        
        #### Direct PPTX processing (fastest)
        - Uses python-pptx to extract slides directly
        - Faster than other methods
        - May not render complex slides as accurately
        - Good for simple presentations with basic elements
        
        #### PDF2Image (most compatible)
        - Uses LibreOffice to convert PPTX to PDF
        - Uses pdf2image library to convert PDF to images
        - More compatible fallback option
        - Slightly lower quality than pdftoppm method
        """)
        
        st.info("💡 Tip: If you encounter issues with PPTX conversion, try another method or enable automatic fallback.")
    
    # Process presentations
    st.header("Extract Presentations")
    
    if st.button("Extract Presentations", disabled=total_presentations == 0):
        if total_presentations == 0:
            st.warning("No presentations found for processing.")
            return
        
        # Process presentations with progress tracking
        progress_bar = st.progress(0)
        status_text = st.empty()
        current_action_text = st.empty()  # For showing the current action
        results_container = st.container()
        
        # Process the presentations with progress tracking
        results = process_presentations(
            input_dir, 
            progress_bar, 
            status_text, 
            create_without_logo_folder,
            conversion_method=conversion_method,
            enable_auto_fallback=enable_auto_fallback,
            current_action_text=current_action_text
        )
        
        # Clear the current action text when done
        current_action_text.empty()
        
        # Complete
        status_text.success(f"Processed {total_presentations} presentations")
        
        # Display results
        with results_container:
            st.subheader("Processing Results")
            
            # Statistics
            success_count = sum(1 for r in results if r["status"] == "success")
            error_count = sum(1 for r in results if r["status"] == "error")
            slide_count = sum(r.get("slide_count", 0) for r in results)
            
            st.write(f"✅ Successfully processed: {success_count}")
            st.write(f"❌ Errors: {error_count}")
            st.write(f"📊 Total slides extracted: {slide_count}")
            
            # Detailed results
            for result in results:
                status_icon = "✅" if result["status"] == "success" else "❌"
                file_name = os.path.basename(result["file_path"])
                
                with st.expander(f"{status_icon} {file_name}"):
                    st.write(f"**Status:** {result['status']}")
                    st.write(f"**Message:** {result['message']}")
                    
                    if result["status"] == "success":
                        st.write(f"**Slides extracted:** {result['slide_count']}")
                        st.write(f"**Output directory:** {result['output_dir']}")

if __name__ == "__main__":
    main()
