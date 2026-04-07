import os
from pathlib import Path
import subprocess
import tempfile
from typing import List, Dict, Optional, Tuple

def extract_slides_from_pptx(pptx_path: str, output_dir: str, resolution: str = "4K") -> Tuple[bool, str, int]:
    """
    Extract high-resolution images from PowerPoint slides
    
    Args:
        pptx_path: Path to PowerPoint file
        output_dir: Directory to save extracted images
        resolution: Image resolution (HD, 2K, 4K)
        
    Returns:
        Tuple of (success, message, slide_count)
    """
    # Create output directory if it doesn't exist
    os.makedirs(output_dir, exist_ok=True)
    
    # Define resolution parameters based on requested quality
    if resolution == "4K":
        width = 3840
        height = 2160
    elif resolution == "2K":
        width = 2560
        height = 1440
    else:  # Default to HD
        width = 1920
        height = 1080
    
    # Maximum pixels to prevent decompression bomb errors
    max_pixels = 178956970
    
    try:
        # Import our improved converter
        try:
            from utils.pptx_converter import extract_slides_from_pptx as improved_extract
            success, message, slide_count = improved_extract(pptx_path, output_dir)
            if success:
                return success, message, slide_count
        except ImportError:
            # If the improved converter isn't available, continue with the fallback method
            pass
            
        # Fallback method using direct LibreOffice + pdftoppm conversion
        # Create a temporary directory for the conversion
        with tempfile.TemporaryDirectory() as temp_dir:
            # Use libreoffice to convert PPTX to PDF
            pdf_output = os.path.join(temp_dir, "presentation.pdf")
            lo_command = [
                "libreoffice", "--headless", "--convert-to", "pdf",
                "--outdir", temp_dir, pptx_path
            ]
            
            try:
                subprocess.run(lo_command, check=True, capture_output=True, timeout=120)
            except FileNotFoundError:
                # If libreoffice is not installed, try soffice
                lo_command[0] = "soffice"
                try:
                    subprocess.run(lo_command, check=True, capture_output=True, timeout=120)
                except FileNotFoundError:
                    return False, "LibreOffice/OpenOffice not found - please install to enable slide extraction", 0
                except subprocess.TimeoutExpired:
                    return False, "LibreOffice process timed out - presentation may be too complex", 0
            except subprocess.TimeoutExpired:
                return False, "LibreOffice process timed out - presentation may be too complex", 0
                
            if not os.path.exists(pdf_output):
                return False, "Failed to convert PPTX to PDF", 0
            
            # Try multiple methods to convert PDF to images
            slide_count = 0
            
            # Method 1: Use pdftoppm if available
            try:
                base_filename = os.path.join(output_dir, "slide")
                pdftoppm_command = [
                    "pdftoppm", "-png", "-r", "300",  # 300 dpi for high quality
                    pdf_output, base_filename
                ]
                subprocess.run(pdftoppm_command, check=True, capture_output=True, timeout=120)
                
                # Count and rename the extracted slides
                extracted_files = sorted([f for f in os.listdir(output_dir) if f.startswith("slide-") and f.endswith(".png")])
                slide_count = len(extracted_files)
                
                # Check if slides were extracted
                if slide_count > 0:
                    # Process images to achieve desired resolution while preventing decompression bomb errors
                    from PIL import Image
                    for i, filename in enumerate(extracted_files, 1):
                        file_path = os.path.join(output_dir, filename)
                        with Image.open(file_path) as img:
                            # Calculate new dimensions while maintaining aspect ratio
                            aspect_ratio = img.width / img.height
                            
                            if aspect_ratio > width/height:  # Wider than target
                                new_width = width
                                new_height = int(width / aspect_ratio)
                            else:  # Taller than or equal to target
                                new_height = height
                                new_width = int(height * aspect_ratio)
                            
                            # Check if the new dimensions would exceed our pixel limit
                            if new_width * new_height > max_pixels:
                                scale_factor = (max_pixels / (new_width * new_height)) ** 0.5
                                new_width = int(new_width * scale_factor)
                                new_height = int(new_height * scale_factor)
                            
                            # Resize the image
                            resized = img.resize((new_width, new_height), Image.LANCZOS)
                            
                            # Create new filename with format 01.png, 02.png, etc.
                            new_filename = os.path.join(output_dir, f"{i:02d}.png")
                            
                            # Save the resized image (replacing the original)
                            resized.save(new_filename, "PNG", optimize=True)
                            
                            # Remove the original file if it's different from the new one
                            if file_path != new_filename and os.path.exists(file_path):
                                os.remove(file_path)
                                
                    return True, f"Successfully extracted {slide_count} slides", slide_count
            except (FileNotFoundError, subprocess.CalledProcessError, subprocess.TimeoutExpired) as e:
                # pdftoppm failed, try the next method
                pass
                
            # Method 2: Use pdf2image as a fallback
            try:
                from pdf2image import convert_from_path
                
                # Convert PDF to images
                images = convert_from_path(pdf_output, dpi=300, output_folder=temp_dir, fmt="png")
                
                # Process each image
                slide_count = len(images)
                
                # Check if we got any images
                if slide_count > 0:
                    # Process each image to the correct resolution
                    for i, img in enumerate(images, 1):
                        # Calculate new dimensions while maintaining aspect ratio
                        aspect_ratio = img.width / img.height
                        
                        if aspect_ratio > width/height:  # Wider than target
                            new_width = width
                            new_height = int(width / aspect_ratio)
                        else:  # Taller than or equal to target
                            new_height = height
                            new_width = int(height * aspect_ratio)
                        
                        # Check if the new dimensions would exceed our pixel limit
                        if new_width * new_height > max_pixels:
                            scale_factor = (max_pixels / (new_width * new_height)) ** 0.5
                            new_width = int(new_width * scale_factor)
                            new_height = int(new_height * scale_factor)
                        
                        # Resize the image
                        resized = img.resize((new_width, new_height), Image.LANCZOS)
                        
                        # Save to output directory
                        output_path = os.path.join(output_dir, f"{i:02d}.png")
                        resized.save(output_path, "PNG", optimize=True)
                        
                    return True, f"Successfully extracted {slide_count} slides", slide_count
            except ImportError:
                # pdf2image not available
                pass
            except Exception as pdf2image_error:
                # pdf2image failed for some other reason
                return False, f"Error extracting slides with pdf2image: {str(pdf2image_error)}", 0
                
            # If we got here, all extraction methods failed
            return False, "Failed to extract slides - required tools not available", 0
            
    except Exception as e:
        return False, f"Error extracting slides: {str(e)}", 0

def count_slides_in_pptx(pptx_path: str) -> Tuple[bool, int]:
    """
    Count the number of slides in a PowerPoint presentation
    
    Args:
        pptx_path: Path to PowerPoint file
        
    Returns:
        Tuple of (success, slide_count)
    """
    try:
        # First, try using python-pptx which is more reliable
        try:
            import pptx
            prs = pptx.Presentation(pptx_path)
            return True, len(prs.slides)
        except (ImportError, Exception) as e:
            # If python-pptx fails, try the LibreOffice approach
            pass
        
        # Create a temporary directory for extraction
        with tempfile.TemporaryDirectory() as temp_dir:
            # Use libreoffice to convert to PDF and count pages
            pdf_output = os.path.join(temp_dir, "presentation.pdf")
            
            # Try with libreoffice command
            lo_command = [
                "libreoffice", "--headless", "--convert-to", "pdf",
                "--outdir", temp_dir, pptx_path
            ]
            
            try:
                # Try libreoffice command
                subprocess.run(lo_command, check=True, capture_output=True, timeout=60)
            except (FileNotFoundError, subprocess.CalledProcessError):
                # If libreoffice fails, try soffice
                lo_command[0] = "soffice"
                try:
                    subprocess.run(lo_command, check=True, capture_output=True, timeout=60)
                except (FileNotFoundError, subprocess.CalledProcessError):
                    # Both failed, return failure
                    return False, 0
            
            # Check if PDF was created
            if not os.path.exists(pdf_output):
                return False, 0
                
            # Try multiple methods to get the page count
            
            # Method 1: Use pdfinfo if available
            try:
                pdfinfo_command = ["pdfinfo", pdf_output]
                result = subprocess.run(pdfinfo_command, check=True, capture_output=True, text=True, timeout=10)
                
                # Parse output to get page count
                for line in result.stdout.splitlines():
                    if line.startswith("Pages:"):
                        return True, int(line.split(":")[1].strip())
            except (FileNotFoundError, subprocess.CalledProcessError, ValueError):
                pass
                
            # Method 2: Use pdf2image if available
            try:
                from pdf2image import convert_from_path
                images = convert_from_path(pdf_output, dpi=72)  # Low DPI just for counting
                return True, len(images)
            except (ImportError, Exception):
                pass
                
            # Method 3: Last resort - try pdftoppm to convert and count files
            try:
                output_prefix = os.path.join(temp_dir, "page")
                pdftoppm_command = ["pdftoppm", "-png", pdf_output, output_prefix]
                subprocess.run(pdftoppm_command, check=True, capture_output=True, timeout=30)
                
                # Count the generated files
                page_count = len([f for f in os.listdir(temp_dir) if f.startswith("page-") and f.endswith(".png")])
                if page_count > 0:
                    return True, page_count
            except (FileNotFoundError, subprocess.CalledProcessError):
                pass
            
            # All methods failed
            return False, 0
    except Exception as e:
        # Fallback to python-pptx as last resort
        try:
            import pptx
            prs = pptx.Presentation(pptx_path)
            return True, len(prs.slides)
        except Exception:
            return False, 0
