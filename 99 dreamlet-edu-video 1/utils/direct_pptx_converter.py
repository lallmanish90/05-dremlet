import os
import io
from PIL import Image
import cv2
import numpy as np
from pptx import Presentation

def extract_slides_directly(pptx_path, output_dir, status_text=None, target_resolution=(3840, 2160)):
    """
    Extract slides directly from PPTX file using python-pptx and convert to 4K images.
    This is faster than the LibreOffice method but may have lower quality for complex slides.
    
    Args:
        pptx_path: Path to the PPTX file
        output_dir: Directory to save the extracted images
        status_text: Streamlit text element for status updates
        target_resolution: Target resolution (width, height) for upscaling
        
    Returns:
        Tuple of (success, message, slide_count)
    """
    try:
        # Ensure output directory exists
        os.makedirs(output_dir, exist_ok=True)
        
        if status_text:
            status_text.info(f"Extracting slides directly from PPTX: {os.path.basename(pptx_path)}")
        
        # Open presentation
        prs = Presentation(pptx_path)
        slide_count = len(prs.slides)
        
        if slide_count == 0:
            if status_text:
                status_text.warning("Presentation contains no slides.")
            return False, "Presentation contains no slides", 0
        
        # Process each slide
        for i, slide in enumerate(prs.slides, 1):
            # Extract slide dimensions
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # Create image with the same aspect ratio as the slide
            slide_aspect_ratio = slide_width / slide_height
            
            # Determine output dimensions while maintaining aspect ratio
            if slide_aspect_ratio > target_resolution[0] / target_resolution[1]:
                # Wider than 16:9
                width = target_resolution[0]
                height = int(width / slide_aspect_ratio)
            else:
                # Taller than or equal to 16:9
                height = target_resolution[1]
                width = int(height * slide_aspect_ratio)
            
            # Create a blank image with white background
            img = np.ones((height, width, 3), np.uint8) * 255
            
            # Calculate scale factors
            scale_x = width / slide_width
            scale_y = height / slide_height
            
            # Process each shape in the slide
            for shape in slide.shapes:
                # Try to get the shape image if it has one
                if hasattr(shape, 'image') and shape.image:
                    try:
                        # Get image data
                        image_stream = io.BytesIO(shape.image.blob)
                        slide_image = Image.open(image_stream)
                        
                        # Convert PIL Image to OpenCV format
                        slide_image_cv = cv2.cvtColor(np.array(slide_image), cv2.COLOR_RGB2BGR)
                        
                        # Calculate the position and size after scaling
                        left = int(shape.left * scale_x)
                        top = int(shape.top * scale_y)
                        img_width = int(shape.width * scale_x)
                        img_height = int(shape.height * scale_y)
                        
                        # Resize image to fit the scaled dimensions
                        resized_image = cv2.resize(slide_image_cv, (img_width, img_height), interpolation=cv2.INTER_LANCZOS4)
                        
                        # Place the image on the slide at the correct position
                        try:
                            # Make sure the region doesn't go outside the image bounds
                            if img_height > 0 and img_width > 0:
                                if top + img_height <= height and left + img_width <= width:
                                    img[top:top+img_height, left:left+img_width] = resized_image
                        except Exception as e:
                            if status_text:
                                status_text.warning(f"Error placing image on slide {i}: {str(e)}")
                    except Exception as e:
                        if status_text:
                            status_text.warning(f"Error processing image in slide {i}: {str(e)}")
            
            # Save the slide as PNG
            output_path = os.path.join(output_dir, f"{i:02d}.png")
            cv2.imwrite(output_path, img, [cv2.IMWRITE_PNG_COMPRESSION, 9])
            
            if status_text:
                status_text.info(f"Processed slide {i}/{slide_count}")
        
        return True, f"Successfully extracted {slide_count} slides", slide_count
        
    except Exception as e:
        error_msg = f"Error extracting slides directly: {str(e)}"
        if status_text:
            status_text.error(error_msg)
        return False, error_msg, 0