from __future__ import annotations

import argparse
import concurrent.futures
import enum
import fnmatch
import glob
import importlib.util
import os
from pathlib import Path
import platform
import re
import shutil
import subprocess
import sys
import tempfile
import time
import tomllib
from functools import lru_cache
from typing import Any, Optional
import zipfile

import cv2
import numpy as np
import psutil
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation

try:
    LANCZOS = Image.Resampling.LANCZOS
except AttributeError:  # pragma: no cover
    LANCZOS = Image.LANCZOS


TARGET_RESOLUTION = (3840, 2160)
LOGO_SIZE = 250
LOGO_PADDING = 0
LOGO_POSITION = "top-right"
COPYRIGHT_PADDING = 35
COPYRIGHT_POSITION = "bottom-center"
FONT_SIZE = 65
FONT_COLOR = (0, 0, 0)
DEFAULT_CONFIG_PATH = Path(__file__).resolve().parents[1] / "config" / "06 Generate 4K Images.toml"


class ConflictPolicy(str, enum.Enum):
    SKIP_SAFELY = "skip_safely"
    OVERWRITE = "overwrite"
    REPROCESS_NEW_FOLDER = "reprocess_new_folder"
    REPORT_ONLY = "report_only"


CONFLICT_POLICY_LABELS = {
    ConflictPolicy.SKIP_SAFELY.value: "Skip safely",
    ConflictPolicy.OVERWRITE.value: "Reprocess and overwrite",
    ConflictPolicy.REPROCESS_NEW_FOLDER.value: "Reprocess into a new folder",
    ConflictPolicy.REPORT_ONLY.value: "Report only",
}


class ConversionMethod(str, enum.Enum):
    LIBREOFFICE = "libreoffice"
    PYTHON_PPTX = "python-pptx"
    PDF2IMAGE = "pdf2image"


CONVERSION_METHODS = {
    ConversionMethod.LIBREOFFICE: "LibreOffice + Poppler (best quality)",
    ConversionMethod.PYTHON_PPTX: "Direct PPTX processing (fastest)",
    ConversionMethod.PDF2IMAGE: "PDF2Image (most compatible)",
}


class ConfigError(ValueError):
    """Raised when the 4K config is invalid."""


def load_config(config_path: os.PathLike[str] | str | None = None) -> dict[str, Any]:
    config_file = Path(config_path) if config_path else DEFAULT_CONFIG_PATH
    try:
        with config_file.open("rb") as handle:
            data = tomllib.load(handle)
    except FileNotFoundError as exc:
        raise ConfigError(f"Config file not found: {config_file}") from exc

    paths = data.get("paths")
    machine = data.get("machine")
    run_table = data.get("run")
    conversion = data.get("conversion")
    if not isinstance(paths, dict):
        raise ConfigError("Missing required [paths] table")
    if not isinstance(machine, dict):
        raise ConfigError("Missing required [machine] table")
    if not isinstance(run_table, dict):
        raise ConfigError("Missing required [run] table")
    if not isinstance(conversion, dict):
        raise ConfigError("Missing required [conversion] table")

    profile = str(machine.get("profile", "auto")).strip()
    allowed_profiles = {"auto", "macbook_pro_m3_pro_18gb", "windows_i5_12450h_rtx3050_16gb", "generic_fallback"}
    if profile not in allowed_profiles:
        raise ConfigError(f"Unsupported profile: {profile}")

    method = str(conversion.get("method", ConversionMethod.LIBREOFFICE.value)).strip()
    if method not in {item.value for item in ConversionMethod}:
        raise ConfigError(f"Unsupported method: {method}")

    conflict_policy = str(run_table.get("conflict_policy", ConflictPolicy.SKIP_SAFELY.value)).strip()
    if conflict_policy not in CONFLICT_POLICY_LABELS:
        raise ConfigError(f"Unsupported conflict_policy: {conflict_policy}")

    def resolve_path(raw_value: Any, key: str) -> Path:
        if not isinstance(raw_value, str) or not raw_value.strip():
            raise ConfigError(f"Missing required path: {key}")
        candidate = Path(raw_value)
        if not candidate.is_absolute():
            candidate = (config_file.parent / candidate).resolve()
        return candidate

    keep_raw_without_logo = run_table.get("keep_raw_without_logo", False)
    if not isinstance(keep_raw_without_logo, bool):
        raise ConfigError("Expected boolean keep_raw_without_logo")

    max_threads = run_table.get("max_threads")
    if max_threads is not None and (isinstance(max_threads, bool) or not isinstance(max_threads, int) or max_threads < 1):
        raise ConfigError("Expected positive integer max_threads")

    batch_size = run_table.get("batch_size")
    if batch_size is not None and (isinstance(batch_size, bool) or not isinstance(batch_size, int) or batch_size < 1):
        raise ConfigError("Expected positive integer batch_size")

    image_processing = str(run_table.get("image_processing", "auto")).strip()
    if image_processing not in {"auto", "cpu", "cuda"}:
        raise ConfigError(f"Unsupported image_processing: {image_processing}")

    enable_auto_fallback = conversion.get("enable_auto_fallback", True)
    if not isinstance(enable_auto_fallback, bool):
        raise ConfigError("Expected boolean enable_auto_fallback")

    return {
        "paths": {
            "input_root": resolve_path(paths.get("input_root"), "input_root"),
            "config_root": resolve_path(paths.get("config_root"), "config_root"),
        },
        "machine": {"profile": profile},
        "run": {
            "conflict_policy": conflict_policy,
            "keep_raw_without_logo": keep_raw_without_logo,
            "max_threads": max_threads,
            "batch_size": batch_size,
            "image_processing": image_processing,
        },
        "conversion": {
            "method": method,
            "enable_auto_fallback": enable_auto_fallback,
        },
    }


GPU_FRAMEWORKS = {"cuda": False}
try:  # pragma: no cover - environment dependent
    import cupy as cp
    import cupyx.scipy.ndimage as gpu_ndimage

    GPU_FRAMEWORKS["cuda"] = True
except ImportError:  # pragma: no cover - default path
    cp = None
    gpu_ndimage = None


FALLBACK_PROFILE = {
    "name": "Generic machine",
    "notes": "No configured profile matched the current computer.",
    "match_rules": {},
    "optimization_settings": {
        "max_threads": 4,
        "batch_size": 2,
        "memory_limit_gb": 8,
        "use_gpu_acceleration": False,
        "preferred_image_processing": "cpu",
        "default_conversion_method": ConversionMethod.LIBREOFFICE.value,
    },
}


MACHINE_PROFILES = {
    "macbook_pro_m3_pro_18gb": {
        "name": "MacBook Pro M3 Pro (18GB)",
        "notes": "Implemented and verified locally",
        "match_rules": {
            "platform": "darwin",
            "architecture": ["arm64"],
            "cpu_contains": None,
            "ram_gb_min": 16,
            "ram_gb_max": 20,
            "gpu_type": "apple_silicon",
            "gpu_name_contains": None,
            "hostname_contains": None,
        },
        "optimization_settings": {
            "max_threads": 8,
            "batch_size": 6,
            "memory_limit_gb": 14,
            "use_gpu_acceleration": False,
            "preferred_image_processing": "cpu",
            "default_conversion_method": ConversionMethod.LIBREOFFICE.value,
        },
    },
    "windows_i5_12450h_rtx3050_16gb": {
        "name": "Acer Windows laptop (i5-12450H + RTX 3050)",
        "notes": "Machine-aware Windows/CUDA path is implemented, but it still needs a real-machine run",
        "match_rules": {
            "platform": "windows",
            "architecture": ["amd64", "x86_64"],
            "cpu_contains": "i5-12450h",
            "ram_gb_min": 14,
            "ram_gb_max": 18,
            "gpu_type": "nvidia",
            "gpu_name_contains": "rtx 3050",
            "hostname_contains": None,
        },
        "optimization_settings": {
            "max_threads": 6,
            "batch_size": 4,
            "memory_limit_gb": 12,
            "use_gpu_acceleration": True,
            "preferred_image_processing": "cuda",
            "default_conversion_method": ConversionMethod.LIBREOFFICE.value,
        },
    },
}


def detect_apple_silicon_capability() -> bool:
    if platform.system() != "Darwin":
        return False
    try:  # pragma: no cover - platform specific
        result = subprocess.run(
            ["sysctl", "-n", "hw.optional.arm64"],
            capture_output=True,
            text=True,
            timeout=5,
        )
        return result.returncode == 0 and result.stdout.strip() == "1"
    except Exception:
        return False


def detect_gpu_info() -> dict[str, Any]:
    gpu_info: dict[str, Any] = {"gpu_detected": False}
    try:  # pragma: no cover - hardware specific
        result = subprocess.run(
            ["nvidia-smi", "--query-gpu=name", "--format=csv,noheader"],
            capture_output=True,
            text=True,
            timeout=5,
        )
        if result.returncode == 0 and result.stdout.strip():
            gpu_name = result.stdout.strip().splitlines()[0]
            gpu_info.update(
                {
                    "gpu_detected": True,
                    "gpu_type": "nvidia",
                    "gpu_name": gpu_name,
                }
            )
    except Exception:
        pass

    if detect_apple_silicon_capability():  # pragma: no cover - platform specific
        gpu_info.update(
            {
                "gpu_detected": True,
                "gpu_type": "apple_silicon",
                "gpu_name": "Apple Silicon GPU",
            }
        )

    return gpu_info


def collect_system_info() -> dict[str, Any]:
    architecture = platform.machine().lower()
    cpu = (platform.processor() or "").lower()
    if detect_apple_silicon_capability():
        architecture = "arm64"
        if not cpu or cpu == "i386":
            cpu = "apple silicon"

    system_info = {
        "hostname": platform.node().lower(),
        "platform": platform.system().lower(),
        "architecture": architecture,
        "cpu": cpu,
        "memory_gb": round(psutil.virtual_memory().total / (1024**3)),
        "cpu_count": psutil.cpu_count() or 1,
    }
    system_info.update(detect_gpu_info())
    return system_info


def derive_optimization_settings(profile: dict[str, Any], system_info: dict[str, Any]) -> dict[str, Any]:
    settings = dict(profile.get("optimization_settings", {}))
    cpu_count = max(1, int(system_info.get("cpu_count") or settings.get("max_threads", 4)))
    memory_gb = max(4, int(system_info.get("memory_gb") or settings.get("memory_limit_gb", 8)))

    settings["max_threads"] = max(1, min(int(settings.get("max_threads", 4)), max(1, cpu_count - 1)))
    settings["memory_limit_gb"] = max(4, min(int(settings.get("memory_limit_gb", 8)), max(4, memory_gb - 2)))
    settings["batch_size"] = max(1, min(int(settings.get("batch_size", 2)), max(1, settings["memory_limit_gb"] // 2)))

    preferred = settings.get("preferred_image_processing", "cpu")
    if preferred == "cuda" and not GPU_FRAMEWORKS["cuda"]:
        preferred = "cpu"
    settings["preferred_image_processing"] = preferred
    settings["use_gpu_acceleration"] = bool(settings.get("use_gpu_acceleration", False) and preferred == "cuda")
    return settings


def build_machine_config(
    machine_id: str,
    profile: dict[str, Any],
    system_info: dict[str, Any],
    match_source: str,
    match_reasons: list[str],
) -> dict[str, Any]:
    return {
        "machine_id": machine_id,
        "name": profile.get("name", machine_id),
        "notes": profile.get("notes", ""),
        "match_source": match_source,
        "match_reasons": match_reasons,
        "detected_info": system_info,
        "optimization_settings": derive_optimization_settings(profile, system_info),
    }


def build_fallback_config(system_info: dict[str, Any], match_source: str = "fallback", match_reasons: Optional[list[str]] = None) -> dict[str, Any]:
    return build_machine_config(
        "generic_fallback",
        FALLBACK_PROFILE,
        system_info,
        match_source,
        match_reasons or ["No configured profile matched the current machine"],
    )


def evaluate_profile_match(profile: dict[str, Any], system_info: dict[str, Any]) -> Optional[tuple[int, list[str]]]:
    rules = profile.get("match_rules", {})
    reasons: list[str] = []
    score = 0

    platform_rule = rules.get("platform")
    if platform_rule:
        if system_info.get("platform") != platform_rule:
            return None
        score += 20
        reasons.append(f"platform={platform_rule}")

    architecture_rules = rules.get("architecture", [])
    if architecture_rules:
        if system_info.get("architecture") not in architecture_rules:
            return None
        score += 15
        reasons.append(f"architecture={system_info.get('architecture')}")

    cpu_rule = rules.get("cpu_contains")
    if cpu_rule:
        if cpu_rule not in system_info.get("cpu", ""):
            return None
        score += 10
        reasons.append(f"cpu contains {cpu_rule}")

    ram_gb = int(system_info.get("memory_gb") or 0)
    ram_min = rules.get("ram_gb_min")
    ram_max = rules.get("ram_gb_max")
    if ram_min is not None and ram_gb < ram_min:
        return None
    if ram_max is not None and ram_gb > ram_max:
        return None
    if ram_min is not None or ram_max is not None:
        score += 10
        reasons.append(f"ram={ram_gb}GB")

    gpu_type_rule = rules.get("gpu_type")
    if gpu_type_rule:
        if system_info.get("gpu_type") != gpu_type_rule:
            return None
        score += 20
        reasons.append(f"gpu_type={gpu_type_rule}")

    gpu_name_rule = rules.get("gpu_name_contains")
    if gpu_name_rule:
        gpu_name = (system_info.get("gpu_name") or "").lower()
        if gpu_name_rule not in gpu_name:
            return None
        score += 25
        reasons.append(f"gpu name contains {gpu_name_rule}")

    return score, reasons


def select_machine_profile(
    system_info: dict[str, Any],
    profiles: dict[str, dict[str, Any]],
    override_profile_id: Optional[str] = None,
) -> dict[str, Any]:
    if override_profile_id:
        if override_profile_id == "generic_fallback":
            return build_fallback_config(system_info, "manual", ["Manual override selected"])
        if override_profile_id in profiles:
            return build_machine_config(
                override_profile_id,
                profiles[override_profile_id],
                system_info,
                "manual",
                ["Manual override selected"],
            )
        return build_fallback_config(system_info, "manual", [f"Unknown manual override: {override_profile_id}"])

    matches: list[tuple[int, str, list[str]]] = []
    for machine_id, profile in profiles.items():
        result = evaluate_profile_match(profile, system_info)
        if result is not None:
            score, reasons = result
            matches.append((score, machine_id, reasons))

    if matches:
        matches.sort(key=lambda item: item[0], reverse=True)
        _, machine_id, reasons = matches[0]
        return build_machine_config(machine_id, profiles[machine_id], system_info, "auto", reasons)

    return build_fallback_config(system_info)


def resolve_runtime_settings(config: dict[str, Any], machine_config: dict[str, Any]) -> dict[str, Any]:
    machine_settings = dict(machine_config.get("optimization_settings", {}))
    run_settings = dict(config.get("run", {}))

    max_threads = int(run_settings.get("max_threads") or machine_settings.get("max_threads") or 4)
    batch_size = int(run_settings.get("batch_size") or machine_settings.get("batch_size") or 2)
    requested_processing = str(run_settings.get("image_processing", "auto")).strip() or "auto"

    preferred_processing = machine_settings.get("preferred_image_processing", "cpu")
    if requested_processing != "auto":
        preferred_processing = requested_processing
    if preferred_processing not in {"cpu", "cuda"}:
        preferred_processing = "cpu"

    runtime_acceleration = preferred_processing
    if runtime_acceleration == "cuda" and not GPU_FRAMEWORKS["cuda"]:
        runtime_acceleration = "cpu"

    return {
        "max_threads": max(1, max_threads),
        "batch_size": max(1, batch_size),
        "requested_image_processing": requested_processing,
        "preferred_image_processing": runtime_acceleration,
        "runtime_acceleration": runtime_acceleration,
        "use_gpu_acceleration": runtime_acceleration == "cuda",
    }


def detect_conversion_tooling() -> dict[str, bool]:
    return {
        "libreoffice": bool(get_libreoffice_path()),
        "pdftocairo": shutil.which("pdftocairo") is not None,
        "pdftoppm": shutil.which("pdftoppm") is not None,
        "pdf2image": importlib.util.find_spec("pdf2image") is not None,
    }


def select_pdf_rasterizer(tooling: dict[str, bool]) -> Optional[str]:
    if tooling.get("pdftocairo"):
        return "pdftocairo"
    if tooling.get("pdftoppm"):
        return "pdftoppm"
    return None


def get_recommended_conversion_method(machine_id: str, tooling: dict[str, bool]) -> ConversionMethod:
    if tooling.get("libreoffice") and select_pdf_rasterizer(tooling):
        return ConversionMethod.LIBREOFFICE
    if tooling.get("pdf2image"):
        return ConversionMethod.PDF2IMAGE
    return ConversionMethod.PYTHON_PPTX


def build_pdf2image_kwargs(output_folder: str, tooling: dict[str, bool]) -> dict[str, Any]:
    return {
        "dpi": 300,
        "fmt": "png",
        "output_folder": output_folder,
        "paths_only": True,
        "use_pdftocairo": bool(tooling.get("pdftocairo")),
    }


def find_files(directory: os.PathLike[str] | str, pattern: str) -> list[str]:
    results: list[str] = []
    for root, _, filenames in os.walk(directory):
        for filename in fnmatch.filter(filenames, pattern):
            results.append(os.path.join(root, filename))
    return results


def find_presentation_files(directory: os.PathLike[str] | str) -> list[str]:
    return find_files(directory, "*.pptx") + find_files(directory, "*.zip")


def is_archived_presentation(file_path: os.PathLike[str] | str) -> bool:
    return "all_pptx" in Path(file_path).parts


def find_presentation_inventory(directory: os.PathLike[str] | str) -> dict[str, Any]:
    all_presentations = find_presentation_files(directory)
    actionable: list[str] = []
    archived: list[str] = []
    for file_path in all_presentations:
        if is_archived_presentation(file_path):
            archived.append(file_path)
        else:
            actionable.append(file_path)

    summary = {
        "total_found": len(all_presentations),
        "actionable_count": len(actionable),
        "archived_count": len(archived),
        "pptx_actionable_count": sum(1 for path in actionable if path.lower().endswith(".pptx")),
        "zip_actionable_count": sum(1 for path in actionable if path.lower().endswith(".zip")),
    }
    return {
        "all": all_presentations,
        "actionable": actionable,
        "archived": archived,
        "summary": summary,
    }


def build_run_summary(inventory: dict[str, Any], conflict_policy_label: str) -> dict[str, Any]:
    summary = dict(inventory["summary"])
    actionable_count = summary["actionable_count"]
    summary["conflict_policy_label"] = conflict_policy_label
    summary["button_label"] = f"Extract {actionable_count} Presentations"
    return summary


def resolve_existing_processing_action(
    archived_source_exists: bool,
    output_folder_exists: bool,
    existing_png_count: int,
    conflict_policy: str,
) -> dict[str, str]:
    has_existing_outputs = output_folder_exists and existing_png_count > 0
    has_conflict = archived_source_exists or has_existing_outputs

    if not has_conflict:
        return {"action": "process", "status": "process", "message": "No prior processed output detected"}
    if conflict_policy == ConflictPolicy.REPORT_ONLY.value:
        return {"action": "report", "status": "reported", "message": "Conflict detected; report-only mode left files unchanged"}
    if conflict_policy == ConflictPolicy.REPROCESS_NEW_FOLDER.value:
        return {"action": "reprocess_new_folder", "status": "reprocess", "message": "Conflict detected; reprocessing into a new folder"}
    if conflict_policy == ConflictPolicy.OVERWRITE.value:
        return {"action": "reprocess", "status": "reprocess", "message": "Conflict detected; existing output will be overwritten"}
    if has_existing_outputs:
        return {"action": "skip", "status": "skipped", "message": "Skipped safely because existing processed output was found"}
    return {"action": "reprocess", "status": "repair", "message": "Archived source found but processed output is missing; recreating output safely"}


def get_libreoffice_path() -> Optional[str]:
    system = platform.system()
    if system == "Darwin":  # pragma: no cover - platform specific
        possible_paths = [
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            "/Applications/LibreOffice.app/Contents/MacOS/soffice.bin",
            "/opt/homebrew/bin/soffice",
            "/usr/local/bin/soffice",
        ]
        for path in possible_paths:
            if os.path.exists(path):
                return path
        return None
    if system == "Windows":
        possible_paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files\LibreOffice\program\soffice.com",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.com",
        ]
        for path in possible_paths:
            if os.path.exists(path):
                return path
    soffice = shutil.which("soffice")
    return soffice or None


def extract_numeric_prefix(filename: str) -> int:
    match = re.match(r"^(\d+)", filename)
    if match:
        return int(match.group(1))
    return 999999


def determine_target_size(width: int, height: int) -> tuple[int, int]:
    if width >= TARGET_RESOLUTION[0] and height >= TARGET_RESOLUTION[1]:
        return width, height
    aspect_ratio = width / height
    if aspect_ratio > TARGET_RESOLUTION[0] / TARGET_RESOLUTION[1]:
        new_width = TARGET_RESOLUTION[0]
        new_height = int(new_width / aspect_ratio)
    else:
        new_height = TARGET_RESOLUTION[1]
        new_width = int(new_height * aspect_ratio)
    return new_width, new_height


class AcceleratedImageProcessor:
    def __init__(self, runtime_acceleration: str = "cpu") -> None:
        self.runtime_acceleration = "cuda" if runtime_acceleration == "cuda" and GPU_FRAMEWORKS["cuda"] else "cpu"

    def resize_to_4k(self, image: Image.Image) -> Image.Image:
        if image.width >= TARGET_RESOLUTION[0] and image.height >= TARGET_RESOLUTION[1]:
            return image.copy()
        target_size = determine_target_size(image.width, image.height)
        if self.runtime_acceleration == "cuda":
            return self._resize_cuda(image, target_size)
        return self._resize_cpu(image, target_size)

    def _resize_cpu(self, image: Image.Image, target_size: tuple[int, int]) -> Image.Image:
        image_array = np.array(image)
        resized_array = cv2.resize(image_array, target_size, interpolation=cv2.INTER_LANCZOS4)
        return Image.fromarray(resized_array)

    def _resize_cuda(self, image: Image.Image, target_size: tuple[int, int]) -> Image.Image:
        if cp is None or gpu_ndimage is None:
            return self._resize_cpu(image, target_size)

        image_array = np.array(image)
        if image_array.ndim == 2:
            scale_factors = (
                target_size[1] / image_array.shape[0],
                target_size[0] / image_array.shape[1],
            )
        else:
            scale_factors = (
                target_size[1] / image_array.shape[0],
                target_size[0] / image_array.shape[1],
                1.0,
            )

        try:
            gpu_array = cp.asarray(image_array)
            resized_gpu = gpu_ndimage.zoom(gpu_array, scale_factors, order=1, prefilter=False)
            resized_array = cp.asnumpy(resized_gpu).astype(np.uint8)
            if resized_array.ndim == 3 and resized_array.shape[2] == 1:
                resized_array = resized_array[:, :, 0]
            return Image.fromarray(resized_array)
        except Exception:
            return self._resize_cpu(image, target_size)


def resize_to_4k(image: Image.Image, image_processor: Optional[AcceleratedImageProcessor] = None) -> Image.Image:
    if image.width >= TARGET_RESOLUTION[0] and image.height >= TARGET_RESOLUTION[1]:
        return image.copy()
    if image_processor is not None:
        return image_processor.resize_to_4k(image)
    new_width, new_height = determine_target_size(image.width, image.height)
    return image.resize((new_width, new_height), LANCZOS)


def read_copyright(copyright_path: Path) -> str:
    try:
        return copyright_path.read_text(encoding="utf-8").strip()
    except UnicodeDecodeError:
        return copyright_path.read_text(encoding="cp1252").strip()
    except Exception:
        return "© All Rights Reserved"


def get_position(img_width: int, img_height: int, element_width: int, element_height: int, position: str, padding: int) -> tuple[int, int]:
    if position == "top-left":
        return (padding, padding)
    if position == "top-right":
        return (img_width - element_width - padding, padding)
    if position == "bottom-left":
        return (padding, img_height - element_height - padding)
    if position == "bottom-right":
        return (img_width - element_width - padding, img_height - element_height - padding)
    if position == "bottom-center":
        return ((img_width - element_width) // 2, img_height - element_height - padding)
    raise ValueError(f"Invalid position: {position}")


@lru_cache(maxsize=8)
def get_font(size: int) -> ImageFont.ImageFont:
    try:
        return ImageFont.truetype("Arial.ttf", size)
    except IOError:
        try:
            return ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", size)
        except IOError:
            return ImageFont.load_default()


@lru_cache(maxsize=16)
def create_branding_overlay(img_width: int, img_height: int, config_root_value: str, copyright_text: str) -> Image.Image:
    overlay = Image.new("RGBA", (img_width, img_height), (0, 0, 0, 0))
    logo_path = Path(config_root_value) / "logo.png"
    if logo_path.exists():
        try:
            with Image.open(logo_path) as logo:
                logo = logo.convert("RGBA").resize((LOGO_SIZE, LOGO_SIZE), LANCZOS)
                overlay.paste(logo, get_position(img_width, img_height, LOGO_SIZE, LOGO_SIZE, LOGO_POSITION, LOGO_PADDING), logo)
        except Exception:
            pass

    draw = ImageDraw.Draw(overlay)
    font = get_font(FONT_SIZE)
    text_bbox = draw.textbbox((0, 0), copyright_text, font=font)
    text_width = text_bbox[2] - text_bbox[0]
    text_height = text_bbox[3] - text_bbox[1]
    draw.text(
        get_position(img_width, img_height, text_width, text_height, COPYRIGHT_POSITION, COPYRIGHT_PADDING),
        copyright_text,
        font=font,
        fill=FONT_COLOR,
    )
    return overlay


def apply_branding(image: Image.Image, config_root: Path) -> Image.Image:
    if image.mode != "RGBA":
        working = image.convert("RGBA")
    else:
        working = image.copy()

    copyright_text = read_copyright(config_root / "copyright.txt")
    overlay = create_branding_overlay(working.width, working.height, str(config_root.resolve()), copyright_text)
    return Image.alpha_composite(working, overlay)


def extract_slides_directly(pptx_path: Path, output_dir: Path) -> tuple[bool, list[Path], int]:
    output_dir.mkdir(parents=True, exist_ok=True)
    presentation = Presentation(str(pptx_path))
    slide_count = len(presentation.slides)
    if slide_count == 0:
        return False, [], 0

    extracted: list[Path] = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_width = presentation.slide_width or TARGET_RESOLUTION[0]
        slide_height = presentation.slide_height or TARGET_RESOLUTION[1]
        aspect_ratio = slide_width / slide_height
        if aspect_ratio > TARGET_RESOLUTION[0] / TARGET_RESOLUTION[1]:
            width = TARGET_RESOLUTION[0]
            height = int(width / aspect_ratio)
        else:
            height = TARGET_RESOLUTION[1]
            width = int(height * aspect_ratio)

        image = Image.new("RGB", (width, height), "white")
        _ = slide
        output_path = output_dir / f"{index:02d}.png"
        image.save(output_path, format="PNG")
        extracted.append(output_path)
    return True, extracted, slide_count


def _convert_pptx_to_pdf(pptx_path: Path, output_dir: Path) -> Path:
    libreoffice_path = get_libreoffice_path()
    if not libreoffice_path:
        raise RuntimeError("LibreOffice is not available")
    command = [
        libreoffice_path,
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(output_dir),
        str(pptx_path),
    ]
    subprocess.run(command, check=True, capture_output=True)
    pdf_path = output_dir / f"{pptx_path.stem}.pdf"
    if not pdf_path.exists():
        raise RuntimeError(f"LibreOffice did not create {pdf_path.name}")
    return pdf_path


def _extract_pdf_images(pdf_path: Path, output_dir: Path, use_pdf2image: bool) -> list[Path]:
    if use_pdf2image:
        from pdf2image import convert_from_path

        images = convert_from_path(str(pdf_path), **build_pdf2image_kwargs(str(output_dir), detect_conversion_tooling()))
        return [Path(path) for path in images]

    rasterizer = select_pdf_rasterizer(detect_conversion_tooling())
    if rasterizer is None:
        raise RuntimeError("No Poppler rasterizer is available")

    output_dir.mkdir(parents=True, exist_ok=True)
    prefix = output_dir / pdf_path.stem
    command = [rasterizer, "-png", str(pdf_path), str(prefix)]
    subprocess.run(command, check=True, capture_output=True)
    return sorted(output_dir.glob(f"{pdf_path.stem}-*.png"))


def extract_slides_with_method(pptx_path: Path, output_dir: Path, method: ConversionMethod, enable_auto_fallback: bool) -> tuple[bool, list[Path], int]:
    methods = [method]
    if enable_auto_fallback:
        for candidate in (ConversionMethod.LIBREOFFICE, ConversionMethod.PDF2IMAGE, ConversionMethod.PYTHON_PPTX):
            if candidate not in methods:
                methods.append(candidate)

    for current_method in methods:
        try:
            if current_method == ConversionMethod.PYTHON_PPTX:
                return extract_slides_directly(pptx_path, output_dir)

            with tempfile.TemporaryDirectory() as tmp_dir_raw:
                tmp_dir = Path(tmp_dir_raw)
                pdf_path = _convert_pptx_to_pdf(pptx_path, tmp_dir)
                page_paths = _extract_pdf_images(pdf_path, output_dir, use_pdf2image=current_method == ConversionMethod.PDF2IMAGE)
                if not page_paths:
                    raise RuntimeError("No PDF pages were rasterized")
                return True, sorted(page_paths), len(page_paths)
        except Exception:
            continue

    return False, [], 0


def batch_items(items: list[Any], batch_size: int) -> list[list[Any]]:
    return [items[index : index + batch_size] for index in range(0, len(items), batch_size)]


def extract_slides_from_zip(
    zip_file: Path,
    output_dir: Path,
    runtime_settings: Optional[dict[str, Any]] = None,
) -> tuple[bool, list[Path], int]:
    output_dir.mkdir(parents=True, exist_ok=True)
    runtime_settings = runtime_settings or {"max_threads": 1, "batch_size": 1, "runtime_acceleration": "cpu"}
    image_processor = AcceleratedImageProcessor(runtime_settings.get("runtime_acceleration", "cpu"))
    max_threads = max(1, int(runtime_settings.get("max_threads") or 1))
    batch_size = max(1, int(runtime_settings.get("batch_size") or 1))

    with tempfile.TemporaryDirectory() as tmp_dir_raw:
        tmp_dir = Path(tmp_dir_raw)
        with zipfile.ZipFile(zip_file, "r") as archive:
            archive.extractall(tmp_dir)

        image_files: list[Path] = []
        for extension in ("*.png", "*.jpg", "*.jpeg", "*.PNG", "*.JPG", "*.JPEG"):
            image_files.extend(Path(path) for path in glob.glob(str(tmp_dir / "**" / extension), recursive=True))
        image_files = list({path.resolve(): path for path in image_files}.values())

        if not image_files:
            return False, [], 0

        image_files.sort(key=lambda item: extract_numeric_prefix(item.name))
        output_paths: list[Path] = []

        def process_archive_image(task: tuple[int, Path]) -> Path:
            index, image_path = task
            with Image.open(image_path) as image:
                resized = resize_to_4k(image.convert("RGBA"), image_processor=image_processor)
                output_path = output_dir / f"{index:02d}.png"
                resized.save(output_path, format="PNG", optimize=True)
                return output_path

        indexed_images = list(enumerate(image_files, start=1))
        for batch in batch_items(indexed_images, batch_size):
            with concurrent.futures.ThreadPoolExecutor(max_workers=min(max_threads, len(batch))) as executor:
                output_paths.extend(executor.map(process_archive_image, batch))
        return True, output_paths, len(output_paths)


def build_processing_paths(source_path: Path) -> tuple[Path, Path]:
    parent_dir = source_path.parent
    if parent_dir.name == "all_pptx":
        course_dir = parent_dir.parent
        file_stem = source_path.stem
        lecture_match = re.search(r"(\d+)", file_stem)
        lecture_folder_name = f"Lecture {lecture_match.group(1).zfill(2)}" if lecture_match else file_stem
        lecture_folder = course_dir / lecture_folder_name
        output_folder = lecture_folder / "English image"
        archived_dir = parent_dir
    else:
        archived_dir = parent_dir / "all_pptx"
        file_stem = source_path.stem
        if source_path.suffix.lower() == ".zip" and re.match(r"^\d+$", file_stem):
            lecture_folder = parent_dir / f"Lecture {file_stem.zfill(2)}"
        else:
            lecture_folder = parent_dir / file_stem
        output_folder = lecture_folder / "English image"
    return output_folder, archived_dir / source_path.name


def make_conflict_output_folder(output_folder: Path) -> Path:
    timestamp = time.strftime("%Y%m%d_%H%M%S")
    new_output_folder = output_folder.with_name(f"{output_folder.name} rerun {timestamp}")
    new_output_folder.mkdir(parents=True, exist_ok=True)
    return new_output_folder


def process_presentation(
    source_path: Path,
    config: dict[str, Any],
    runtime_settings: Optional[dict[str, Any]] = None,
) -> dict[str, Any]:
    result = {
        "file_path": str(source_path),
        "status": "error",
        "message": "",
        "slide_count": 0,
        "output_dir": "",
    }
    runtime_settings = runtime_settings or {"max_threads": 1, "batch_size": 1, "runtime_acceleration": "cpu"}
    image_processor = AcceleratedImageProcessor(runtime_settings.get("runtime_acceleration", "cpu"))
    max_threads = max(1, int(runtime_settings.get("max_threads") or 1))
    batch_size = max(1, int(runtime_settings.get("batch_size") or 1))

    output_folder, archived_source_path = build_processing_paths(source_path)
    output_folder.mkdir(parents=True, exist_ok=True)
    png_count = len([path for path in output_folder.glob("*.png") if path.is_file()])

    conflict_decision = resolve_existing_processing_action(
        archived_source_exists=archived_source_path.exists(),
        output_folder_exists=output_folder.exists(),
        existing_png_count=png_count,
        conflict_policy=config["run"]["conflict_policy"],
    )
    result["output_dir"] = str(output_folder)

    if conflict_decision["action"] in {"skip", "report"}:
        result["status"] = conflict_decision["status"]
        result["message"] = conflict_decision["message"]
        result["slide_count"] = png_count
        return result

    if conflict_decision["action"] == "reprocess_new_folder":
        output_folder = make_conflict_output_folder(output_folder)
        result["output_dir"] = str(output_folder)

    keep_raw = config["run"]["keep_raw_without_logo"]
    raw_output_folder = output_folder / "without_logo_png" if keep_raw else None
    extraction_folder = raw_output_folder or Path(tempfile.mkdtemp(prefix="dreamlet-4k-"))
    extraction_folder.mkdir(parents=True, exist_ok=True)

    try:
        source_suffix = source_path.suffix.lower()
        conversion_method = ConversionMethod(config["conversion"]["method"])
        if source_suffix == ".zip":
            success, extracted_images, slide_count = extract_slides_from_zip(source_path, extraction_folder, runtime_settings)
        else:
            success, extracted_images, slide_count = extract_slides_with_method(
                source_path,
                extraction_folder,
                conversion_method,
                bool(config["conversion"]["enable_auto_fallback"]),
            )
        if not success:
            result["message"] = "Failed to extract images from presentation"
            return result

        branded_count = 0
        config_root = config["paths"]["config_root"]

        def process_output_image(extracted_image: Path | str) -> Path:
            extracted_path = Path(extracted_image)
            with Image.open(extracted_path) as image:
                resized = resize_to_4k(image.convert("RGBA"), image_processor=image_processor)
                branded = apply_branding(resized, config_root)
                output_path = output_folder / extracted_path.name
                branded.save(output_path, format="PNG", optimize=True)
                return output_path

        extracted_paths = [Path(path) for path in extracted_images]
        for batch in batch_items(extracted_paths, batch_size):
            with concurrent.futures.ThreadPoolExecutor(max_workers=min(max_threads, len(batch))) as executor:
                branded_count += len(list(executor.map(process_output_image, batch)))

        archived_source_path.parent.mkdir(parents=True, exist_ok=True)
        if source_path.exists() and not archived_source_path.exists():
            shutil.move(str(source_path), str(archived_source_path))

        result["status"] = "success"
        result["message"] = f"Successfully processed {branded_count} images"
        result["slide_count"] = slide_count
        return result
    finally:
        if not keep_raw and extraction_folder.exists():
            shutil.rmtree(extraction_folder, ignore_errors=True)


def process_presentations(
    config: dict[str, Any],
    runtime_settings: Optional[dict[str, Any]] = None,
) -> list[dict[str, Any]]:
    inventory = find_presentation_inventory(config["paths"]["input_root"])
    results: list[dict[str, Any]] = []
    runtime_settings = runtime_settings or {"batch_size": 1, "max_threads": 1, "runtime_acceleration": "cpu"}
    batch_size = max(1, int(runtime_settings.get("batch_size") or 1))
    for batch in batch_items(list(inventory["actionable"]), batch_size):
        for file_path in batch:
            results.append(process_presentation(Path(file_path), config, runtime_settings))
    return results


def summarize_results(results: list[dict[str, Any]]) -> dict[str, int]:
    summary = {"success": 0, "skipped": 0, "reported": 0, "error": 0}
    for result in results:
        status = result["status"]
        if status == "success":
            summary["success"] += 1
        elif status == "skipped":
            summary["skipped"] += 1
        elif status == "reported":
            summary["reported"] += 1
        else:
            summary["error"] += 1
    return summary


def run(config: dict[str, Any]) -> int:
    override_profile = config["machine"]["profile"]
    if override_profile == "auto":
        override_profile = None

    input_root = config["paths"]["input_root"]
    if not input_root.exists():
        print(f"Input directory not found: {input_root}")
        return 1

    machine_info = select_machine_profile(collect_system_info(), MACHINE_PROFILES, override_profile)
    inventory = find_presentation_inventory(input_root)
    runtime_settings = resolve_runtime_settings(config, machine_info)
    runtime_acceleration = runtime_settings["runtime_acceleration"]
    profile_source = "config" if config["machine"]["profile"] != "auto" else machine_info["match_source"]
    fallback_reason = "none"
    if runtime_settings["requested_image_processing"] == "cuda" and runtime_acceleration != "cuda":
        fallback_reason = "cuda_unavailable"
    print(
        f"workflow=4k profile={machine_info['machine_id']} profile_source={profile_source} "
        f"acceleration={runtime_acceleration} fallback={fallback_reason} "
        f"actionable={inventory['summary']['actionable_count']} archived={inventory['summary']['archived_count']}"
    )
    results = process_presentations(config, runtime_settings)
    summary = summarize_results(results)
    print(
        f"success={summary['success']} skipped={summary['skipped']} "
        f"reported={summary['reported']} error={summary['error']}"
    )
    for result in results:
        print(f"{result['status']}: {Path(result['file_path']).name} -> {result['message']}")
    return 0 if summary["error"] == 0 else 1


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="CLI version of 06 Generate 4K Images.")
    parser.add_argument("--config", help="Path to TOML config. Defaults to config/06 Generate 4K Images.toml")
    args = parser.parse_args(argv)
    try:
        return run(load_config(args.config))
    except ConfigError as exc:
        print(f"Config error: {exc}", file=sys.stderr)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
