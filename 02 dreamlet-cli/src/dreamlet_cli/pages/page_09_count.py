
"""
CODING CONVENTION: NO SHARED CODE
- All code for this page must be contained entirely within this single file
- Never import from other page files or create shared utilities
- Copy any needed functions directly into this file
- Each page is completely self-contained and independent

STATUS: LEGACY
PURPOSE: Older lecture artifact verification and count page.
MAIN INPUTS:
- lecture folders and generated assets under `input/`
MAIN OUTPUTS:
- count and discrepancy summaries shown in the UI
REQUIRED CONFIG / ASSETS:
- `input/` directory
EXTERNAL SERVICES:
- optional document-processing tools used during validation
HARDWARE ASSUMPTIONS:
- none
REPLACED BY:
- `pages/08 Validate File Counts.py`
"""

from dreamlet_cli.compat import st
import os
import re
import time
import glob
import fnmatch
import pandas as pd
import subprocess
import tempfile
from typing import Dict, List, Tuple, Optional, Any
import pptx
from enum import Enum
from pathlib import Path

# Local utility functions (moved from multiple utils modules)
def get_input_directory() -> str:
    """Get the path to the input directory"""
    return os.path.join(os.getcwd(), "input")

def find_files(directory: str, pattern: str) -> List[str]:
    """Find all files matching a pattern in a directory (recursively)"""
    result = []
    for root, _, filenames in os.walk(directory):
        for filename in fnmatch.filter(filenames, pattern):
            result.append(os.path.join(root, filename))
    return result

def find_transcript_files(directory: str) -> List[str]:
    """Find all transcript files in a directory"""
    transcripts = []
    patterns = ["Lecture*.txt", "Lecture*.md", "*lecture*.txt", "*lecture*.md", "*transcript*.txt", "*transcript*.md"]
    for pattern in patterns:
        transcripts.extend(find_files(directory, pattern))
    return list(set([f for f in transcripts if "slide" not in os.path.basename(f).lower()]))

def find_slide_files(directory: str) -> List[str]:
    """Find all slide description files in a directory"""
    slides = []
    patterns = ["*-slides.txt", "*-slides.md", "*slide*.txt", "*slide*.md"]
    for pattern in patterns:
        slides.extend(find_files(directory, pattern))
    return list(set(slides))

def find_presentation_files(directory: str) -> List[str]:
    """Find all presentation files in a directory"""
    return find_files(directory, "*.pptx")

def find_pptx_for_lecture(course_dir: str, lecture_num: str) -> Optional[str]:
    """
    Find PPTX file for a specific lecture, checking multiple locations
    
    Args:
        course_dir: Course directory path
        lecture_num: Lecture number (e.g., "01", "1")
        
    Returns:
        Path to PPTX file if found, None otherwise
    """
    # Standardize lecture number
    std_lecture = lecture_num.zfill(2)
    
    # Possible PPTX filenames
    possible_names = [
        f"Lecture {std_lecture}.pptx",
        f"Lecture{std_lecture}.pptx", 
        f"lecture {std_lecture}.pptx",
        f"lecture{std_lecture}.pptx",
        f"Lecture {int(lecture_num)}.pptx",
        f"lecture {int(lecture_num)}.pptx"
    ]
    
    # Check in lecture-specific directory first
    lecture_dir = os.path.join(course_dir, f"Lecture {std_lecture}")
    if os.path.exists(lecture_dir):
        for name in possible_names:
            pptx_path = os.path.join(lecture_dir, name)
            if os.path.exists(pptx_path):
                return pptx_path
    
    # Check in course root directory
    for name in possible_names:
        pptx_path = os.path.join(course_dir, name)
        if os.path.exists(pptx_path):
            return pptx_path
    
    # Check in all_pptx folder as fallback
    all_pptx_dir = os.path.join(course_dir, "all_pptx")
    if os.path.exists(all_pptx_dir):
        for name in possible_names:
            pptx_path = os.path.join(all_pptx_dir, name)
            if os.path.exists(pptx_path):
                return pptx_path
        
        # Also check for any PPTX with the lecture number in the filename
        for file in os.listdir(all_pptx_dir):
            if file.endswith('.pptx') and (std_lecture in file or lecture_num in file):
                return os.path.join(all_pptx_dir, file)
    
    return None

def extract_course_lecture_section(file_path: str) -> Tuple[Optional[str], Optional[str], Optional[str]]:
    """Extract course, lecture, and section information from a file path"""
    dir_parts = os.path.normpath(file_path).split(os.sep)
    course, lecture, section = None, None, None
    
    for part in dir_parts:
        if not course:
            course_match = re.search(r'course\s*(\d+)', part.lower())
            if course_match:
                course = course_match.group(1)
    
    filename = os.path.basename(file_path)
    for pattern in [r'lecture\s*(\d+)', r'lec\s*(\d+)', r'^(\d+)[-\s]']:
        lecture_match = re.search(pattern, filename.lower())
        if lecture_match:
            lecture = lecture_match.group(1)
            break
    
    return course, lecture, section

def group_files_by_lecture(files: List[str]) -> Dict[str, List[str]]:
    """Group files by lecture number"""
    groups = {}
    for file_path in files:
        _, lecture, _ = extract_course_lecture_section(file_path)
        if lecture:
            key = f"Lecture {lecture.zfill(2)}"
            if key not in groups:
                groups[key] = []
            groups[key].append(file_path)
    return groups

def count_slides_in_transcript(text: str) -> int:
    """Count the number of slides in a transcript file"""
    patterns = [r'\[Slide\s+(\d+)', r'Slide\s+(\d+)', r'##\s*Slide\s+(\d+)']
    slides = set()
    for pattern in patterns:
        matches = re.findall(pattern, text, re.IGNORECASE)
        slides.update(matches)
    return len(slides)

def count_slides_in_slide_file(text: str) -> int:
    """Count the number of slides in a slide description file"""
    patterns = [r'Slide\s+\d+\s*[:-]', r'^\s*\d+\s*[:-]', r'^\s*\[Slide\s+\d+\]']
    count = 0
    for line in text.split('\n'):
        for pattern in patterns:
            if re.search(pattern, line, re.IGNORECASE):
                count += 1
                break
    return count

def count_slides_in_pptx(pptx_path: str) -> int:
    """Count slides in a PPTX presentation"""
    try:
        if not os.path.exists(pptx_path):
            return 0
        presentation = pptx.Presentation(pptx_path)
        slide_count = len(presentation.slides)
        return slide_count
    except Exception as e:
        # Try alternative method if python-pptx fails
        try:
            import zipfile
            with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
                # Count slide XML files in the presentation
                slide_files = [f for f in zip_ref.namelist() if f.startswith('ppt/slides/slide') and f.endswith('.xml')]
                return len(slide_files)
        except Exception:
            return 0

st.set_page_config(page_title="09 Count - Dreamlet", page_icon="🔢", layout="wide")

class SlideCounting:
    """Enumeration of slide counting methods"""
    PYTHON_PPTX = "python-pptx"
    LIBREOFFICE = "libreoffice"
    PDF_TOOLS = "pdf-tools"

def get_full_course_name(path: str) -> str:
    """
    Extract the full course name from the file path
    
    Args:
        path: File path
        
    Returns:
        Full course name
    """
    # Extract course from path
    parts = path.split(os.path.sep)
    
    # Look for course pattern
    for part in parts:
        if "course" in part.lower():
            return part
        if re.search(r"^\d{1,3}\s+[\w\s]+", part, re.IGNORECASE):
            return part
            
    # No course found
    return "Uncategorized"

def extract_lecture_number(path: str) -> str:
    """
    Extract lecture number from the path
    
    Args:
        path: File path
        
    Returns:
        Lecture number (as string)
    """
    # Try to extract from filename directly
    filename = os.path.basename(path)
    match = re.search(r"lecture[\_\s-]*(\d+)", filename, re.IGNORECASE)
    if match:
        return match.group(1)
    
    # Try to extract from parent directory names
    parts = path.split(os.path.sep)
    for part in parts:
        if "lecture" in part.lower():
            match = re.search(r"lecture[\_\s-]*(\d+)", part, re.IGNORECASE)
            if match:
                return match.group(1)
    
    # No lecture number found
    return "0"

def find_english_text_files(input_dir: str) -> List[str]:
    """
    Find all files in 'English text' folders
    
    Args:
        input_dir: Input directory
        
    Returns:
        List of file paths
    """
    files = []
    for root, dirs, filenames in os.walk(input_dir):
        # Skip hidden directories
        dirs[:] = [d for d in dirs if not d.startswith('.')]
        
        # Check if this is an "English text" folder
        if os.path.basename(root) == "English text":
            for filename in filenames:
                if filename.endswith(('.txt', '.md')):
                    files.append(os.path.join(root, filename))
    
    return files

def find_english_summary_text_files(input_dir: str) -> List[str]:
    """
    Find all files in 'English Summary text' folders
    
    Args:
        input_dir: Input directory
        
    Returns:
        List of file paths
    """
    files = []
    for root, dirs, filenames in os.walk(input_dir):
        # Skip hidden directories
        dirs[:] = [d for d in dirs if not d.startswith('.')]
        
        # Check if this is an "English Summary text" folder
        if os.path.basename(root) == "English Summary text":
            for filename in filenames:
                if filename.endswith(('.txt', '.md')):
                    files.append(os.path.join(root, filename))
    
    return files

def count_files_in_directory(directory: str, extensions: List[str] = ['.txt', '.md']) -> int:
    """
    Count number of files with specified extensions in a directory
    
    Args:
        directory: Directory to search in
        extensions: File extensions to count
        
    Returns:
        Number of files
    """
    if not os.path.exists(directory):
        return 0
    
    count = 0
    for filename in os.listdir(directory):
        if any(filename.lower().endswith(ext) for ext in extensions):
            count += 1
    
    return count

def count_slides_in_pptx_advanced(
    pptx_path: str, 
    status_text: Optional[Any] = None
) -> Tuple[bool, int, str]:
    """
    Enhanced version of slide counting with multiple methods and detailed reporting
    
    Args:
        pptx_path: Path to PowerPoint file
        status_text: Streamlit text element for status updates
        
    Returns:
        Tuple of (success, slide_count, method_used)
    """
    # Check if file exists
    if not os.path.exists(pptx_path):
        if status_text:
            status_text.error(f"File not found: {pptx_path}")
        return False, 0, ""
    
    # Method 1: Use python-pptx (fastest and most reliable)
    try:
        prs = pptx.Presentation(pptx_path)
        slide_count = len(prs.slides)
        return True, slide_count, SlideCounting.PYTHON_PPTX
    except Exception as e:
        if status_text:
            status_text.warning(f"Python-PPTX failed for {os.path.basename(pptx_path)}: {str(e)}")
        
        # Method 1.5: Try zipfile approach
        try:
            import zipfile
            with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
                slide_files = [f for f in zip_ref.namelist() if f.startswith('ppt/slides/slide') and f.endswith('.xml')]
                if len(slide_files) > 0:
                    return True, len(slide_files), "zipfile"
        except Exception as zip_e:
            if status_text:
                status_text.warning(f"Zipfile method failed: {str(zip_e)}")
    
    # Method 2: Try with LibreOffice
    try:
        with tempfile.TemporaryDirectory() as temp_dir:
            # Use libreoffice to convert to PDF
            pdf_output = os.path.join(temp_dir, "presentation.pdf")
            
            # Try with libreoffice command
            lo_command = [
                "libreoffice", "--headless", "--convert-to", "pdf",
                "--outdir", temp_dir, pptx_path
            ]
            
            try:
                # Try libreoffice command
                subprocess.run(lo_command, check=True, capture_output=True, timeout=60)
            except (FileNotFoundError, subprocess.CalledProcessError):
                # If libreoffice fails, try soffice
                lo_command[0] = "soffice"
                subprocess.run(lo_command, check=True, capture_output=True, timeout=60)
            
            if not os.path.exists(pdf_output):
                if status_text:
                    status_text.warning("LibreOffice failed to convert PPTX to PDF")
                return False, 0, ""
            
            # Method 2.1: Use pdftoppm to count pages
            try:
                output_prefix = os.path.join(temp_dir, "page-")
                pdftoppm_command = ["pdftoppm", "-png", pdf_output, output_prefix]
                subprocess.run(pdftoppm_command, check=True, capture_output=True, timeout=30)
                
                # Count the generated files
                page_count = len([f for f in os.listdir(temp_dir) if f.startswith("page-") and f.endswith(".png")])
                if page_count > 0:
                    return True, page_count, SlideCounting.LIBREOFFICE
            except Exception as e:
                if status_text:
                    status_text.warning(f"pdftoppm failed: {str(e)}")
            
            # Method 2.2: Use pdfinfo to count pages
            try:
                pdfinfo_command = ["pdfinfo", pdf_output]
                result = subprocess.run(pdfinfo_command, check=True, capture_output=True, text=True, timeout=10)
                
                # Extract page count from output
                for line in result.stdout.split("\n"):
                    if "Pages:" in line:
                        page_count = int(line.split("Pages:")[1].strip())
                        return True, page_count, SlideCounting.LIBREOFFICE
            except Exception as e:
                if status_text:
                    status_text.warning(f"pdfinfo failed: {str(e)}")
    except Exception as e:
        if status_text:
            status_text.warning(f"LibreOffice conversion failed: {str(e)}")
    
    # All methods failed
    return False, 0, ""

def count_slides_in_artifact_c(directory: str) -> int:
    """
    Count slide descriptions in artifact C files
    
    Args:
        directory: Directory to search in
        
    Returns:
        Number of slide descriptions
    """
    count = 0
    
    # Define patterns for artifact C files
    artifact_c_patterns = [
        r'.*-c\.(txt|md)$',  # Files ending with -c.txt or -c.md
        r'.*-artifact-c\.(txt|md)$',  # Files with -artifact-c.txt or -artifact-c.md
        r'.*slide.*\.(txt|md)$'  # Files with 'slide' in the name
    ]
    
    # Find all potential artifact C files
    artifact_c_files = []
    for root, _, files in os.walk(directory):
        for file in files:
            file_path = os.path.join(root, file)
            if any(re.match(pattern, file, re.IGNORECASE) for pattern in artifact_c_patterns):
                artifact_c_files.append(file_path)
    
    # Count slides in each file
    for file_path in artifact_c_files:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            count += count_slides_in_slide_file(content)
        except Exception:
            pass
    
    return count

def group_files_by_course_lecture(
    english_text_files: List[str],
    english_summary_files: List[str],
    presentations: List[str],
    status_text: Optional[Any] = None
) -> Dict:
    """
    Group files by course and lecture for comprehensive comparison
    
    Args:
        english_text_files: List of English text files
        english_summary_files: List of English summary text files
        presentations: List of presentation files
        status_text: Streamlit text element for status updates
        
    Returns:
        Dictionary with files grouped by course and lecture
    """
    grouped_data = {}
    
    # First pass - organize presentations
    for pptx_path in presentations:
        try:
            # Extract course and lecture info
            course_name = get_full_course_name(pptx_path)
            lecture_num = extract_lecture_number(pptx_path)
            
            # Handle course numbering for sorting
            course_num_match = re.search(r'\d+', course_name)
            course_num = course_num_match.group() if course_num_match else "999"
            
            # Create unique key for this lecture
            key = f"{course_num}:{lecture_num}"
            
            if key not in grouped_data:
                grouped_data[key] = {
                    "course_name": course_name,
                    "course_num": course_num,
                    "lecture_num": lecture_num,
                    "presentation_path": pptx_path,
                    "english_text_dir": None,
                    "english_summary_dir": None,
                    "english_text_files": [],
                    "english_summary_files": []
                }
            else:
                grouped_data[key]["presentation_path"] = pptx_path
        except Exception as e:
            if status_text:
                status_text.warning(f"Error processing {os.path.basename(pptx_path)}: {str(e)}")
    
    # Second pass - find matching English text files
    for file_path in english_text_files:
        try:
            # Extract course and lecture
            parts = file_path.split(os.path.sep)
            
            # Look for "Lecture X" directory which should be the parent of "English text"
            lecture_dir_idx = -1
            for i, part in enumerate(parts):
                if "lecture" in part.lower():
                    lecture_dir_idx = i
                    break
            
            if lecture_dir_idx == -1:
                continue
                
            lecture_dir = parts[lecture_dir_idx]
            lecture_num = extract_lecture_number(lecture_dir)
            
            # Course should be the parent of the lecture directory
            if lecture_dir_idx > 0:
                course_name = parts[lecture_dir_idx - 1]
            else:
                course_name = "Uncategorized"
            
            # Handle course numbering for sorting
            course_num_match = re.search(r'\d+', course_name)
            course_num = course_num_match.group() if course_num_match else "999"
            
            # Create key
            key = f"{course_num}:{lecture_num}"
            
            # Store the directory containing all English text files
            english_text_dir = os.path.dirname(file_path)
            
            # Either add to existing entry or create new one
            if key in grouped_data:
                grouped_data[key]["english_text_dir"] = english_text_dir
                grouped_data[key]["english_text_files"].append(file_path)
            else:
                grouped_data[key] = {
                    "course_name": course_name,
                    "course_num": course_num,
                    "lecture_num": lecture_num,
                    "presentation_path": None,
                    "english_text_dir": english_text_dir,
                    "english_summary_dir": None,
                    "english_text_files": [file_path],
                    "english_summary_files": []
                }
        except Exception as e:
            if status_text:
                status_text.warning(f"Error processing {os.path.basename(file_path)}: {str(e)}")
    
    # Third pass - find matching English summary text files
    for file_path in english_summary_files:
        try:
            # Extract course and lecture
            parts = file_path.split(os.path.sep)
            
            # Look for "Lecture X" directory which should be the parent of "English Summary text"
            lecture_dir_idx = -1
            for i, part in enumerate(parts):
                if "lecture" in part.lower():
                    lecture_dir_idx = i
                    break
            
            if lecture_dir_idx == -1:
                continue
                
            lecture_dir = parts[lecture_dir_idx]
            lecture_num = extract_lecture_number(lecture_dir)
            
            # Course should be the parent of the lecture directory
            if lecture_dir_idx > 0:
                course_name = parts[lecture_dir_idx - 1]
            else:
                course_name = "Uncategorized"
            
            # Handle course numbering for sorting
            course_num_match = re.search(r'\d+', course_name)
            course_num = course_num_match.group() if course_num_match else "999"
            
            # Create key
            key = f"{course_num}:{lecture_num}"
            
            # Store the directory containing all English summary text files
            english_summary_dir = os.path.dirname(file_path)
            
            # Either add to existing entry or create new one
            if key in grouped_data:
                grouped_data[key]["english_summary_dir"] = english_summary_dir
                grouped_data[key]["english_summary_files"].append(file_path)
            else:
                grouped_data[key] = {
                    "course_name": course_name,
                    "course_num": course_num,
                    "lecture_num": lecture_num,
                    "presentation_path": None,
                    "english_text_dir": None,
                    "english_summary_dir": english_summary_dir,
                    "english_text_files": [],
                    "english_summary_files": [file_path]
                }
        except Exception as e:
            if status_text:
                status_text.warning(f"Error processing {os.path.basename(file_path)}: {str(e)}")
    
    return grouped_data

def find_presentation_files_with_all_pptx(input_dir: str) -> List[str]:
    """
    Find all presentation files in a directory, including in all_pptx folders
    
    Args:
        input_dir: Directory to search in
        
    Returns:
        List of presentation file paths
    """
    # Get all presentations from regular paths
    presentations = find_presentation_files(input_dir)
    
    # Find all all_pptx directories
    all_pptx_dirs = []
    for root, dirs, _ in os.walk(input_dir):
        if "all_pptx" in dirs:
            all_pptx_dirs.append(os.path.join(root, "all_pptx"))
    
    # Find presentations in all_pptx directories
    for all_pptx_dir in all_pptx_dirs:
        # Use find_files to maintain consistency with how other files are found
        presentations.extend(find_files(all_pptx_dir, "*.pptx"))
    
    # Remove duplicates while preserving order
    seen = set()
    unique_presentations = []
    for p in presentations:
        if p not in seen:
            seen.add(p)
            unique_presentations.append(p)
    
    return unique_presentations

def check_all_files(status_text: Optional[Any] = None) -> Dict:
    """
    Comprehensive check of all file sources
    
    Args:
        status_text: Streamlit text element for status updates
        
    Returns:
        Dictionary with results
    """
    input_dir = get_input_directory()
    
    if status_text:
        status_text.info("Finding files...")
    
    # Find all relevant files
    english_text_files = find_english_text_files(input_dir)
    english_summary_files = find_english_summary_text_files(input_dir)
    presentations = find_presentation_files_with_all_pptx(input_dir)
    
    if status_text:
        status_text.info(f"Found {len(english_text_files)} English text files, {len(english_summary_files)} English summary text files, and {len(presentations)} presentations")
    
    # Group files by course and lecture
    grouped_data = group_files_by_course_lecture(
        english_text_files,
        english_summary_files,
        presentations,
        status_text
    )
    
    # Process all entries to count slides
    results = []
    
    for key, data in grouped_data.items():
        # Initialize counts
        presentation_count = 0
        english_text_count = 0
        english_summary_count = 0
        artifact_c_count = 0
        
        # Extract metadata
        course_name = data["course_name"]
        lecture_num = data["lecture_num"]
        
        # Get file paths
        presentation_path = data["presentation_path"]
        english_text_dir = data["english_text_dir"]
        english_summary_dir = data["english_summary_dir"]
        
        # Count English text files
        if english_text_dir:
            english_text_count = len(data["english_text_files"])
        
        # Count English summary text files
        if english_summary_dir:
            english_summary_count = len(data["english_summary_files"])
        
        # Count slides in presentation - try to find PPTX if not already found
        pptx_method = ""
        if not presentation_path and english_text_dir:
            # Try to find PPTX file using enhanced search including all_pptx folder
            lecture_dir = os.path.dirname(english_text_dir)
            course_dir = os.path.dirname(lecture_dir)
            presentation_path = find_pptx_for_lecture(course_dir, lecture_num)
        
        if presentation_path:
            success, count, method = count_slides_in_pptx_advanced(presentation_path, status_text)
            if success:
                presentation_count = count
                pptx_method = method
        
        # Artifact C count - look for slide description files in the lecture parent directory
        if english_text_dir:
            # Going up one level from "English text" should get us to the lecture directory
            lecture_dir = os.path.dirname(english_text_dir)
            artifact_c_count = count_slides_in_artifact_c(lecture_dir)
        
        # Check for discrepancies - improved logic to better detect issues
        has_discrepancy = False
        
        # Group counts by their types for better comparison
        non_zero_counts = {
            "PPTX": presentation_count if presentation_count > 0 else None,
            "Text": english_text_count if english_text_count > 0 else None,
            "Summary": english_summary_count if english_summary_count > 0 else None,
            "Artifact": artifact_c_count if artifact_c_count > 0 else None
        }
        
        # Remove None values
        actual_counts = {k: v for k, v in non_zero_counts.items() if v is not None}
        
        # If we have at least 2 different file sources with counts
        if len(actual_counts) >= 2:
            # Get unique count values
            unique_counts = set(actual_counts.values())
            
            # If not all counts are the same, we have a discrepancy
            if len(unique_counts) > 1:
                has_discrepancy = True
            
        # Special case: If we have PPTX = 0 but text/summary files > 0, that's also a discrepancy
        if (presentation_count == 0 and (english_text_count > 0 or english_summary_count > 0)):
            has_discrepancy = True
        
        # Special case: If we have text files but no summary files (or vice versa)
        if (english_text_count > 0 and english_summary_count == 0) or \
           (english_text_count == 0 and english_summary_count > 0):
            has_discrepancy = True
        
        # Store result
        result = {
            "Course Name": course_name,
            "Lecture": f"Lecture {lecture_num}",
            "PPTX Slides": presentation_count,
            "English Text Files": english_text_count,
            "English Summary Files": english_summary_count,
            "Artifact C Slides": artifact_c_count,
            "Has Discrepancy": has_discrepancy,
            "PPTX Path": presentation_path if presentation_path else "",
            "English Text Path": english_text_dir if english_text_dir else "",
            "English Summary Path": english_summary_dir if english_summary_dir else "",
            "PPTX Method": pptx_method
        }
        
        results.append(result)
    
    # Calculate statistics
    total_lectures = len(results)
    total_discrepancies = sum(1 for r in results if r["Has Discrepancy"])
    
    # Count by type of mismatch
    pptx_text_mismatch = 0
    pptx_summary_mismatch = 0
    text_summary_mismatch = 0
    artifact_mismatch = 0
    missing_pptx = 0
    missing_text = 0
    missing_summary = 0
    
    for r in results:
        pptx = r["PPTX Slides"]
        text = r["English Text Files"]
        summary = r["English Summary Files"]
        artifact = r["Artifact C Slides"]
        
        # Count mismatches between non-zero values
        if pptx > 0 and text > 0 and pptx != text:
            pptx_text_mismatch += 1
        
        if pptx > 0 and summary > 0 and pptx != summary:
            pptx_summary_mismatch += 1
        
        if text > 0 and summary > 0 and text != summary:
            text_summary_mismatch += 1
        
        if artifact > 0 and ((pptx > 0 and artifact != pptx) or 
                             (text > 0 and artifact != text) or
                             (summary > 0 and artifact != summary)):
            artifact_mismatch += 1
            
        # Count special cases where one source is missing (has 0 count)
        if pptx == 0 and (text > 0 or summary > 0):
            missing_pptx += 1
            
        if text == 0 and (pptx > 0 or summary > 0):
            missing_text += 1
            
        if summary == 0 and (pptx > 0 or text > 0):
            missing_summary += 1
    
    return {
        "results": results,
        "stats": {
            "total_lectures": total_lectures,
            "total_discrepancies": total_discrepancies,
            "pptx_text_mismatch": pptx_text_mismatch,
            "pptx_summary_mismatch": pptx_summary_mismatch,
            "text_summary_mismatch": text_summary_mismatch,
            "artifact_mismatch": artifact_mismatch,
            "missing_pptx": missing_pptx,
            "missing_text": missing_text,
            "missing_summary": missing_summary
        }
    }

def main():
    st.title("Count - artifacts/pptx")
    st.write("Enhanced verification system for slide counts across multiple file sources")
    
    # Instructions
    with st.expander("Instructions & Information", expanded=False):
        st.markdown("""
        This tool performs comprehensive verification by counting:
        1. **PPTX Slides**: Number of slides in presentation files
        2. **English Text Files**: Number of text files in "English text" folders
        3. **English Summary Files**: Number of text files in "English Summary text" folders
        4. **Artifact C Slides**: Number of slide descriptions in slide description files
        
        All four counts should match for each lecture.
        """)
    
    input_dir = get_input_directory()
    
    if not os.path.exists(input_dir):
        st.error(f"Input directory not found: {input_dir}")
        st.info("Please create an 'input' directory in the project root and add your files.")
        return
    
    # Count Button
    if st.button("Count All Files"):
        with st.spinner("Counting files in all sources..."):
            # Display progress
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            # Initialize progress steps
            status_text.text("Finding files...")
            progress_bar.progress(0.2)
            
            # Count files
            results = check_all_files(status_text)
            
            # Update progress
            status_text.text("Processing results...")
            progress_bar.progress(0.9)
            time.sleep(0.5)
            
            # Complete
            progress_bar.progress(1.0)
            status_text.text("Counting complete")
            
            # Store results in session state
            st.session_state.count_new_results = results
    
    # Display results if available
    if hasattr(st.session_state, 'count_new_results'):
        results = st.session_state.count_new_results
        data = results["results"]
        stats = results["stats"]
        
        # Summary dashboard
        st.header("Status Dashboard")
        
        col1, col2, col3, col4 = st.columns(4)
        with col1:
            st.metric("Total Lectures", stats["total_lectures"])
        with col2:
            st.metric("Total Discrepancies", stats["total_discrepancies"])
        with col3:
            st.metric("Discrepancy Rate", f"{stats['total_discrepancies']/max(1, stats['total_lectures'])*100:.1f}%")
        with col4:
            st.metric("Missing Sources", sum(1 for r in data if any(c == 0 for c in [r["PPTX Slides"], r["English Text Files"], r["English Summary Files"]])))
        
        # Breakdown of mismatches
        st.subheader("Discrepancy Breakdown")
        
        # Row 1: Count mismatches between different sources
        st.write("**Count Mismatches:**")
        col1, col2, col3, col4 = st.columns(4)
        with col1:
            st.metric("PPTX vs Text", stats["pptx_text_mismatch"])
        with col2:
            st.metric("PPTX vs Summary", stats["pptx_summary_mismatch"])
        with col3:
            st.metric("Text vs Summary", stats["text_summary_mismatch"])
        with col4:
            st.metric("Artifact C Issues", stats["artifact_mismatch"])
        
        # Row 2: Missing sources
        st.write("**Missing Sources:**")
        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric("Missing PPTX", stats["missing_pptx"], 
                      help="Lectures with English text or summary files but no PPTX file (or zero slides)")
        with col2:
            st.metric("Missing Text", stats["missing_text"],
                     help="Lectures with PPTX or summary files but no English text files")
        with col3:
            st.metric("Missing Summary", stats["missing_summary"],
                     help="Lectures with PPTX or text files but no English summary files")
        
        # Full results table
        st.header("All Lectures")
        
        # Create DataFrame for display
        df_display = pd.DataFrame(data)

        if df_display.empty:
            st.info("No lecture rows were discovered in the current input directory.")
            return
        
        # Sort the DataFrame by Course Name and Lecture
        df_display = df_display.sort_values(by=["Course Name", "Lecture"])
        
        # Create a styled version for display with conditional formatting
        def highlight_discrepancies(val):
            if val is True:
                return 'background-color: #ffcccc'  # Light red
            return ''
        
        # Add styling based on the Has Discrepancy column
        st.dataframe(
            df_display.style.apply(
                lambda row: ['background-color: #ffcccc' if row['Has Discrepancy'] else '' for _ in row], 
                axis=1
            ),
            use_container_width=True
        )
        
        # Display discrepancies only if there are any
        if stats["total_discrepancies"] > 0:
            st.header("Discrepancies Only")
            
            # Filter for discrepancies
            discrepancies = [r for r in data if r["Has Discrepancy"]]
            df_discrepancies = pd.DataFrame(discrepancies)
            
            # Sort the DataFrame by Course Name and Lecture
            df_discrepancies = df_discrepancies.sort_values(by=["Course Name", "Lecture"])
            
            st.dataframe(df_discrepancies, use_container_width=True)
        
        # Detailed view option
        st.header("Detailed File Paths")
        with st.expander("Show File Paths", expanded=False):
            # Create a DataFrame with just the file paths and counts
            paths_df = pd.DataFrame([{
                "Course Name": r["Course Name"],
                "Lecture": r["Lecture"],
                "PPTX Path": r["PPTX Path"],
                "PPTX Slides": r["PPTX Slides"],
                "PPTX Method": r["PPTX Method"],
                "English Text Path": r["English Text Path"],
                "English Text Files": r["English Text Files"],
                "English Summary Path": r["English Summary Path"],
                "English Summary Files": r["English Summary Files"]
            } for r in data])
            
            # Sort by course and lecture
            paths_df = paths_df.sort_values(by=["Course Name", "Lecture"])
            
            st.dataframe(paths_df, use_container_width=True)

if __name__ == "__main__":
    main()
