"""
CLI Command: 4K Image Processing (Page 06)

Converts the Streamlit page 06_4K_Image_pptx_zip.py to CLI interface
while maintaining 100% functional parity.

This command processes PPTX presentations and ZIP archives containing images,
extracting slides/images and upscaling them to 4K resolution with logo and copyright overlay.
"""

import click
import os
import sys
import re
import time
import cv2
import numpy as np
import io
import enum
import fnmatch
import platform
import tempfile
import glob
import shutil
import subprocess
import zipfile
from pathlib import Path
from typing import Dict, Any, List, Tuple, Optional
from datetime import datetime, timedelta
from PIL import Image, ImageDraw, ImageFont

# Add project root to path for imports
project_root = Path(__file__).parent.parent.parent
sys.path.insert(0, str(project_root))

from cli.progress import DreamletProgress
from cli.reports import generate_report
from cli.config import load_config

# Try to import LANCZOS for different PIL versions
try:
    LANCZOS = Image.Resampling.LANCZOS
except AttributeError:
    LANCZOS = Image.LANCZOS

# Import PPTX library
from pptx import Presentation

from rich.console import Console
console = Console()

class StatusManager:
    """CLI-compatible status manager for 4K image processing"""
    
    def __init__(self, verbose=False):
        self.verbose = verbose
        self.console = Console()
        
    def info(self, message, verbose_only=False):
        """Display info message"""
        if not verbose_only or self.verbose:
            self.console.print(f"[cyan]ℹ[/cyan] {message}")
    
    def warning(self, message):
        """Display warning message"""
        self.console.print(f"[yellow]⚠[/yellow] {message}")
    
    def error(self, message):
        """Display error message"""
        self.console.print(f"[red]✗[/red] {message}")
    
    def success(self, message):
        """Display success message"""
        self.console.print(f"[green]✓[/green] {message}")

# Constants for image processing
LOGO_PATH = 'config/logo.png'
COPYRIGHT_PATH = 'config/copyright.txt'
LOGO_SIZE = 250
LOGO_PADDING = 0
LOGO_POSITION = 'top-right'
COPYRIGHT_PADDING = 35
COPYRIGHT_POSITION = 'bottom-center'
FONT_SIZE = 65
FONT_COLOR = (0, 0, 0)  # Black
TARGET_RESOLUTION = (3840, 2160)  # 4K

# Conversion methods enum
class ConversionMethod(str, enum.Enum):
    LIBREOFFICE = "libreoffice"
    PYTHON_PPTX = "python-pptx"
    PDF2IMAGE = "pdf2image"

CONVERSION_METHODS = {
    ConversionMethod.LIBREOFFICE: "LibreOffice + pdftoppm (best quality)",
    ConversionMethod.PYTHON_PPTX: "Direct PPTX processing (fastest)",
    ConversionMethod.PDF2IMAGE: "PDF2Image (most compatible)"
}

def increase_image_decompression_limit():
    """Increases the Pillow image decompression bomb limit to avoid errors with very large images"""
    new_limit = 3840 * 2160 * 50  # Approximately 415 million pixels
    Image.MAX_IMAGE_PIXELS = new_limit

def get_input_directory() -> str:
    """Get the path to the input directory"""
    input_dir = os.path.join(os.getcwd(), "input")
    return input_dir

def find_files(directory: str, pattern: str) -> List[str]:
    """Find all files matching a pattern in a directory (recursively)"""
    result = []
    for root, _, filenames in os.walk(directory):
        for filename in fnmatch.filter(filenames, pattern):
            result.append(os.path.join(root, filename))
    return result

def find_presentation_files(directory: str) -> List[str]:
    """Find all presentation files (PPTX and ZIP) in a directory"""
    pptx_files = find_files(directory, "*.pptx")
    zip_files = find_files(directory, "*.zip")
    return pptx_files + zip_files

def extract_numeric_prefix(filename: str) -> int:
    """Extract numeric prefix from filename like '1_Title.png' -> 1"""
    match = re.match(r'^(\d+)', filename)
    if match:
        return int(match.group(1))
    return 999999  # Put files without numeric prefix at the end

def read_copyright(file_path):
    """Read copyright text from file"""
    try:
        # Try UTF-8 encoding first
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read().strip()
            
        # Clean up any encoding issues with the copyright symbol
        # Remove any 'A' characters that appear immediately before '©'
        text = re.sub(r'A©', '©', text)
        
        return text
    except UnicodeDecodeError:
        # If UTF-8 fails, try with different encoding
        try:
            with open(file_path, 'r', encoding='cp1252') as f:
                text = f.read().strip()
                # Clean up any encoding issues
                text = re.sub(r'A©', '©', text)
                return text
        except Exception:
            return "© All Rights Reserved"
    except Exception:
        return "© All Rights Reserved"

def get_position(img_width, img_height, element_width, element_height, position, padding):
    """Calculate position based on alignment and padding"""
    if position == 'top-left':
        return (padding, padding)
    elif position == 'top-right':
        return (img_width - element_width - padding, padding)
    elif position == 'bottom-left':
        return (padding, img_height - element_height - padding)
    elif position == 'bottom-right':
        return (img_width - element_width - padding, img_height - element_height - padding)
    elif position == 'bottom-center':
        return ((img_width - element_width) // 2, img_height - element_height - padding)
    else:
        raise ValueError(f"Invalid position: {position}")

def get_font(size):
    """Get a font with fallback options"""
    try:
        # Try to use a common sans-serif font
        return ImageFont.truetype("Arial.ttf", size)
    except IOError:
        try:
            # Fallback to DejaVuSans if Arial is not available
            return ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", size)
        except IOError:
            # If both fail, use the default font
            return ImageFont.load_default()

def process_image(input_path, output_path, logo, copyright_text, status_manager: StatusManager):
    """Process an image by adding logo and copyright text"""
    try:
        with Image.open(input_path) as img:
            # Check if the image is already 4K or larger
            if img.width >= TARGET_RESOLUTION[0] and img.height >= TARGET_RESOLUTION[1]:
                processed_img = img.copy()
            else:
                # Upscale the image to 4K while maintaining aspect ratio
                aspect_ratio = img.width / img.height
                if aspect_ratio > 16/9:
                    new_width = TARGET_RESOLUTION[0]
                    new_height = int(TARGET_RESOLUTION[0] / aspect_ratio)
                else:
                    new_height = TARGET_RESOLUTION[1]
                    new_width = int(TARGET_RESOLUTION[1] * aspect_ratio)
                processed_img = img.resize((new_width, new_height), LANCZOS)
            
            # Convert to RGBA if necessary
            if processed_img.mode != 'RGBA':
                processed_img = processed_img.convert('RGBA')

            # Add logo if it exists
            if logo:
                logo_resized = logo.resize((LOGO_SIZE, LOGO_SIZE), LANCZOS)
                logo_pos = get_position(processed_img.width, processed_img.height, 
                                       LOGO_SIZE, LOGO_SIZE, LOGO_POSITION, LOGO_PADDING)
                
                # Create a new image for the logo with an alpha channel
                logo_img = Image.new('RGBA', processed_img.size, (0, 0, 0, 0))
                logo_img.paste(logo_resized, logo_pos, logo_resized)
                
                # Composite the logo onto the processed image
                processed_img = Image.alpha_composite(processed_img, logo_img)

            # Add copyright
            draw = ImageDraw.Draw(processed_img)
            font = get_font(FONT_SIZE)
            
            # Get text size
            text_bbox = draw.textbbox((0, 0), copyright_text, font=font)
            text_width = text_bbox[2] - text_bbox[0]
            text_height = text_bbox[3] - text_bbox[1]
            
            text_pos = get_position(processed_img.width, processed_img.height, 
                                   text_width, text_height, COPYRIGHT_POSITION, COPYRIGHT_PADDING)
            draw.text(text_pos, copyright_text, font=font, fill=FONT_COLOR)

            # Save the processed image with the same format and quality as the original
            if img.format == 'JPEG':
                processed_img = processed_img.convert('RGB')
                processed_img.save(output_path, format=img.format, quality=95, optimize=True, progressive=True)
            elif img.format == 'PNG':
                processed_img.save(output_path, format=img.format, optimize=True)
            else:
                processed_img.save(output_path, format=img.format)
        
        return True
    except Exception as e:
        status_manager.error(f"Error processing {input_path}: {str(e)}")
        return False

def extract_slides_from_zip(zip_file: str, output_folder: str, status_mgr: StatusManager, target_resolution=(3840, 2160)) -> Tuple[bool, List[str], int]:
    """
    Extract images from a ZIP file, upscale them to 4K, and save them WITHOUT logo/copyright.
    The main processing loop will add logo and copyright later.
    """
    temp_dir = tempfile.mkdtemp()
    
    try:
        os.makedirs(output_folder, exist_ok=True)
        
        status_mgr.info(f"Extracting ZIP file: {os.path.basename(zip_file)}")
        
        # Extract ZIP contents
        try:
            with zipfile.ZipFile(zip_file, 'r') as zip_ref:
                zip_ref.extractall(temp_dir)
        except Exception as e:
            error_msg = f"Error extracting ZIP file: {str(e)}"
            status_mgr.error(error_msg)
            return False, [], 0
        
        # Find all image files in the extracted contents
        image_extensions = ['*.png', '*.jpg', '*.jpeg', '*.PNG', '*.JPG', '*.JPEG']
        image_files = []
        for ext in image_extensions:
            image_files.extend(glob.glob(os.path.join(temp_dir, '**', ext), recursive=True))
        
        if not image_files:
            error_msg = "No image files found in ZIP archive"
            status_mgr.error(error_msg)
            return False, [], 0
        
        # Sort images by numeric prefix to maintain slide order
        image_files.sort(key=lambda x: extract_numeric_prefix(os.path.basename(x)))
        
        status_mgr.info(f"Found {len(image_files)} images in ZIP, processing...")
        
        # Process each image
        extracted_images = []
        for i, img_path in enumerate(image_files, 1):
            try:
                with Image.open(img_path) as img:
                    # Check if upscaling is needed
                    aspect_ratio = img.width / img.height
                    
                    if img.width >= target_resolution[0] and img.height >= target_resolution[1]:
                        # Image is already 4K or larger, just copy it
                        processed_img = img.copy()
                    else:
                        # Upscale to 4K while maintaining aspect ratio
                        if aspect_ratio > target_resolution[0] / target_resolution[1]:
                            new_width = target_resolution[0]
                            new_height = int(target_resolution[0] / aspect_ratio)
                        else:
                            new_height = target_resolution[1]
                            new_width = int(target_resolution[1] * aspect_ratio)
                        
                        processed_img = img.resize((new_width, new_height), LANCZOS)
                    
                    # Save with sequential naming
                    output_path = os.path.join(output_folder, f"{i:02d}.png")
                    processed_img.save(output_path, "PNG", optimize=True)
                    extracted_images.append(output_path)
                    
            except Exception as e:
                status_mgr.warning(f"Error processing image {os.path.basename(img_path)}: {str(e)}")
                continue
        
        image_count = len(extracted_images)
        
        if image_count == 0:
            error_msg = "Failed to process any images from ZIP"
            status_mgr.error(error_msg)
            return False, [], 0
        
        status_mgr.info(f"Successfully extracted {image_count} images from ZIP")
        
        return True, extracted_images, image_count
        
    finally:
        # Clean up temporary directory
        shutil.rmtree(temp_dir, ignore_errors=True)

def get_libreoffice_path():
    """Get the correct LibreOffice path based on the operating system"""
    if platform.system() == 'Darwin':  # macOS
        possible_paths = [
            '/Applications/LibreOffice.app/Contents/MacOS/soffice',
            '/Applications/LibreOffice.app/Contents/MacOS/soffice.bin',
            '/opt/homebrew/bin/soffice',
            '/usr/local/bin/soffice'
        ]
        for path in possible_paths:
            if os.path.exists(path):
                return path
        return None
    else:
        return 'soffice'

def extract_slides_directly(pptx_path, output_dir, status_mgr: StatusManager, target_resolution=(3840, 2160)):
    """Extract slides directly from PPTX file using python-pptx"""
    try:
        os.makedirs(output_dir, exist_ok=True)
        
        status_mgr.info(f"Extracting slides directly from PPTX: {os.path.basename(pptx_path)}")
        
        prs = Presentation(pptx_path)
        slide_count = len(prs.slides)
        
        if slide_count == 0:
            status_mgr.warning("Presentation contains no slides.")
            return False, "Presentation contains no slides", 0
        
        for i, slide in enumerate(prs.slides, 1):
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            slide_aspect_ratio = slide_width / slide_height
            
            if slide_aspect_ratio > target_resolution[0] / target_resolution[1]:
                width = target_resolution[0]
                height = int(width / slide_aspect_ratio)
            else:
                height = target_resolution[1]
                width = int(height * slide_aspect_ratio)
            
            img = np.ones((height, width, 3), np.uint8) * 255
            
            scale_x = width / slide_width
            scale_y = height / slide_height
            
            for shape in slide.shapes:
                if hasattr(shape, 'image') and shape.image:
                    try:
                        image_stream = io.BytesIO(shape.image.blob)
                        slide_image = Image.open(image_stream)
                        
                        slide_image_cv = cv2.cvtColor(np.array(slide_image), cv2.COLOR_RGB2BGR)
                        
                        left = int(shape.left * scale_x)
                        top = int(shape.top * scale_y)
                        img_width = int(shape.width * scale_x)
                        img_height = int(shape.height * scale_y)
                        
                        resized_image = cv2.resize(slide_image_cv, (img_width, img_height), interpolation=cv2.INTER_LANCZOS4)
                        
                        try:
                            if img_height > 0 and img_width > 0:
                                if top + img_height <= height and left + img_width <= width:
                                    img[top:top+img_height, left:left+img_width] = resized_image
                        except Exception as e:
                            status_mgr.warning(f"Error placing image on slide {i}: {str(e)}")
                    except Exception as e:
                        status_mgr.warning(f"Error processing image in slide {i}: {str(e)}")
            
            output_path = os.path.join(output_dir, f"{i:02d}.png")
            cv2.imwrite(output_path, img, [cv2.IMWRITE_PNG_COMPRESSION, 9])
        
        return True, f"Successfully extracted {slide_count} slides", slide_count
    
    except Exception as e:
        error_msg = f"Error in direct PPTX extraction: {str(e)}"
        status_mgr.error(error_msg)
        return False, error_msg, 0

def convert_with_pdf2image(pptx_path, output_dir, status_mgr: StatusManager, target_resolution=(3840, 2160)):
    """Convert PPTX to images using LibreOffice for PPTX→PDF and pdf2image for PDF→PNG"""
    temp_dir = tempfile.mkdtemp()
    pdf_path = os.path.join(temp_dir, os.path.splitext(os.path.basename(pptx_path))[0] + '.pdf')
    
    try:
        os.makedirs(output_dir, exist_ok=True)
        
        status_mgr.info(f"Converting presentation using pdf2image: {os.path.basename(pptx_path)}")
        
        libreoffice_path = get_libreoffice_path()
        if not libreoffice_path:
            error_msg = "LibreOffice not found. Please make sure LibreOffice is installed."
            status_mgr.error(error_msg)
            return False, error_msg, 0
        
        try:
            status_mgr.info("Converting presentation to PDF...")
            
            process = subprocess.run(
                [libreoffice_path, '--headless', '--convert-to', 'pdf', 
                 '--outdir', temp_dir, pptx_path], 
                check=True, capture_output=True, timeout=120
            )
            
            if not os.path.exists(pdf_path):
                error_msg = f"PDF was not created at expected path: {pdf_path}"
                status_mgr.error(error_msg)
                return False, error_msg, 0
                
        except Exception as e:
            error_msg = f"Error converting to PDF: {str(e)}"
            status_mgr.error(error_msg)
            return False, error_msg, 0
        
        try:
            status_mgr.info("Converting PDF to images using pdf2image...")
            
            from pdf2image import convert_from_path
            
            images = convert_from_path(pdf_path, dpi=300, fmt="png")
            slide_count = len(images)
            
            if slide_count == 0:
                error_msg = "No slides were extracted from the PDF"
                status_mgr.error(error_msg)
                return False, error_msg, 0
            
            for i, img in enumerate(images, 1):
                aspect_ratio = img.width / img.height
                
                if aspect_ratio > target_resolution[0] / target_resolution[1]:
                    new_width = target_resolution[0]
                    new_height = int(target_resolution[0] / aspect_ratio)
                else:
                    new_height = target_resolution[1]
                    new_width = int(target_resolution[1] * aspect_ratio)
                
                resized_img = img.resize((new_width, new_height), LANCZOS)
                
                output_path = os.path.join(output_dir, f"{i:02d}.png")
                resized_img.save(output_path, "PNG", optimize=True)
            
            return True, f"Successfully extracted {slide_count} slides", slide_count
        
        except ImportError:
            error_msg = "pdf2image module not available. Please install it with: pip install pdf2image"
            status_mgr.error(error_msg)
            return False, error_msg, 0
        except Exception as e:
            error_msg = f"Error converting PDF to images: {str(e)}"
            status_mgr.error(error_msg)
            return False, error_msg, 0
            
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)

def extract_slides_from_pptx(pptx_path, output_dir, status_mgr: StatusManager):
    """Extract slides from PPTX using LibreOffice + pdftoppm method"""
    temp_dir = tempfile.mkdtemp()
    
    try:
        os.makedirs(output_dir, exist_ok=True)
        
        status_mgr.info(f"Converting presentation to PDF using LibreOffice...")
        
        libreoffice_path = get_libreoffice_path()
        if not libreoffice_path:
            error_msg = "LibreOffice not found. Please make sure LibreOffice is installed."
            status_mgr.error(error_msg)
            return False, error_msg, 0
            
        try:
            process = subprocess.run(
                [libreoffice_path, '--headless', '--convert-to', 'pdf',
                 '--outdir', temp_dir, pptx_path], 
                check=True, capture_output=True, timeout=120
            )
            
            expected_pdf_name = os.path.splitext(os.path.basename(pptx_path))[0] + '.pdf'
            pdf_path = os.path.join(temp_dir, expected_pdf_name)
            
            if not os.path.exists(pdf_path):
                error_msg = f"PDF was not created at expected path: {pdf_path}"
                status_mgr.error(error_msg)
                return False, error_msg, 0
                
        except Exception as e:
            error_msg = f"Error converting to PDF: {str(e)}"
            status_mgr.error(error_msg)
            return False, error_msg, 0
            
        # Convert PDF to images using pdftoppm
        status_mgr.info(f"Converting PDF to images...")
            
        temp_img_dir = tempfile.mkdtemp()
        success = False
        
        try:
            base_filename = os.path.join(temp_img_dir, "slide")
            pdftoppm_command = [
                "pdftoppm", "-png", "-r", "300",
                pdf_path, base_filename
            ]
            subprocess.run(pdftoppm_command, check=True, capture_output=True, timeout=120)
            
            extracted_files = sorted([f for f in os.listdir(temp_img_dir) if f.startswith("slide-") and f.endswith(".png")])
            slide_count = len(extracted_files)
            
            if slide_count > 0:
                for i, filename in enumerate(extracted_files, 1):
                    file_path = os.path.join(temp_img_dir, filename)
                    with Image.open(file_path) as img:
                        aspect_ratio = img.width / img.height
                        
                        if aspect_ratio > 3840/2160:
                            new_width = 3840
                            new_height = int(3840 / aspect_ratio)
                        else:
                            new_height = 2160
                            new_width = int(2160 * aspect_ratio)
                        
                        resized = img.resize((new_width, new_height), LANCZOS)
                        new_filename = os.path.join(output_dir, f"{i:02d}.png")
                        resized.save(new_filename, "PNG", optimize=True)
                        
                return True, f"Successfully extracted {slide_count} slides", slide_count
        except Exception:
            pass
            
        # Fallback to pdf2image
        try:
            from pdf2image import convert_from_path
            
            images = convert_from_path(pdf_path, dpi=300, output_folder=temp_img_dir, fmt="png")
            slide_count = len(images)
            
            if slide_count > 0:
                for i, img in enumerate(images, 1):
                    aspect_ratio = img.width / img.height
                    
                    if aspect_ratio > 3840/2160:
                        new_width = 3840
                        new_height = int(3840 / aspect_ratio)
                    else:
                        new_height = 2160
                        new_width = int(2160 * aspect_ratio)
                    
                    resized = img.resize((new_width, new_height), LANCZOS)
                    output_path = os.path.join(output_dir, f"{i:02d}.png")
                    resized.save(output_path, "PNG", optimize=True)
                    
                return True, f"Successfully extracted {slide_count} slides", slide_count
        except Exception:
            pass
            
        return False, "Failed to extract slides - required tools not available", 0
        
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)

def extract_slides_with_method(
    pptx_path: str, 
    output_dir: str, 
    method: ConversionMethod,
    status_mgr: StatusManager,
    enable_auto_fallback: bool = True,
    target_resolution: Tuple[int, int] = (3840, 2160)
) -> Tuple[bool, str, int]:
    """Extract slides from a PPTX file using the specified method"""
    os.makedirs(output_dir, exist_ok=True)
    
    status_mgr.info(f"Attempting extraction using {CONVERSION_METHODS[method]} method")
    
    success, message, slide_count = _try_method(pptx_path, output_dir, method, status_mgr, target_resolution)
    
    if not success and enable_auto_fallback:
        status_mgr.warning(f"Primary method failed: {message}")
        status_mgr.info("Trying fallback methods...")
        
        fallback_methods = get_fallback_methods(method)
        
        for fallback_method in fallback_methods:
            status_mgr.info(f"Trying {CONVERSION_METHODS[fallback_method]} method...")
                
            success, message, slide_count = _try_method(pptx_path, output_dir, fallback_method, status_mgr, target_resolution)
            
            if success:
                status_mgr.success(f"Fallback method {CONVERSION_METHODS[fallback_method]} succeeded")
                break
        
    return success, message, slide_count

def get_fallback_methods(primary_method: ConversionMethod) -> List[ConversionMethod]:
    """Get a list of fallback methods to try if the primary method fails"""
    all_methods = list(ConversionMethod)
    all_methods.remove(primary_method)
    
    if primary_method == ConversionMethod.LIBREOFFICE:
        return [ConversionMethod.PDF2IMAGE, ConversionMethod.PYTHON_PPTX]
    elif primary_method == ConversionMethod.PYTHON_PPTX:
        return [ConversionMethod.LIBREOFFICE, ConversionMethod.PDF2IMAGE]
    else:  # PDF2IMAGE
        return [ConversionMethod.LIBREOFFICE, ConversionMethod.PYTHON_PPTX]

def _try_method(
    pptx_path: str, 
    output_dir: str, 
    method: ConversionMethod,
    status_mgr: StatusManager,
    target_resolution: Tuple[int, int]
) -> Tuple[bool, str, int]:
    """Try to extract slides using a specific method"""
    try:
        if method == ConversionMethod.LIBREOFFICE:
            return extract_slides_from_pptx(pptx_path, output_dir, status_mgr)
        elif method == ConversionMethod.PYTHON_PPTX:
            return extract_slides_directly(pptx_path, output_dir, status_mgr, target_resolution)
        elif method == ConversionMethod.PDF2IMAGE:
            return convert_with_pdf2image(pptx_path, output_dir, status_mgr, target_resolution)
        else:
            error_msg = f"Unknown conversion method: {method}"
            status_mgr.error(error_msg)
            return False, error_msg, 0
            
    except ImportError as e:
        error_msg = f"Required module not available for {method} method: {str(e)}"
        status_mgr.error(error_msg)
        return False, error_msg, 0
        
    except Exception as e:
        error_msg = f"Error in {method} method: {str(e)}"
        status_mgr.error(error_msg)
        return False, error_msg, 0

def extract_and_upscale_images(pptx_file, output_folder, conversion_method=None, enable_auto_fallback=True, status_mgr=None):
    """
    Extract images from a PPTX file, upscale them to 4K, and save them in the output folder.
    """
    # Use the multi-method converter
    success, message, slide_count = extract_slides_with_method(
        pptx_file, 
        output_folder, 
        method=conversion_method or ConversionMethod.LIBREOFFICE,
        status_mgr=status_mgr,
        enable_auto_fallback=enable_auto_fallback,
        target_resolution=TARGET_RESOLUTION
    )
    
    if not success:
        status_mgr.error(message)
        return False, [], 0
    
    # Gather the paths of the created images
    extracted_images = []
    if slide_count > 0:
        for i in range(1, slide_count + 1):
            image_path = os.path.join(output_folder, f"{i:02d}.png")
            if os.path.exists(image_path):
                extracted_images.append(image_path)
    
    return success, extracted_images, slide_count

def find_presentations() -> List[str]:
    """
    Find all presentation files in the input directory
    
    Returns:
        List of presentation file paths
    """
    input_dir = get_input_directory()
    
    # Find all presentation files
    all_presentations = find_presentation_files(input_dir)
    
    return all_presentations

def process_presentations(input_dir: str, status_manager: StatusManager, create_without_logo_folder=False, 
                 conversion_method=None, enable_auto_fallback=True) -> List[Dict]:
    """
    Process all PPTX and ZIP files in the input directory structure
    """
    results = []
    
    # Try to load logo and copyright text
    try:
        logo = Image.open(LOGO_PATH).convert("RGBA") if os.path.exists(LOGO_PATH) else None
        copyright_text = read_copyright(COPYRIGHT_PATH) if os.path.exists(COPYRIGHT_PATH) else "© All Rights Reserved"
    except Exception as e:
        status_manager.warning(f"Could not load logo or copyright: {str(e)}")
        logo = None
        copyright_text = "© All Rights Reserved"
    
    # Find all presentations in the input directory
    all_presentations = find_presentations()
    total_presentations = len(all_presentations)
    
    # Process each presentation
    for i, pptx_file in enumerate(all_presentations):
        # Update progress with file info
        processed_count = i + 1
        file_name = os.path.basename(pptx_file)
        status_manager.info(f"Processing {processed_count}/{total_presentations}: {file_name}")
        
        result = {
            "file_path": pptx_file,
            "status": "error",
            "message": "",
            "slide_count": 0,
            "output_dir": ""
        }
        
        try:
            # Determine file type first
            file_extension = os.path.splitext(pptx_file)[1].lower()
            is_zip_file = file_extension == '.zip'
            
            # Get the directory containing the presentation
            pptx_dir = os.path.dirname(pptx_file)
            
            # Check if this is inside an all_pptx directory
            parent_dir = os.path.basename(pptx_dir)
            if parent_dir == 'all_pptx':
                # If file is in all_pptx, create lecture folder in parent directory
                course_dir = os.path.dirname(pptx_dir)  # Go up one level from all_pptx
                pptx_name = os.path.splitext(os.path.basename(pptx_file))[0]
                
                # Extract lecture number from filename for proper naming
                lecture_match = re.search(r'(\d+)', pptx_name)
                if lecture_match:
                    lecture_num = lecture_match.group(1).zfill(2)
                    lecture_folder = os.path.join(course_dir, f"Lecture {lecture_num}")
                else:
                    lecture_folder = os.path.join(course_dir, pptx_name)
                    
                os.makedirs(lecture_folder, exist_ok=True)
                
                # Create "English image" folder inside the lecture folder
                output_folder = os.path.join(lecture_folder, "English image")
                os.makedirs(output_folder, exist_ok=True)
                
                # The all_pptx folder already exists since the file is in it
                all_pptx_folder = pptx_dir
            else:
                # Create 'all_pptx' directory in the same folder as the presentation
                all_pptx_folder = os.path.join(pptx_dir, 'all_pptx')
                os.makedirs(all_pptx_folder, exist_ok=True)
                
                # Create folder with same name as file for images
                file_name = os.path.splitext(os.path.basename(pptx_file))[0]
                
                # For ZIP files with numeric names, add "Lecture" prefix to match PPTX structure
                if is_zip_file and re.match(r'^\d+$', file_name):
                    lecture_folder = os.path.join(pptx_dir, f"Lecture {file_name.zfill(2)}")
                else:
                    lecture_folder = os.path.join(pptx_dir, file_name)
                    
                os.makedirs(lecture_folder, exist_ok=True)
                
                # Create "English image" folder inside the lecture folder
                output_folder = os.path.join(lecture_folder, "English image")
                os.makedirs(output_folder, exist_ok=True)
            
            # Create 'without_logo_png' folder inside the output folder only if requested
            without_logo_folder = os.path.join(output_folder, 'without_logo_png')
            if create_without_logo_folder:
                os.makedirs(without_logo_folder, exist_ok=True)
            
            # Check if file is already processed and in all_pptx
            new_pptx_path = os.path.join(all_pptx_folder, os.path.basename(pptx_file))
            if os.path.exists(new_pptx_path):
                # If already processed, just update the result
                result["message"] = "Presentation already processed"
                
                # Count existing images
                png_count = len([f for f in os.listdir(output_folder) 
                               if f.endswith('.png') and not os.path.isdir(os.path.join(output_folder, f))])
                
                result["status"] = "success"
                result["slide_count"] = png_count
                result["output_dir"] = output_folder
                results.append(result)
                continue
                
            # Extract and upscale images from the presentation or ZIP file
            temp_dir = None
            
            try:
                # Only create and use without_logo_folder if requested
                if create_without_logo_folder:
                    # Ensure the folder exists
                    os.makedirs(without_logo_folder, exist_ok=True)
                    
                    if is_zip_file:
                        # Extract from ZIP file
                        success, extracted_images, slide_count = extract_slides_from_zip(
                            pptx_file,
                            without_logo_folder,
                            status_mgr=status_manager,
                            target_resolution=TARGET_RESOLUTION
                        )
                    else:
                        # Extract from PPTX file
                        success, extracted_images, slide_count = extract_and_upscale_images(
                            pptx_file, 
                            without_logo_folder,
                            conversion_method=conversion_method,
                            enable_auto_fallback=enable_auto_fallback,
                            status_mgr=status_manager
                        )
                    
                    # Process each extracted image (add logo and copyright)
                    processed_img_count = 0
                    for img_path in extracted_images:
                        # Original image is in without_logo_folder
                        # Processed image will be in output_folder
                        output_filename = os.path.basename(img_path)
                        output_path = os.path.join(output_folder, output_filename)
                        
                        # Process the image
                        if process_image(img_path, output_path, logo, copyright_text, status_manager):
                            processed_img_count += 1
                else:
                    # Use a temporary directory that we manage explicitly
                    temp_dir = tempfile.mkdtemp()
                    
                    if is_zip_file:
                        # Extract from ZIP file
                        success, extracted_images, slide_count = extract_slides_from_zip(
                            pptx_file,
                            temp_dir,
                            status_mgr=status_manager,
                            target_resolution=TARGET_RESOLUTION
                        )
                    else:
                        # Extract from PPTX file
                        success, extracted_images, slide_count = extract_and_upscale_images(
                            pptx_file, 
                            temp_dir,
                            conversion_method=conversion_method,
                            enable_auto_fallback=enable_auto_fallback,
                            status_mgr=status_manager
                        )
                    
                    # Process each extracted image (add logo and copyright)
                    processed_img_count = 0
                    for img_path in extracted_images:
                        # Temporary image is in temp_dir
                        # Processed image will go directly to output_folder
                        output_filename = os.path.basename(img_path)
                        output_path = os.path.join(output_folder, output_filename)
                        
                        # Process the image
                        if process_image(img_path, output_path, logo, copyright_text, status_manager):
                            processed_img_count += 1
                
                if not success:
                    result["message"] = "Failed to extract images from presentation"
                    results.append(result)
                    continue
            finally:
                # Clean up the temporary directory if it was created
                if temp_dir and os.path.exists(temp_dir):
                    shutil.rmtree(temp_dir, ignore_errors=True)
            
            # Move the presentation file to all_pptx folder
            if os.path.exists(pptx_file) and not os.path.exists(new_pptx_path):
                os.rename(pptx_file, new_pptx_path)
            
            # Update result
            result["status"] = "success"
            result["message"] = f"Successfully processed {processed_img_count} images"
            result["slide_count"] = slide_count
            result["output_dir"] = output_folder
        
        except Exception as e:
            result["message"] = f"Error: {str(e)}"
        
        results.append(result)
        
        # Small delay for UI updates
        time.sleep(0.1)
    
    return results

def run_4k_image_processing(ctx_obj: Dict[str, Any], conversion_method: str = "libreoffice", 
                           enable_auto_fallback: bool = True, create_without_logo: bool = False) -> Dict[str, Any]:
    """
    Main function to run the 4K image processing operation
    This replaces the Streamlit page's main() function
    """
    # Increase image decompression limit
    increase_image_decompression_limit()
    
    # Get configuration from context
    config = ctx_obj.get('config')
    
    # Create status manager
    status_manager = StatusManager(verbose=ctx_obj.get('verbose', False))
    
    # Validate input directory
    input_dir = config.input_dir
    if not os.path.exists(input_dir):
        error_msg = f"Input directory not found: {input_dir}"
        status_manager.error(error_msg)
        return {
            "status": "error",
            "message": error_msg,
            "statistics": {"total_files": 0, "processed_count": 0, "error_count": 1}
        }
    
    status_manager.info(f"Scanning input directory: {input_dir}")
    
    # Find presentation files
    presentation_files = find_presentations()
    
    if not presentation_files:
        warning_msg = "No presentation files (.pptx or .zip) found in input directory"
        status_manager.warning(warning_msg)
        return {
            "status": "warning",
            "message": warning_msg,
            "statistics": {"total_files": 0, "processed_count": 0, "error_count": 0}
        }
    
    # Count file types
    pptx_count = sum(1 for f in presentation_files if f.lower().endswith('.pptx'))
    zip_count = sum(1 for f in presentation_files if f.lower().endswith('.zip'))
    
    status_manager.info(f"Found {len(presentation_files)} files: {pptx_count} PPTX, {zip_count} ZIP")
    
    # Convert method string to enum
    try:
        method_enum = ConversionMethod(conversion_method)
    except ValueError:
        method_enum = ConversionMethod.LIBREOFFICE
        status_manager.warning(f"Invalid conversion method '{conversion_method}', using default: libreoffice")
    
    # Process files with progress tracking
    with DreamletProgress(description="Processing presentations", total=len(presentation_files)) as progress:
        
        # Process presentations
        results = process_presentations(
            input_dir, 
            status_manager, 
            create_without_logo_folder=create_without_logo,
            conversion_method=method_enum,
            enable_auto_fallback=enable_auto_fallback
        )
        
        # Update progress
        progress.update(completed=len(presentation_files))
    
    # Calculate statistics
    success_count = sum(1 for r in results if r["status"] == "success")
    error_count = sum(1 for r in results if r["status"] == "error")
    slide_count = sum(r.get("slide_count", 0) for r in results)
    
    # Determine final status
    if error_count > 0 and success_count == 0:
        final_status = "error"
        status_message = f"Failed to process any presentations ({error_count} errors)"
    elif error_count > 0:
        final_status = "warning"
        status_message = f"Processed {success_count} presentations with {error_count} errors"
    elif success_count == 0:
        final_status = "warning"
        status_message = "No presentations were processed"
    else:
        final_status = "success"
        status_message = f"Successfully processed {success_count} presentations ({slide_count} total slides)"
    
    # Show summary
    status_manager.success(status_message) if final_status == "success" else \
    status_manager.warning(status_message) if final_status == "warning" else \
    status_manager.error(status_message)
    
    # Prepare results for report generation
    report_results = {
        "status": final_status,
        "message": status_message,
        "input_stats": {
            "input_directory": input_dir,
            "presentation_files_found": len(presentation_files),
            "pptx_files": pptx_count,
            "zip_files": zip_count
        },
        "processing_results": results,
        "statistics": {
            "total_files": len(presentation_files),
            "processed_count": success_count,
            "error_count": error_count,
            "total_slides": slide_count
        },
        "settings": {
            "conversion_method": conversion_method,
            "enable_auto_fallback": enable_auto_fallback,
            "create_without_logo": create_without_logo
        },
        "errors": [r for r in results if r["status"] == "error"],
        "output_files": [r["output_dir"] for r in results if r["status"] == "success"]
    }
    
    # Generate report
    report_path = generate_report("06", "4K Image Processing", report_results)
    status_manager.info(f"Report saved to: {report_path}", verbose_only=True)
    
    report_results["report_path"] = report_path
    return report_results

@click.command()
@click.pass_context
def process_4k_images(ctx):
    """
    Generate high-resolution 4K images from presentation slides (PPTX) or ZIP archives
    
    This command processes PowerPoint presentations and ZIP archives containing images,
    extracting slides/images and upscaling them to 4K resolution (3840x2160) with
    logo and copyright overlay.
    
    All settings are configured in config.json under "page_06_4k_image":
    - conversion_method: Method for PPTX conversion (libreoffice, python-pptx, pdf2image)
    - enable_auto_fallback: Try other methods if primary fails
    - create_without_logo: Create folder with original images without logo/copyright
    - target_resolution: Output resolution [width, height]
    - logo_path: Path to logo file
    - copyright_path: Path to copyright text file
    
    Examples:
        dreamlet run 06                    # Process with settings from config.json
        dreamlet config show               # View current configuration
        dreamlet config create             # Create default config.json
    """
    
    # Get configuration
    config = ctx.obj['config']
    from cli.config import get_page_config
    page_config = get_page_config(config, 'page_06_4k_image')
    
    # Extract settings from config
    conversion_method = page_config.get('conversion_method', 'libreoffice')
    enable_auto_fallback = page_config.get('enable_auto_fallback', True)
    create_without_logo = page_config.get('create_without_logo', False)
    
    # Check for dry run mode
    if config.dry_run:
        from rich.console import Console
        console = Console()
        console.print("[yellow]DRY RUN MODE - No files will be processed[/yellow]")
        console.print(f"Would process with settings: method={conversion_method}, fallback={enable_auto_fallback}, without_logo={create_without_logo}")
        return
    
    # Run the 4K image processing operation
    try:
        results = run_4k_image_processing(
            ctx.obj, 
            conversion_method=conversion_method,
            enable_auto_fallback=enable_auto_fallback,
            create_without_logo=create_without_logo
        )
        
        # Exit with appropriate code based on results
        if results["status"] == "error":
            sys.exit(1)
        elif results["status"] == "warning":
            sys.exit(2)
        else:
            sys.exit(0)
    
    except KeyboardInterrupt:
        from rich.console import Console
        console = Console()
        console.print("\n[yellow]Operation cancelled by user[/yellow]")
        sys.exit(130)
    except Exception as e:
        from rich.console import Console
        console = Console()
        console.print(f"[red]Unexpected error: {e}[/red]")
        sys.exit(1)

if __name__ == "__main__":
    process_4k_images()